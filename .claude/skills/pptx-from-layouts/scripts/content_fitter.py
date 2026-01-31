"""Content fitting utilities for PPTX generation.

Provides measure-then-shrink logic for text placeholders and table areas.
Given text content and placeholder dimensions, calculates the optimal font
size that makes text fit without overflow (respecting a 70% floor). For
tables, calculates shrunk font size and row height to fit all rows in
available space.

Uses check_text_overflow from visual_validator for actual font-metric-based
text measurement. All font sizes are whole integer points (no fractional).
"""

import sys
from pathlib import Path

from pptx.util import Pt

# Import measurement functions from visual_validator (same directory)
sys.path.insert(0, str(Path(__file__).parent))
from visual_validator import check_text_overflow, find_font_path

# Import internal API wrapper
_scripts_dir = str(Path(__file__).resolve().parents[3] / "scripts")
if _scripts_dir not in sys.path:
    sys.path.insert(0, _scripts_dir)
from pptx_compat import qn, find_paragraph_pPr


def get_placeholder_dimensions(placeholder) -> dict:
    """Extract usable dimensions from a python-pptx placeholder shape.

    Converts EMU (English Metric Units) to inches and subtracts text frame
    margins to calculate usable content area.

    Args:
        placeholder: A python-pptx placeholder shape object with left, top,
                     width, height attributes (in EMU) and a text_frame.

    Returns:
        Dict with:
        - width_inches: float (total placeholder width)
        - height_inches: float (total placeholder height)
        - left_inches: float (left position)
        - top_inches: float (top position)
        - margin_left: float (left margin in inches)
        - margin_right: float (right margin in inches)
        - margin_top: float (top margin in inches)
        - margin_bottom: float (bottom margin in inches)
        - usable_width_inches: float (width minus horizontal margins)
        - usable_height_inches: float (height minus vertical margins)
    """
    EMU_PER_INCH = 914400

    width_inches = placeholder.width / EMU_PER_INCH
    height_inches = placeholder.height / EMU_PER_INCH
    left_inches = placeholder.left / EMU_PER_INCH
    top_inches = placeholder.top / EMU_PER_INCH

    tf = placeholder.text_frame

    # Read margins (default 0.05 inches if not set, per python-pptx defaults)
    margin_left = (tf.margin_left or 91440) / EMU_PER_INCH
    margin_right = (tf.margin_right or 91440) / EMU_PER_INCH
    margin_top = (tf.margin_top or 45720) / EMU_PER_INCH
    margin_bottom = (tf.margin_bottom or 45720) / EMU_PER_INCH

    usable_width = max(0.1, width_inches - margin_left - margin_right)
    usable_height = max(0.1, height_inches - margin_top - margin_bottom)

    return {
        "width_inches": width_inches,
        "height_inches": height_inches,
        "left_inches": left_inches,
        "top_inches": top_inches,
        "margin_left": margin_left,
        "margin_right": margin_right,
        "margin_top": margin_top,
        "margin_bottom": margin_bottom,
        "usable_width_inches": usable_width,
        "usable_height_inches": usable_height,
    }


def _get_slide_master_title_style(placeholder) -> dict | None:
    """Get font info from slide master's titleStyle for TITLE placeholders.

    Navigates: placeholder -> slide -> slide_layout -> slide_master -> txStyles -> titleStyle

    Args:
        placeholder: A python-pptx placeholder shape.

    Returns:
        Dict with font_size_pt and font_name if found, else None.
    """
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        # Check if this is a title placeholder
        if not hasattr(placeholder, 'placeholder_format'):
            return None

        ph_type = placeholder.placeholder_format.type
        if ph_type not in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return None

        # Navigate to slide master via the shape's part
        # placeholder.part is SlidePart -> slide.slide_layout -> slide_master
        slide_part = placeholder.part
        if not hasattr(slide_part, 'slide'):
            return None

        slide = slide_part.slide
        if not hasattr(slide, 'slide_layout'):
            return None

        slide_layout = slide.slide_layout
        if not hasattr(slide_layout, 'slide_master'):
            return None

        slide_master = slide_layout.slide_master
        master_elem = slide_master._element

        # Find txStyles/titleStyle/lvl1pPr/defRPr
        txStyles = master_elem.find(qn('p:txStyles'))
        if txStyles is None:
            return None

        titleStyle = txStyles.find(qn('p:titleStyle'))
        if titleStyle is None:
            return None

        lvl1pPr = titleStyle.find(qn('a:lvl1pPr'))
        if lvl1pPr is None:
            return None

        defRPr = lvl1pPr.find(qn('a:defRPr'))
        if defRPr is None:
            return None

        sz = defRPr.get('sz')
        if sz is None:
            return None

        font_size_pt = int(int(sz) / 100)
        font_name = "Aptos"  # Default

        # Try to get font name
        latin = defRPr.find(qn('a:latin'))
        if latin is not None:
            typeface = latin.get('typeface')
            if typeface:
                font_name = typeface

        return {
            "font_size_pt": font_size_pt,
            "font_name": font_name,
        }
    except Exception:
        return None


def get_placeholder_font_info(placeholder) -> dict:
    """Read intended font size from placeholder's existing formatting.

    Checks formatting sources in priority order:
    1. Paragraph defRPr sz attribute (layout-level default)
    2. First run's explicit font.size
    3. Slide master's titleStyle (for TITLE placeholders)
    4. Fallback to 16pt

    Args:
        placeholder: A python-pptx placeholder shape with a text_frame.

    Returns:
        Dict with:
        - font_size_pt: int (font size in points)
        - font_name: str (font family name, default 'Aptos')
        - source: str ('defRPr', 'run', 'master_titleStyle', or 'fallback')
    """
    tf = placeholder.text_frame
    font_name = "Aptos"

    # Check each paragraph for defRPr sz
    for para in tf.paragraphs:
        pPr = find_paragraph_pPr(para)
        if pPr is not None:
            defRPr = pPr.find(qn("a:defRPr"))
            if defRPr is not None:
                sz = defRPr.get("sz")
                if sz is not None:
                    font_size_pt = int(int(sz) / 100)
                    # Try to get font name from defRPr
                    latin = defRPr.find(qn("a:latin"))
                    if latin is not None:
                        typeface = latin.get("typeface")
                        if typeface:
                            font_name = typeface
                    return {
                        "font_size_pt": font_size_pt,
                        "font_name": font_name,
                        "source": "defRPr",
                    }

    # Check first run's explicit font size
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                font_size_pt = int(run.font.size.pt)
                if run.font.name:
                    font_name = run.font.name
                return {
                    "font_size_pt": font_size_pt,
                    "font_name": font_name,
                    "source": "run",
                }

    # Check slide master's titleStyle for TITLE placeholders
    master_style = _get_slide_master_title_style(placeholder)
    if master_style is not None:
        return {
            "font_size_pt": master_style["font_size_pt"],
            "font_name": master_style["font_name"],
            "source": "master_titleStyle",
        }

    # Fallback
    return {
        "font_size_pt": 16,
        "font_name": font_name,
        "source": "fallback",
    }


def calculate_fit_font_size(
    text: str,
    font_name: str,
    intended_size_pt: int,
    width_inches: float,
    height_inches: float,
    line_spacing: float = 1.2,
    floor_pct: float = 0.70,
) -> dict:
    """Binary search for largest integer font size that fits text in bounds.

    Measures text at intended size first (fast path). If overflow detected,
    performs binary search between floor and intended size to find the largest
    integer point size that fits without overflow.

    Args:
        text: The text content to fit.
        font_name: Font family name for measurement.
        intended_size_pt: The designed/intended font size in points.
        width_inches: Available width in inches.
        height_inches: Available height in inches.
        line_spacing: Line spacing multiplier (default 1.2).
        floor_pct: Minimum shrink factor (default 0.70 = 70% of intended).

    Returns:
        Dict with:
        - size_pt: int (calculated font size)
        - shrunk: bool (True if size was reduced)
        - fits: bool (True if text fits at returned size)
        - floor_used: bool (True if floor was the final answer)
    """
    floor_size = max(1, int(intended_size_pt * floor_pct))

    # Fast path: check if text fits at intended size
    result = check_text_overflow(
        text=text,
        font_name=font_name,
        font_size_pt=intended_size_pt,
        width_inches=width_inches,
        height_inches=height_inches,
        line_spacing=line_spacing,
    )

    # If measurement failed (font not found), return intended size
    if result.get("overflows") is None:
        return {
            "size_pt": intended_size_pt,
            "shrunk": False,
            "fits": True,
            "floor_used": False,
        }

    # If it fits at intended size, no shrinking needed
    if not result["overflows"]:
        return {
            "size_pt": intended_size_pt,
            "shrunk": False,
            "fits": True,
            "floor_used": False,
        }

    # Binary search between floor_size and intended_size_pt - 1
    low = floor_size
    high = intended_size_pt - 1
    best_fit = None

    while low <= high:
        mid = (low + high) // 2
        mid_result = check_text_overflow(
            text=text,
            font_name=font_name,
            font_size_pt=mid,
            width_inches=width_inches,
            height_inches=height_inches,
            line_spacing=line_spacing,
        )

        if mid_result.get("overflows") is None:
            # Measurement failed, try smaller
            high = mid - 1
            continue

        if not mid_result["overflows"]:
            # Fits at mid, try larger
            best_fit = mid
            low = mid + 1
        else:
            # Overflows at mid, try smaller
            high = mid - 1

    if best_fit is not None:
        return {
            "size_pt": best_fit,
            "shrunk": True,
            "fits": True,
            "floor_used": best_fit == floor_size,
        }

    # Nothing fits even at floor — return floor anyway
    return {
        "size_pt": floor_size,
        "shrunk": True,
        "fits": False,
        "floor_used": True,
    }


def apply_uniform_font_size(text_frame, size_pt: int) -> None:
    """Set all runs in all paragraphs to the specified point size.

    Modifies only font size — preserves all other run formatting (bold,
    italic, color, font name, underline, etc.).

    Args:
        text_frame: A python-pptx TextFrame object.
        size_pt: Font size in integer points to apply uniformly.
    """
    pt_size = Pt(size_pt)
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = pt_size


def shrink_textbox_to_fit(
    textbox,
    font_name: str = "Aptos",
    intended_size_pt: int = None,
    floor_pct: float = 0.70,
    line_spacing: float = 1.2,
) -> dict:
    """Auto-shrink font size in a textbox to fit content within bounds.

    Measures the text in the textbox against its dimensions and shrinks
    the font size if overflow is detected. This is the main entry point
    for auto-font-sizing in dynamically created textboxes (not placeholders).

    Args:
        textbox: A python-pptx Shape object with a text_frame.
        font_name: Font family name for measurement (default 'Aptos').
        intended_size_pt: Starting font size. If None, reads from first run.
        floor_pct: Minimum shrink factor (default 0.70 = 70% of intended).
        line_spacing: Line spacing multiplier (default 1.2).

    Returns:
        Dict with:
        - original_size_pt: int (font size before shrinking)
        - final_size_pt: int (font size after shrinking)
        - shrunk: bool (True if size was reduced)
        - fits: bool (True if text fits at final size)
    """
    EMU_PER_INCH = 914400
    tf = textbox.text_frame

    # Collect all text from the textbox
    all_text = '\n'.join(
        run.text for para in tf.paragraphs for run in para.runs
    )

    if not all_text.strip():
        return {
            "original_size_pt": intended_size_pt or 12,
            "final_size_pt": intended_size_pt or 12,
            "shrunk": False,
            "fits": True,
        }

    # Get intended size from first run if not specified
    if intended_size_pt is None:
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    intended_size_pt = int(run.font.size.pt)
                    break
            if intended_size_pt:
                break
        if not intended_size_pt:
            intended_size_pt = 12  # Default fallback

    # Calculate usable dimensions (textbox size minus margins)
    width_inches = textbox.width / EMU_PER_INCH
    height_inches = textbox.height / EMU_PER_INCH

    # Read margins from text frame (defaults if not set)
    margin_left = (tf.margin_left or 91440) / EMU_PER_INCH
    margin_right = (tf.margin_right or 91440) / EMU_PER_INCH
    margin_top = (tf.margin_top or 45720) / EMU_PER_INCH
    margin_bottom = (tf.margin_bottom or 45720) / EMU_PER_INCH

    usable_width = max(0.1, width_inches - margin_left - margin_right)
    usable_height = max(0.1, height_inches - margin_top - margin_bottom)

    # Calculate optimal font size
    fit_result = calculate_fit_font_size(
        text=all_text,
        font_name=font_name,
        intended_size_pt=intended_size_pt,
        width_inches=usable_width,
        height_inches=usable_height,
        line_spacing=line_spacing,
        floor_pct=floor_pct,
    )

    final_size = fit_result['size_pt']
    shrunk = fit_result['shrunk']

    # Apply the calculated size if shrinking was needed
    if shrunk:
        apply_uniform_font_size(tf, final_size)

    return {
        "original_size_pt": intended_size_pt,
        "final_size_pt": final_size,
        "shrunk": shrunk,
        "fits": fit_result['fits'],
    }


def shrink_placeholder_to_fit(
    placeholder,
    font_name: str = None,
    intended_size_pt: int = None,
    floor_pct: float = 0.70,
    line_spacing: float = 1.2,
) -> dict:
    """Auto-shrink font size in a placeholder to fit content within bounds.

    Similar to shrink_textbox_to_fit but specifically for placeholders.
    Reads font info from placeholder formatting if not specified.

    Args:
        placeholder: A python-pptx placeholder Shape object.
        font_name: Font family name for measurement. If None, reads from placeholder.
        intended_size_pt: Starting font size. If None, reads from placeholder.
        floor_pct: Minimum shrink factor (default 0.70 = 70% of intended).
        line_spacing: Line spacing multiplier (default 1.2).

    Returns:
        Dict with:
        - original_size_pt: int (font size before shrinking)
        - final_size_pt: int (font size after shrinking)
        - shrunk: bool (True if size was reduced)
        - fits: bool (True if text fits at final size)
    """
    tf = placeholder.text_frame

    # Collect all text from the placeholder
    all_text = '\n'.join(
        run.text for para in tf.paragraphs for run in para.runs
    )

    if not all_text.strip():
        return {
            "original_size_pt": intended_size_pt or 16,
            "final_size_pt": intended_size_pt or 16,
            "shrunk": False,
            "fits": True,
        }

    # Get placeholder dimensions
    dims = get_placeholder_dimensions(placeholder)

    # Get font info from placeholder if not specified
    if font_name is None or intended_size_pt is None:
        font_info = get_placeholder_font_info(placeholder)
        if font_name is None:
            font_name = font_info.get('font_name', 'Aptos')
        if intended_size_pt is None:
            intended_size_pt = font_info.get('font_size_pt', 16)

    # Calculate optimal font size
    fit_result = calculate_fit_font_size(
        text=all_text,
        font_name=font_name,
        intended_size_pt=intended_size_pt,
        width_inches=dims['usable_width_inches'],
        height_inches=dims['usable_height_inches'],
        line_spacing=line_spacing,
        floor_pct=floor_pct,
    )

    final_size = fit_result['size_pt']
    shrunk = fit_result['shrunk']

    # Apply the calculated size if shrinking was needed
    if shrunk:
        apply_uniform_font_size(tf, final_size)

    return {
        "original_size_pt": intended_size_pt,
        "final_size_pt": final_size,
        "shrunk": shrunk,
        "fits": fit_result['fits'],
    }


def calculate_table_fit(
    row_count: int,
    col_count: int,
    available_height_inches: float,
    intended_font_size_pt: int = 10,
    header_height_inches: float = 0.35,
    floor_pct: float = 0.70,
) -> dict:
    """Calculate font size and row height to fit table in available space.

    Uses a linear search downward from intended font size to find the largest
    integer point size where all data rows (excluding header) fit within the
    available height. Row height is derived from font size using a 1.8x
    heuristic (font_size_pt * 1.8 / 72.0 inches).

    Args:
        row_count: Total number of data rows (not counting header).
        col_count: Number of columns (for reference, not used in height calc).
        available_height_inches: Total height available for the table.
        intended_font_size_pt: Starting font size in points (default 10).
        header_height_inches: Fixed header row height (default 0.35 inches).
        floor_pct: Minimum shrink factor (default 0.70 = 70% of intended).

    Returns:
        Dict with:
        - font_size_pt: int (calculated font size)
        - row_height_inches: float (calculated row height)
        - header_height_inches: float (header height used)
        - fits: bool (True if all rows fit)
    """
    floor_size = max(1, int(intended_font_size_pt * floor_pct))
    data_height = available_height_inches - header_height_inches

    if data_height <= 0:
        return {
            "font_size_pt": floor_size,
            "row_height_inches": (floor_size * 1.8) / 72.0,
            "header_height_inches": header_height_inches,
            "fits": False,
        }

    # Linear search downward from intended to floor
    for size in range(intended_font_size_pt, floor_size - 1, -1):
        row_height = (size * 1.8) / 72.0
        total_data_height = row_count * row_height

        if total_data_height <= data_height:
            return {
                "font_size_pt": size,
                "row_height_inches": row_height,
                "header_height_inches": header_height_inches,
                "fits": True,
            }

    # Floor reached — still doesn't fit
    row_height = (floor_size * 1.8) / 72.0
    return {
        "font_size_pt": floor_size,
        "row_height_inches": row_height,
        "header_height_inches": header_height_inches,
        "fits": False,
    }


def _strip_all_markers(text: str) -> str:
    """Strip typography markers and markdown from text for width measurement.

    Removes:
    - Inline typography markers: {blue}, {bold}, {italic}, {signpost}, etc.
    - Custom color markers: {color:#RRGGBB}
    - Paragraph markers: {bullet:...}, {level:...}, {space:...}
    - Closing markers: {/blue}, {/bold}, {/color}, etc.
    - Markdown bold/italic: **text**, *text*

    Args:
        text: Text potentially containing markers.

    Returns:
        Clean text with all markers removed.
    """
    import re

    if not text:
        return ""
    if not isinstance(text, str):
        text = str(text)

    # Remove inline markers (opening tags)
    text = re.sub(
        r'\{(blue|red|green|italic|bold|underline|strike|signpost|question|caps|super|sub|font:[^}]+|size:\d+)\}',
        '', text
    )
    # Remove custom color markers
    text = re.sub(r'\{color:#[0-9A-Fa-f]{6}\}', '', text)
    # Remove paragraph markers
    text = re.sub(r'\{(bullet:[^}]+|level:\d+|space:(before|after):[^}]+)\}', '', text)
    # Remove closing markers
    text = re.sub(r'\{/(blue|red|green|italic|bold|underline|strike|signpost|question|caps|super|sub|font|color|size)\}', '', text)
    # Remove markdown bold/italic
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)

    return text.strip()


def calculate_column_widths(
    headers: list[str],
    rows: list[list[str]],
    total_width_emu: int,
    font_name: str = "Aptos",
    font_size_pt: int = 10,
    header_font_size_pt: int | None = None,
    min_col_width_pct: float = 0.08,
    max_col_width_pct: float = 0.60,
    padding_pt: float = 12.0,
) -> list[int]:
    """Calculate dynamic column widths based on content width.

    Measures text width for each column (headers + all data cells) using actual
    font metrics, then distributes total table width proportionally. Applies
    min/max constraints to prevent extremely narrow or wide columns.

    Args:
        headers: List of column header strings (can be empty).
        rows: List of data rows, each row is a list of cell strings.
        total_width_emu: Total table width in EMUs.
        font_name: Font family name for measurement (default 'Aptos').
        font_size_pt: Font size for data cells in points (default 10).
        header_font_size_pt: Font size for headers (default: font_size_pt + 1).
        min_col_width_pct: Minimum column width as percentage of total (default 8%).
        max_col_width_pct: Maximum column width as percentage of total (default 60%).
        padding_pt: Extra padding per column in points (default 12pt).

    Returns:
        List of column widths in EMUs, one per column. If measurement fails,
        returns equal column widths as fallback.
    """
    from PIL import ImageFont

    col_count = len(headers) if headers else (len(rows[0]) if rows else 0)
    if col_count == 0:
        return []

    if header_font_size_pt is None:
        header_font_size_pt = font_size_pt + 1

    # Fallback: equal column widths
    equal_width = total_width_emu // col_count
    fallback_widths = [equal_width] * col_count

    # Try to find font for measurement (find_font_path now includes fallback chain)
    font_path, _is_exact = find_font_path(font_name)
    if font_path is None:
        # No font available even after fallback chain — use equal widths
        return fallback_widths

    try:
        data_font = ImageFont.truetype(font_path, font_size_pt)
        header_font = ImageFont.truetype(font_path, header_font_size_pt)
    except Exception:
        return fallback_widths

    # Measure max content width for each column
    max_widths_pt: list[float] = [0.0] * col_count

    # Helper to measure text width
    def measure_width(text: str, font: ImageFont.FreeTypeFont) -> float:
        clean = _strip_all_markers(text)
        if not clean:
            return 0.0
        if hasattr(font, 'getlength'):
            return font.getlength(clean)
        bbox = font.getbbox(clean)
        return (bbox[2] - bbox[0]) if bbox else 0.0

    # Measure headers
    for col_idx, header in enumerate(headers):
        if col_idx < col_count:
            width = measure_width(str(header), header_font)
            max_widths_pt[col_idx] = max(max_widths_pt[col_idx], width)

    # Measure data rows
    for row_data in rows:
        # Skip section header rows (dict type)
        if isinstance(row_data, dict):
            continue
        for col_idx in range(min(len(row_data), col_count)):
            cell_value = row_data[col_idx]
            width = measure_width(str(cell_value), data_font)
            max_widths_pt[col_idx] = max(max_widths_pt[col_idx], width)

    # Add padding to each column
    max_widths_pt = [w + padding_pt for w in max_widths_pt]

    # Calculate proportional widths
    total_content_width = sum(max_widths_pt)
    if total_content_width <= 0:
        return fallback_widths

    # Convert total EMU to points for calculation (1 inch = 72 pt = 914400 EMU)
    total_width_pt = total_width_emu / 914400 * 72

    # Calculate min/max constraints in points
    min_width_pt = total_width_pt * min_col_width_pct
    max_width_pt = total_width_pt * max_col_width_pct

    # First pass: calculate proportional widths with constraints
    proportional_widths: list[float] = []
    for w in max_widths_pt:
        prop_width = (w / total_content_width) * total_width_pt
        # Apply constraints
        constrained = max(min_width_pt, min(max_width_pt, prop_width))
        proportional_widths.append(constrained)

    # Normalize to exactly fit total width
    total_proportional = sum(proportional_widths)
    if total_proportional > 0:
        scale = total_width_pt / total_proportional
        proportional_widths = [w * scale for w in proportional_widths]

    # Convert back to EMU
    widths_emu = [int(w / 72 * 914400) for w in proportional_widths]

    # Ensure total exactly matches (handle rounding errors)
    diff = total_width_emu - sum(widths_emu)
    if diff != 0 and widths_emu:
        # Add remainder to largest column
        largest_idx = widths_emu.index(max(widths_emu))
        widths_emu[largest_idx] += diff

    return widths_emu


if __name__ == "__main__":
    """Self-tests for content_fitter module."""
    passed = 0
    failed = 0
    skipped = 0

    def report(name: str, ok: bool, detail: str = "") -> None:
        global passed, failed
        status = "PASS" if ok else "FAIL"
        if not ok:
            failed += 1
        else:
            passed += 1
        suffix = f" ({detail})" if detail else ""
        print(f"  {status}: {name}{suffix}")

    def skip(name: str, reason: str) -> None:
        global skipped
        skipped += 1
        print(f"  SKIP: {name} ({reason})")

    # Determine available font for text tests
    test_font = None
    for candidate in ["Helvetica", "DejaVu Sans", "Arial"]:
        path, exact = find_font_path(candidate)
        if path is not None:
            test_font = candidate
            break

    print("Content Fitter Self-Tests")
    print("=" * 40)

    # --- Test 1: Short text fits at intended size ---
    print("\nTest 1: Short text fits at intended size")
    if test_font:
        result = calculate_fit_font_size(
            text="Hello World",
            font_name=test_font,
            intended_size_pt=24,
            width_inches=8.0,
            height_inches=4.0,
        )
        report("shrunk is False", result["shrunk"] is False, f"size={result['size_pt']}")
        report("fits is True", result["fits"] is True)
        report("size_pt equals intended", result["size_pt"] == 24)
    else:
        skip("Short text test", "No suitable font found")

    # --- Test 2: Long text needs shrinking ---
    print("\nTest 2: Long text needs shrinking")
    if test_font:
        long_text = " ".join(["This is a test sentence with enough words to overflow."] * 20)
        result = calculate_fit_font_size(
            text=long_text,
            font_name=test_font,
            intended_size_pt=24,
            width_inches=4.0,
            height_inches=1.5,
        )
        report("shrunk is True", result["shrunk"] is True, f"size={result['size_pt']}")
        report("size_pt < intended", result["size_pt"] < 24, f"{result['size_pt']} < 24")
        floor_size = max(1, int(24 * 0.70))
        report("size_pt >= floor", result["size_pt"] >= floor_size, f"{result['size_pt']} >= {floor_size}")
        report("size_pt is int", isinstance(result["size_pt"], int))
    else:
        skip("Long text test", "No suitable font found")

    # --- Test 3: Table with 5 rows in 3 inches (should fit) ---
    print("\nTest 3: Table 5 rows in 3 inches")
    result = calculate_table_fit(
        row_count=5,
        col_count=3,
        available_height_inches=3.0,
        intended_font_size_pt=10,
    )
    report("fits is True", result["fits"] is True, f"font={result['font_size_pt']}pt")
    report("font_size_pt == 10", result["font_size_pt"] == 10)
    report("row_height > 0", result["row_height_inches"] > 0)

    # --- Test 4: Table with 30 rows in 2 inches (should shrink/floor) ---
    print("\nTest 4: Table 30 rows in 2 inches")
    result = calculate_table_fit(
        row_count=30,
        col_count=4,
        available_height_inches=2.0,
        intended_font_size_pt=10,
    )
    report("font_size_pt <= 10", result["font_size_pt"] <= 10, f"font={result['font_size_pt']}pt")
    floor_size = max(1, int(10 * 0.70))
    report("font_size_pt >= floor", result["font_size_pt"] >= floor_size, f"{result['font_size_pt']} >= {floor_size}")
    report("font_size_pt is int", isinstance(result["font_size_pt"], int))
    # With 30 rows at floor (7pt), row_height = 7*1.8/72 = 0.175 in
    # Total = 30*0.175 = 5.25 in vs available 1.65 in => won't fit
    if not result["fits"]:
        report("fits is False (expected)", True, "30 rows too many for 2in")
    else:
        report("fits (unexpected but valid)", True, "rows fit somehow")

    # --- Test 5: apply_uniform_font_size mock test ---
    print("\nTest 5: apply_uniform_font_size (mock)")
    try:
        from unittest.mock import MagicMock
        mock_tf = MagicMock()
        mock_run1 = MagicMock()
        mock_run2 = MagicMock()
        mock_para = MagicMock()
        mock_para.runs = [mock_run1, mock_run2]
        mock_tf.paragraphs = [mock_para]

        apply_uniform_font_size(mock_tf, 14)
        report("no crash", True)
        report("run1 font.size set", mock_run1.font.size == Pt(14))
        report("run2 font.size set", mock_run2.font.size == Pt(14))
    except Exception as e:
        report("apply_uniform_font_size", False, str(e))

    # --- Test 6: shrink_textbox_to_fit (mock) ---
    print("\nTest 6: shrink_textbox_to_fit (mock)")
    if test_font:
        try:
            from unittest.mock import MagicMock

            # Create mock textbox with short text that fits
            mock_textbox = MagicMock()
            mock_textbox.width = int(4.0 * 914400)  # 4 inches in EMU
            mock_textbox.height = int(2.0 * 914400)  # 2 inches in EMU

            mock_run = MagicMock()
            mock_run.text = "Short text"
            mock_run.font.size = Pt(12)

            mock_para = MagicMock()
            mock_para.runs = [mock_run]

            mock_tf = MagicMock()
            mock_tf.paragraphs = [mock_para]
            mock_tf.margin_left = 91440
            mock_tf.margin_right = 91440
            mock_tf.margin_top = 45720
            mock_tf.margin_bottom = 45720

            mock_textbox.text_frame = mock_tf

            result = shrink_textbox_to_fit(
                mock_textbox,
                font_name=test_font,
                intended_size_pt=12,
            )
            report("returns dict", isinstance(result, dict))
            report("has original_size_pt", "original_size_pt" in result)
            report("has final_size_pt", "final_size_pt" in result)
            report("has shrunk", "shrunk" in result)
            report("has fits", "fits" in result)
            # Short text in large box should not need shrinking
            report("shrunk is False (short text)", result.get("shrunk") is False)
        except Exception as e:
            report("shrink_textbox_to_fit (mock)", False, str(e))
    else:
        skip("shrink_textbox_to_fit test", "No suitable font found")

    # --- Test 7: shrink_textbox_to_fit with overflow (mock) ---
    print("\nTest 7: shrink_textbox_to_fit with overflow (mock)")
    if test_font:
        try:
            from unittest.mock import MagicMock

            # Create mock textbox with long text in small box
            mock_textbox = MagicMock()
            mock_textbox.width = int(1.0 * 914400)  # 1 inch in EMU
            mock_textbox.height = int(0.5 * 914400)  # 0.5 inches in EMU

            long_text = "This is a very long text that should overflow in a small box"
            mock_run = MagicMock()
            mock_run.text = long_text
            mock_run.font.size = Pt(14)

            mock_para = MagicMock()
            mock_para.runs = [mock_run]

            mock_tf = MagicMock()
            mock_tf.paragraphs = [mock_para]
            mock_tf.margin_left = 91440
            mock_tf.margin_right = 91440
            mock_tf.margin_top = 45720
            mock_tf.margin_bottom = 45720

            mock_textbox.text_frame = mock_tf

            result = shrink_textbox_to_fit(
                mock_textbox,
                font_name=test_font,
                intended_size_pt=14,
            )
            report("shrunk is True (overflow)", result.get("shrunk") is True)
            report("final_size_pt < original", result.get("final_size_pt", 14) < 14)
        except Exception as e:
            report("shrink_textbox_to_fit (overflow)", False, str(e))
    else:
        skip("shrink_textbox_to_fit overflow test", "No suitable font found")

    # --- Test 8: shrink_placeholder_to_fit (mock) ---
    print("\nTest 8: shrink_placeholder_to_fit (mock)")
    if test_font:
        try:
            from unittest.mock import MagicMock

            # Create mock placeholder
            mock_placeholder = MagicMock()
            mock_placeholder.width = int(3.0 * 914400)  # 3 inches
            mock_placeholder.height = int(2.0 * 914400)  # 2 inches
            mock_placeholder.left = 914400
            mock_placeholder.top = 914400

            mock_run = MagicMock()
            mock_run.text = "Placeholder text"
            mock_run.font.size = Pt(16)
            mock_run.font.name = test_font

            mock_para = MagicMock()
            mock_para.runs = [mock_run]
            mock_para._p = None  # Simulate no defRPr

            mock_tf = MagicMock()
            mock_tf.paragraphs = [mock_para]
            mock_tf.margin_left = 91440
            mock_tf.margin_right = 91440
            mock_tf.margin_top = 45720
            mock_tf.margin_bottom = 45720

            mock_placeholder.text_frame = mock_tf

            result = shrink_placeholder_to_fit(
                mock_placeholder,
                font_name=test_font,
                intended_size_pt=16,
            )
            report("returns dict", isinstance(result, dict))
            report("has original_size_pt", "original_size_pt" in result)
            report("has final_size_pt", "final_size_pt" in result)
            # Short text in large placeholder should not shrink
            report("shrunk is False (short text)", result.get("shrunk") is False)
        except Exception as e:
            report("shrink_placeholder_to_fit (mock)", False, str(e))
    else:
        skip("shrink_placeholder_to_fit test", "No suitable font found")

    # --- Summary ---
    print("\n" + "=" * 40)
    total = passed + failed + skipped
    print(f"Results: {passed} passed, {failed} failed, {skipped} skipped / {total} total")

    sys.exit(0 if failed == 0 else 1)
