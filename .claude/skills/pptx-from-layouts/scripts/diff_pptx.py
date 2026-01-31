#!/usr/bin/env python3
"""
PPTX Diff Report Generator - Create human-readable difference reports.

Generates a markdown diff report comparing two PowerPoint presentations,
showing exactly what differs between them for easy debugging.

Usage:
    python diff_pptx.py output.pptx reference.pptx --output diff-report.md
"""

import argparse
import json
import sys
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor


def rgb_to_hex(rgb) -> str:
    """Convert RGBColor to hex string."""
    if rgb is None:
        return "None"
    try:
        return f"#{rgb}"
    except:
        return str(rgb)


def normalize_bool(val) -> bool:
    """Normalize boolean values - treat None as False for bold/italic comparisons.

    In PowerPoint, None means "inherit from theme" which typically renders
    as not bold/not italic. For comparison purposes, treat None as False.
    """
    return bool(val) if val is not None else False


def extract_slide_info(prs: Presentation, slide_idx: int) -> Dict[str, Any]:
    """Extract detailed info from a slide."""
    slide = prs.slides[slide_idx]
    layout = slide.slide_layout

    # Get layout index
    layout_idx = -1
    for i, l in enumerate(prs.slide_layouts):
        if l == layout:
            layout_idx = i
            break

    info = {
        'layout_name': layout.name if layout else "None",
        'layout_index': layout_idx,
        'shapes': [],
        'tables': [],
        'texts': [],
        'placeholders': []  # Track placeholders by idx for comparison
    }

    for shape in slide.shapes:
        shape_info = {
            'type': str(shape.shape_type),
            'name': shape.name,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
        }

        # Get placeholder index if this is a placeholder
        placeholder_idx = None
        placeholder_type = None
        try:
            if shape.placeholder_format:
                placeholder_idx = shape.placeholder_format.idx
                placeholder_type = str(shape.placeholder_format.type)
                shape_info['placeholder_idx'] = placeholder_idx
                shape_info['placeholder_type'] = placeholder_type
        except (ValueError, AttributeError):
            pass

        # Extract text content
        if shape.has_text_frame:
            text_runs = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        font = run.font
                        # Safely get color
                        color = None
                        try:
                            if font.color and font.color.type is not None and font.color.rgb:
                                color = rgb_to_hex(font.color.rgb)
                        except (AttributeError, Exception):
                            pass

                        text_runs.append({
                            'text': run.text,
                            'bold': font.bold,
                            'italic': font.italic,
                            'size': font.size.pt if font.size else None,
                            'color': color,
                            'font': font.name
                        })
            shape_info['text_runs'] = text_runs

            # Use placeholder_idx as key if available, otherwise shape_name
            text_key = f"placeholder_{placeholder_idx}" if placeholder_idx is not None else shape.name
            info['texts'].append({
                'shape_name': shape.name,
                'text_key': text_key,
                'placeholder_idx': placeholder_idx,
                'placeholder_type': placeholder_type,
                'full_text': shape.text_frame.text[:200],
                'runs': text_runs
            })

            # Also add to placeholders list if it's a placeholder
            if placeholder_idx is not None:
                info['placeholders'].append({
                    'idx': placeholder_idx,
                    'type': placeholder_type,
                    'shape_name': shape.name,
                    'full_text': shape.text_frame.text[:200],
                    'runs': text_runs
                })

        # Extract table content
        if shape.has_table:
            table = shape.table
            cells = []
            for r, row in enumerate(table.rows):
                row_cells = []
                for c, cell in enumerate(row.cells):
                    cell_info = {
                        'row': r,
                        'col': c,
                        'text': cell.text,
                    }
                    # Get fill color if any
                    try:
                        if cell.fill and cell.fill.fore_color:
                            cell_info['fill'] = rgb_to_hex(cell.fill.fore_color.rgb)
                    except:
                        pass
                    row_cells.append(cell_info)
                cells.append(row_cells)
            info['tables'].append({
                'rows': len(table.rows),
                'cols': len(table.columns),
                'cells': cells
            })

        info['shapes'].append(shape_info)

    return info


def compare_slides(out_info: Dict, ref_info: Dict, slide_num: int) -> List[str]:
    """Compare two slide info dicts and return list of differences."""
    diffs = []

    # Layout check
    if out_info['layout_name'] != ref_info['layout_name']:
        diffs.append(f"**Layout:** Expected `{ref_info['layout_name']}` (idx {ref_info['layout_index']}), got `{out_info['layout_name']}` (idx {out_info['layout_index']})")

    # Shape count
    if len(out_info['shapes']) != len(ref_info['shapes']):
        diffs.append(f"**Shape count:** Expected {len(ref_info['shapes'])}, got {len(out_info['shapes'])}")

    # Table check
    if len(out_info['tables']) != len(ref_info['tables']):
        diffs.append(f"**Table count:** Expected {len(ref_info['tables'])}, got {len(out_info['tables'])}")
    else:
        for i, (out_tbl, ref_tbl) in enumerate(zip(out_info['tables'], ref_info['tables'])):
            if out_tbl['rows'] != ref_tbl['rows'] or out_tbl['cols'] != ref_tbl['cols']:
                diffs.append(f"**Table {i+1} size:** Expected {ref_tbl['rows']}x{ref_tbl['cols']}, got {out_tbl['rows']}x{out_tbl['cols']}")
            else:
                # Compare cell content
                for r, (out_row, ref_row) in enumerate(zip(out_tbl['cells'], ref_tbl['cells'])):
                    for c, (out_cell, ref_cell) in enumerate(zip(out_row, ref_row)):
                        if out_cell['text'] != ref_cell['text']:
                            diffs.append(f"**Table {i+1} cell ({r},{c}):** Expected `{ref_cell['text'][:30]}...`, got `{out_cell['text'][:30]}...`")

    # Text content check - prefer matching by placeholder_idx over shape_name
    # Build dictionaries for both matching strategies
    out_by_placeholder = {t['placeholder_idx']: t for t in out_info['texts'] if t.get('placeholder_idx') is not None}
    ref_by_placeholder = {t['placeholder_idx']: t for t in ref_info['texts'] if t.get('placeholder_idx') is not None}

    # Track which refs we've already matched
    matched_refs = set()

    # First, match by placeholder_idx (most reliable for templates)
    for idx, ref_text in ref_by_placeholder.items():
        if idx in out_by_placeholder:
            matched_refs.add(ref_text['shape_name'])
            out_text = out_by_placeholder[idx]
            ref_full = ref_text['full_text'].strip()
            out_full = out_text['full_text'].strip()

            if ref_full != out_full:
                # Show truncated diff
                ref_show = ref_full[:100] + "..." if len(ref_full) > 100 else ref_full
                out_show = out_full[:100] + "..." if len(out_full) > 100 else out_full

                diffs.append(f"**Text mismatch in placeholder {idx} (`{ref_text['shape_name']}`):**")
                diffs.append(f"  - Expected: `{ref_show}`")
                diffs.append(f"  - Got: `{out_show}`")

            # Check typography
            for ref_run in ref_text['runs']:
                matching_run = None
                for out_run in out_text['runs']:
                    if out_run['text'] == ref_run['text']:
                        matching_run = out_run
                        break

                if matching_run:
                    # Normalize None to False for bold/italic comparisons
                    ref_bold = normalize_bool(ref_run['bold'])
                    out_bold = normalize_bool(matching_run['bold'])
                    ref_italic = normalize_bool(ref_run['italic'])
                    out_italic = normalize_bool(matching_run['italic'])

                    if ref_bold != out_bold:
                        diffs.append(f"**Bold mismatch:** `{ref_run['text'][:30]}...` expected bold={ref_run['bold']}, got {matching_run['bold']}")
                    if ref_italic != out_italic:
                        diffs.append(f"**Italic mismatch:** `{ref_run['text'][:30]}...` expected italic={ref_run['italic']}, got {matching_run['italic']}")
                    if ref_run['color'] and matching_run['color'] and ref_run['color'] != matching_run['color']:
                        diffs.append(f"**Color mismatch:** `{ref_run['text'][:30]}...` expected {ref_run['color']}, got {matching_run['color']}")
                    if ref_run['size'] and matching_run['size'] and abs(ref_run['size'] - matching_run['size']) > 1:
                        diffs.append(f"**Size mismatch:** `{ref_run['text'][:30]}...` expected {ref_run['size']}pt, got {matching_run['size']}pt")
        else:
            # Placeholder in ref but not in output
            diffs.append(f"**Missing placeholder {idx}:** `{ref_text['shape_name']}` with text `{ref_text['full_text'][:50]}...`")

    # Now handle non-placeholder shapes (match by shape_name)
    out_texts = {t['shape_name']: t for t in out_info['texts'] if t.get('placeholder_idx') is None}
    ref_texts = {t['shape_name']: t for t in ref_info['texts'] if t.get('placeholder_idx') is None}

    for name, ref_text in ref_texts.items():
        if name in matched_refs:
            continue
        if name not in out_texts:
            diffs.append(f"**Missing shape:** `{name}` with text `{ref_text['full_text'][:50]}...`")
        else:
            out_text = out_texts[name]
            ref_full = ref_text['full_text'].strip()
            out_full = out_text['full_text'].strip()

            if ref_full != out_full:
                ref_show = ref_full[:100] + "..." if len(ref_full) > 100 else ref_full
                out_show = out_full[:100] + "..." if len(out_full) > 100 else out_full

                diffs.append(f"**Text mismatch in `{name}`:**")
                diffs.append(f"  - Expected: `{ref_show}`")
                diffs.append(f"  - Got: `{out_show}`")

            # Check typography
            for ref_run in ref_text['runs']:
                matching_run = None
                for out_run in out_text['runs']:
                    if out_run['text'] == ref_run['text']:
                        matching_run = out_run
                        break

                if matching_run:
                    # Normalize None to False for bold/italic comparisons
                    ref_bold = normalize_bool(ref_run['bold'])
                    out_bold = normalize_bool(matching_run['bold'])
                    ref_italic = normalize_bool(ref_run['italic'])
                    out_italic = normalize_bool(matching_run['italic'])

                    if ref_bold != out_bold:
                        diffs.append(f"**Bold mismatch:** `{ref_run['text'][:30]}...` expected bold={ref_run['bold']}, got {matching_run['bold']}")
                    if ref_italic != out_italic:
                        diffs.append(f"**Italic mismatch:** `{ref_run['text'][:30]}...` expected italic={ref_run['italic']}, got {matching_run['italic']}")
                    if ref_run['color'] and matching_run['color'] and ref_run['color'] != matching_run['color']:
                        diffs.append(f"**Color mismatch:** `{ref_run['text'][:30]}...` expected {ref_run['color']}, got {matching_run['color']}")
                    if ref_run['size'] and matching_run['size'] and abs(ref_run['size'] - matching_run['size']) > 1:
                        diffs.append(f"**Size mismatch:** `{ref_run['text'][:30]}...` expected {ref_run['size']}pt, got {matching_run['size']}pt")

    return diffs


def generate_diff_report(output_path: str, reference_path: str) -> str:
    """Generate a markdown diff report."""
    out_prs = Presentation(output_path)
    ref_prs = Presentation(reference_path)

    lines = [
        f"# PPTX Diff Report",
        f"",
        f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        f"",
        f"**Output:** `{output_path}`",
        f"",
        f"**Reference:** `{reference_path}`",
        f"",
        f"---",
        f"",
        f"## Summary",
        f"",
        f"| Metric | Output | Reference | Match |",
        f"|--------|--------|-----------|-------|",
        f"| Slides | {len(out_prs.slides)} | {len(ref_prs.slides)} | {'Yes' if len(out_prs.slides) == len(ref_prs.slides) else 'No'} |",
        f"| Layouts | {len(out_prs.slide_layouts)} | {len(ref_prs.slide_layouts)} | {'Yes' if len(out_prs.slide_layouts) == len(ref_prs.slide_layouts) else 'No'} |",
        f"",
        f"---",
        f"",
        f"## Slide-by-Slide Comparison",
        f"",
    ]

    total_diffs = 0
    min_slides = min(len(out_prs.slides), len(ref_prs.slides))

    for i in range(min_slides):
        out_info = extract_slide_info(out_prs, i)
        ref_info = extract_slide_info(ref_prs, i)

        diffs = compare_slides(out_info, ref_info, i + 1)
        total_diffs += len(diffs)

        status = "MATCH" if not diffs else f"DIFFERENCES ({len(diffs)})"

        lines.append(f"### Slide {i + 1}: {status}")
        lines.append(f"")
        lines.append(f"**Layout:** `{ref_info['layout_name']}` (index {ref_info['layout_index']})")
        lines.append(f"")

        if diffs:
            for diff in diffs:
                lines.append(f"- {diff}")
            lines.append(f"")
        else:
            lines.append(f"All content matches.")
            lines.append(f"")

        lines.append(f"---")
        lines.append(f"")

    # Handle extra slides
    if len(out_prs.slides) > len(ref_prs.slides):
        lines.append(f"### Extra Slides in Output")
        lines.append(f"")
        for i in range(len(ref_prs.slides), len(out_prs.slides)):
            out_info = extract_slide_info(out_prs, i)
            lines.append(f"- Slide {i + 1}: `{out_info['layout_name']}`")
        lines.append(f"")

    if len(ref_prs.slides) > len(out_prs.slides):
        lines.append(f"### Missing Slides from Output")
        lines.append(f"")
        for i in range(len(out_prs.slides), len(ref_prs.slides)):
            ref_info = extract_slide_info(ref_prs, i)
            lines.append(f"- Slide {i + 1}: `{ref_info['layout_name']}`")
        lines.append(f"")

    # Final summary
    lines.append(f"## Conclusion")
    lines.append(f"")
    if total_diffs == 0 and len(out_prs.slides) == len(ref_prs.slides):
        lines.append(f"**RESULT: EXACT MATCH**")
    else:
        lines.append(f"**RESULT: {total_diffs} differences found**")
    lines.append(f"")

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(description='Generate PPTX diff report')
    parser.add_argument('output', help='Output PPTX to validate')
    parser.add_argument('reference', help='Reference PPTX to compare against')
    parser.add_argument('--output', '-o', dest='report_path', default='diff-report.md',
                        help='Output report path (default: diff-report.md)')

    args = parser.parse_args()

    report = generate_diff_report(args.output, args.reference)

    # Save report
    with open(args.report_path, 'w') as f:
        f.write(report)

    print(f"Diff report saved to: {args.report_path}")

    # Print summary to stdout
    print("\n" + "="*60)
    print("DIFF SUMMARY")
    print("="*60)

    # Count differences
    out_prs = Presentation(args.output)
    ref_prs = Presentation(args.reference)

    print(f"Output slides: {len(out_prs.slides)}")
    print(f"Reference slides: {len(ref_prs.slides)}")

    if len(out_prs.slides) == len(ref_prs.slides):
        print("Slide count: MATCH")
    else:
        print(f"Slide count: MISMATCH ({len(out_prs.slides)} vs {len(ref_prs.slides)})")


if __name__ == '__main__':
    main()
