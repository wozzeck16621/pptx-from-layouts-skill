#!/usr/bin/env python3
"""
generate_config.py: Generate template configuration from a PowerPoint template.

Creates a template config JSON file that maps abstract layout capabilities
to concrete layout indices in the template. This enables template-agnostic
slide generation.

Usage:
    # Generate config from template (extracts profile first)
    python generate_config.py template.pptx --output my-config.json

    # Generate config from existing profile JSON
    python generate_config.py --from-profile my-profile.json --output my-config.json

    # Generate Inner Chapter config (captures hardcoded defaults)
    python generate_config.py --generate-ic-defaults --output inner-chapter-config.json
"""

import argparse
import json
import sys
from pathlib import Path
from collections import defaultdict

# Add schemas to path
_claude_dir = Path(__file__).resolve().parents[3]  # .claude/skills/template-profiler/scripts -> .claude
sys.path.insert(0, str(_claude_dir))
_scripts_dir = str(_claude_dir / "scripts")
if _scripts_dir not in sys.path:
    sys.path.insert(0, _scripts_dir)

from pptx import Presentation
from pptx_compat import get_theme_part, get_part_blob, etree


# =============================================================================
# PROFILE EXTRACTION (minimal subset for config generation)
# =============================================================================

def extract_color_scheme(prs) -> dict:
    """Extract theme color scheme from presentation."""
    colors = {}
    try:
        theme_part = get_theme_part(prs)
        if theme_part:
            theme_xml = etree.fromstring(get_part_blob(theme_part))
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            clrScheme = theme_xml.find('.//a:clrScheme', ns)
            if clrScheme is not None:
                colors["scheme_name"] = (clrScheme.get('name') or '').strip()
                for child in clrScheme:
                    tag_name = child.tag.split('}')[-1]
                    for color_elem in child:
                        color_tag = color_elem.tag.split('}')[-1]
                        if color_tag == 'srgbClr':
                            colors[tag_name] = color_elem.get('val')
                        elif color_tag == 'sysClr':
                            colors[tag_name] = color_elem.get('lastClr', color_elem.get('val'))
    except Exception as e:
        colors["_error"] = str(e)
    return colors


def extract_font_scheme(prs) -> dict:
    """Extract font scheme from presentation."""
    fonts = {}
    try:
        theme_part = get_theme_part(prs)
        if theme_part:
            theme_xml = etree.fromstring(get_part_blob(theme_part))
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            fontScheme = theme_xml.find('.//a:fontScheme', ns)
            if fontScheme is not None:
                fonts["scheme_name"] = fontScheme.get('name')
                majorFont = fontScheme.find('a:majorFont', ns)
                if majorFont is not None:
                    latin = majorFont.find('a:latin', ns)
                    if latin is not None:
                        fonts["major_font"] = latin.get('typeface')
                minorFont = fontScheme.find('a:minorFont', ns)
                if minorFont is not None:
                    latin = minorFont.find('a:latin', ns)
                    if latin is not None:
                        fonts["minor_font"] = latin.get('typeface')
    except Exception as e:
        fonts["_error"] = str(e)
    return fonts


def extract_layout_info(prs) -> list[dict]:
    """Extract layout signatures from presentation."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        signature = defaultdict(int)
        for shape in layout.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                type_map = {
                    PP_PLACEHOLDER.TITLE: "title",
                    PP_PLACEHOLDER.SUBTITLE: "subtitle",
                    PP_PLACEHOLDER.BODY: "body",
                    PP_PLACEHOLDER.CENTER_TITLE: "center_title",
                    PP_PLACEHOLDER.PICTURE: "picture",
                    PP_PLACEHOLDER.TABLE: "table",
                    PP_PLACEHOLDER.CHART: "chart",
                }
                ph_type = type_map.get(ph.type, "other")
                signature[ph_type] += 1

        name = layout.name.lower().replace(' ', '-').replace('_', '-')
        category = get_category(layout.name)

        layouts.append({
            "index": idx,
            "name": layout.name,
            "name_normalized": name,
            "category": category,
            "signature": dict(signature),
        })
    return layouts


def get_category(layout_name: str) -> str:
    """Determine category based on layout name."""
    name = layout_name.lower()
    if name == "master-base" or "master" in name:
        return "master"
    if name.startswith("title"):
        return "title"
    if name.startswith("team-"):
        return "team"
    if name.startswith("contact"):
        return "contact"
    if name.startswith("legal"):
        return "legal"
    if name.startswith("grid-"):
        return "grid"
    if name.startswith("column-"):
        return "column"
    if name.startswith("image-"):
        return "image"
    return "content"


# =============================================================================
# CONFIG GENERATION
# =============================================================================

def infer_brand_config(colors: dict, fonts: dict) -> dict:
    """Infer brand config from color and font schemes.

    Args:
        colors: Color scheme dict from extract_color_scheme
        fonts: Font scheme dict from extract_font_scheme

    Returns:
        BrandConfig-compatible dict
    """
    # Map theme colors to brand colors
    # accent1 is typically the primary brand color
    # dk1 (dark 1) is typically black/dark gray for text
    # dk2 is typically a secondary dark color
    primary = colors.get("accent1", "0196FF")
    secondary = colors.get("dk1", "000000")
    accent = colors.get("dk2") or colors.get("accent2", "595959")

    # Normalize colors (strip # if present, uppercase)
    def normalize_color(c):
        if not c:
            return "000000"
        c = c.lstrip("#").upper()
        return c if len(c) == 6 else "000000"

    # Fonts - major is typically for headers, minor for body
    header_font = fonts.get("major_font", "Aptos")
    body_font = fonts.get("minor_font", "Aptos")

    return {
        "primary_color": normalize_color(primary),
        "secondary_color": normalize_color(secondary),
        "accent_color": normalize_color(accent),
        "header_font": header_font or "Aptos",
        "body_font": body_font or "Aptos",
    }


def infer_layout_mappings(layouts: list[dict]) -> tuple[dict, int, int]:
    """Infer layout capability mappings from layout list.

    Returns:
        Tuple of (layout_mappings dict, branding_layout_index, fallback_layout_index)
    """
    mappings = {}

    # Organize layouts by category and body count
    by_category = defaultdict(list)
    by_body_count = defaultdict(list)

    for layout in layouts:
        cat = layout["category"]
        body_count = layout["signature"].get("body", 0)
        by_category[cat].append(layout)
        if body_count > 0:
            by_body_count[body_count].append(layout)

    # Find specific layouts by name patterns and category
    def find_layout(name_patterns: list[str], category: str = None, body_count: int = None) -> dict | None:
        """Find first matching layout by name pattern or criteria."""
        for layout in layouts:
            name = layout["name_normalized"]
            # Check name patterns
            for pattern in name_patterns:
                if pattern in name:
                    return layout
            # Check category + body count
            if category and layout["category"] == category:
                if body_count is None or layout["signature"].get("body", 0) == body_count:
                    return layout
        return None

    def add_mapping(key: str, layout: dict, use_case: str = ""):
        """Add a layout mapping."""
        mappings[key] = {
            "layout_index": layout["index"],
            "layout_name": layout["name"],
            "use_case": use_case,
        }

    # Title layouts
    title_cover = find_layout(["title-cover", "titlecover", "cover"])
    if title_cover:
        add_mapping("title_cover", title_cover, "Cover page with client, date, logos")

    title_centered = find_layout(["title-centered", "titlecentered", "centered-title"])
    if title_centered:
        add_mapping("title_centered", title_centered, "Section dividers, closing statements")

    # Content layouts
    content_image = find_layout(["content-image-right", "contentimageright", "content-with-image"])
    if content_image:
        add_mapping("content_with_image", content_image, "Bullets with image placeholder")

    content_centered = find_layout(["content-centered", "contentcentered"])
    if content_centered:
        add_mapping("content_centered", content_centered, "Centered content without image")

    story_card = find_layout(["story", "image-right-text-left", "text-left"])
    if story_card:
        add_mapping("story_card", story_card, "Story card with image")

    # Column layouts - find by category and body count
    column_layouts = sorted(by_category.get("column", []), key=lambda x: x["signature"].get("body", 0))
    for layout in column_layouts:
        body_count = layout["signature"].get("body", 0)
        if body_count >= 2 and body_count <= 5:
            key = f"column_{body_count}"
            if key not in mappings:
                add_mapping(key, layout, f"{body_count}-column layout")

    # Contact layout
    contact = find_layout(["contact-black", "contactblack", "contact"])
    if contact:
        add_mapping("contact", contact, "Contact information")
    contact_white = find_layout(["contact-white", "contactwhite"])
    if contact_white:
        add_mapping("contact_white", contact_white, "Contact information (light)")

    # Grid layouts
    for layout in by_category.get("grid", []):
        body_count = layout["signature"].get("body", 0)
        picture_count = layout["signature"].get("picture", 0)
        if "3x2" in layout["name_normalized"] or (picture_count == 3 and body_count in [3, 6]):
            if body_count == 3 and "grid_3x2_3body" not in mappings:
                add_mapping("grid_3x2_3body", layout, "3 columns with images top, 3 body areas")
            elif body_count == 6 and "grid_3x2_6body" not in mappings:
                add_mapping("grid_3x2_6body", layout, "3 columns with images top, 6 body areas")
        elif "2x2" in layout["name_normalized"] or (picture_count == 2 and body_count == 2):
            if "grid_2x2_2body" not in mappings:
                add_mapping("grid_2x2_2body", layout, "2 columns with images top, 2 body areas")

    # Master/branding layout
    master = find_layout(["master-base", "masterbase", "master"])
    branding_index = None
    if master:
        add_mapping("master_base", master, "Branding slide")
        branding_index = master["index"]
    elif layouts:
        # Use last layout as branding if no master found
        branding_index = layouts[-1]["index"]

    # Fallback layout - prefer content-centered or first content layout
    fallback_index = None
    if "content_centered" in mappings:
        fallback_index = mappings["content_centered"]["layout_index"]
        add_mapping("fallback", {
            "index": fallback_index,
            "name": mappings["content_centered"]["layout_name"],
        }, "Fallback for unknown content types")
    elif "content_with_image" in mappings:
        fallback_index = mappings["content_with_image"]["layout_index"]
        add_mapping("fallback", {
            "index": fallback_index,
            "name": mappings["content_with_image"]["layout_name"],
        }, "Fallback for unknown content types")
    elif layouts:
        # Just use first content layout
        for layout in layouts:
            if layout["category"] == "content":
                fallback_index = layout["index"]
                add_mapping("fallback", layout, "Fallback for unknown content types")
                break
        if fallback_index is None:
            fallback_index = layouts[0]["index"]
            add_mapping("fallback", layouts[0], "Fallback for unknown content types")

    return mappings, branding_index, fallback_index


def generate_config(
    template_path: str | None = None,
    profile_data: dict | None = None,
    template_name: str | None = None,
) -> dict:
    """Generate template config from a template file or existing profile.

    Args:
        template_path: Path to .pptx template file
        profile_data: Existing profile dict (alternative to template_path)
        template_name: Name for the template (auto-derived if not provided)

    Returns:
        Template config dict compatible with TemplateConfig schema
    """
    if template_path:
        # Extract from template
        prs = Presentation(template_path)
        colors = extract_color_scheme(prs)
        fonts = extract_font_scheme(prs)
        layouts = extract_layout_info(prs)
        if not template_name:
            template_name = Path(template_path).stem
    elif profile_data:
        # Use existing profile
        colors = profile_data.get("master_styles", {}).get("color_scheme", {})
        fonts = profile_data.get("master_styles", {}).get("font_scheme", {})
        layouts = profile_data.get("layouts", [])
        if not template_name:
            template_name = Path(profile_data.get("template_file", "unknown")).stem
    else:
        raise ValueError("Either template_path or profile_data must be provided")

    # Infer brand config
    brand = infer_brand_config(colors, fonts)

    # Infer layout mappings
    mappings, branding_index, fallback_index = infer_layout_mappings(layouts)

    # Default content type routing
    content_type_routing = {
        "title_slide": "title_cover",
        "section_divider": "title_centered",
        "closing": "title_centered",
        "thank_you": "title_centered",
        "quote": "title_centered",
        "contact": "contact",
        "framework_2col": "column_2",
        "framework_3col": "column_3",
        "framework_4col": "column_4",
        "framework_5col": "column_5",
        "framework_cards": "column_4",
        "comparison": "column_2",
        "content": "content_with_image",
        "deliverables": "content_with_image",
        "timeline": "content_with_image",
        "content_no_image": "content_centered",
        "table": "title_centered",
        "pricing": "title_centered",
        "comparison_tables": "title_centered",
        "table_with_image": "content_with_image",
        "story_card": "story_card",
        "grid_3x2_3body": "grid_3x2_3body",
        "grid_3x2_6body": "grid_3x2_6body",
        "grid_2x2_2body": "grid_2x2_2body",
        "content_image_top_4body": "content_image_top_4body",
    }

    # Build config
    fallback_name = mappings.get("fallback", {}).get("layout_name", "content-centered-a")

    config = {
        "template_name": template_name,
        "template_path": str(template_path) if template_path else None,
        "brand": brand,
        "layout_mappings": mappings,
        "content_type_routing": content_type_routing,
        "requires_branding_slide": branding_index is not None,
        "branding_layout_index": branding_index,
        "fallback_layout_index": fallback_index or 0,
        "fallback_layout_name": fallback_name,
    }

    return config


def generate_inner_chapter_config() -> dict:
    """Generate the Inner Chapter config that captures exact hardcoded behavior.

    This creates a config that matches the IC_LAYOUTS dict exactly.
    """
    return {
        "template_name": "inner-chapter",
        "template_path": "template/inner-chapter.pptx",
        "brand": {
            "primary_color": "0196FF",
            "secondary_color": "000000",
            "accent_color": "595959",
            "header_font": "Aptos",
            "body_font": "Aptos",
        },
        "layout_mappings": {
            "master_base": {
                "layout_index": 58,
                "layout_name": "master-base",
                "use_case": "Branding slide (always Slide 1)",
            },
            "title_cover": {
                "layout_index": 0,
                "layout_name": "title-cover",
                "use_case": "Cover page with client, date, logos",
            },
            "title_centered": {
                "layout_index": 2,
                "layout_name": "title-centered",
                "use_case": "Section dividers, closing statements",
            },
            "content_with_image": {
                "layout_index": 3,
                "layout_name": "content-image-right-a",
                "use_case": "Bullets with image placeholder",
            },
            "content_centered": {
                "layout_index": 45,
                "layout_name": "content-centered-a",
                "use_case": "Centered content without image",
            },
            "story_card": {
                "layout_index": 8,
                "layout_name": "content-image-right-text-left",
                "use_case": "Story card with image right, text left",
            },
            "column_2": {
                "layout_index": 6,
                "layout_name": "column-2-centered",
                "use_case": "2-column comparisons",
            },
            "column_3": {
                "layout_index": 35,
                "layout_name": "column-3-centered-a",
                "use_case": "3-column comparisons",
            },
            "column_4": {
                "layout_index": 31,
                "layout_name": "column-4-centered",
                "use_case": "4-column framework/flow diagrams",
            },
            "column_5": {
                "layout_index": 27,
                "layout_name": "column-5-centered",
                "use_case": "5-column cards/comparisons",
            },
            "contact": {
                "layout_index": 52,
                "layout_name": "contact-black",
                "use_case": "Contact information (dark)",
            },
            "contact_white": {
                "layout_index": 51,
                "layout_name": "contact-white",
                "use_case": "Contact information (light)",
            },
            "grid_3x2_3body": {
                "layout_index": 36,
                "layout_name": "grid-3x2-image-top-3-body",
                "use_case": "3 columns with images top, 3 body areas",
            },
            "grid_3x2_6body": {
                "layout_index": 37,
                "layout_name": "grid-3x2-image-top-6-body-a",
                "use_case": "3 columns with images top, 6 body areas",
            },
            "grid_2x2_2body": {
                "layout_index": 42,
                "layout_name": "grid-2x2-image-top-2-body-a",
                "use_case": "2 columns with images top, 2 body areas",
            },
            "grid_2x2_2body_b": {
                "layout_index": 44,
                "layout_name": "grid-2x2-image-top-2-body-b",
                "use_case": "2 columns with images top, 2 body areas (variant B)",
            },
            "content_image_top_4body": {
                "layout_index": 32,
                "layout_name": "content-image-top-4-body",
                "use_case": "4 columns with images top",
            },
            "fallback": {
                "layout_index": 45,
                "layout_name": "content-centered-a",
                "use_case": "Fallback for unknown content types",
            },
        },
        "content_type_routing": {
            "title_slide": "title_cover",
            "section_divider": "title_centered",
            "closing": "title_centered",
            "thank_you": "title_centered",
            "quote": "title_centered",
            "contact": "contact",
            "framework_2col": "column_2",
            "framework_3col": "column_3",
            "framework_4col": "column_4",
            "framework_5col": "column_5",
            "framework_cards": "column_4",
            "comparison": "column_2",
            "content": "content_with_image",
            "deliverables": "content_with_image",
            "timeline": "content_with_image",
            "content_no_image": "content_centered",
            "table": "title_centered",
            "pricing": "title_centered",
            "comparison_tables": "title_centered",
            "table_with_image": "content_with_image",
            "story_card": "story_card",
            "grid_3x2_3body": "grid_3x2_3body",
            "grid_3x2_6body": "grid_3x2_6body",
            "grid_2x2_2body": "grid_2x2_2body",
            "grid_2x2_2body_b": "grid_2x2_2body_b",
            "content_image_top_4body": "content_image_top_4body",
        },
        "requires_branding_slide": True,
        "branding_layout_index": 58,
        "fallback_layout_index": 45,
        "fallback_layout_name": "content-centered-a",
    }


# =============================================================================
# MAIN
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate template configuration for template-agnostic PPTX generation"
    )
    parser.add_argument(
        "template",
        nargs="?",
        help="Path to .pptx template file",
    )
    parser.add_argument(
        "--from-profile",
        help="Generate config from existing profile JSON instead of template",
    )
    parser.add_argument(
        "--output", "-o",
        required=True,
        help="Output path for config JSON",
    )
    parser.add_argument(
        "--name", "-n",
        help="Template name (auto-derived from filename if not provided)",
    )
    parser.add_argument(
        "--generate-ic-defaults",
        action="store_true",
        help="Generate Inner Chapter config capturing exact hardcoded defaults",
    )

    args = parser.parse_args()

    # Handle IC defaults mode
    if args.generate_ic_defaults:
        config = generate_inner_chapter_config()
        output_path = Path(args.output)
        output_path.write_text(json.dumps(config, indent=2, ensure_ascii=False), encoding="utf-8")
        print(f"Created Inner Chapter config: {output_path}", file=sys.stderr)
        print(f"  Layout mappings: {len(config['layout_mappings'])}", file=sys.stderr)
        print(f"  Content type routes: {len(config['content_type_routing'])}", file=sys.stderr)
        return

    # Validate inputs
    if not args.template and not args.from_profile:
        parser.error("Either template path or --from-profile is required")

    # Generate config
    if args.from_profile:
        profile_path = Path(args.from_profile)
        if not profile_path.exists():
            print(f"Error: Profile not found: {args.from_profile}", file=sys.stderr)
            sys.exit(1)
        profile_data = json.loads(profile_path.read_text(encoding="utf-8"))
        config = generate_config(profile_data=profile_data, template_name=args.name)
    else:
        template_path = Path(args.template)
        if not template_path.exists():
            print(f"Error: Template not found: {args.template}", file=sys.stderr)
            sys.exit(1)
        config = generate_config(template_path=str(template_path), template_name=args.name)

    # Write output
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(config, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"Created template config: {output_path}", file=sys.stderr)
    print(f"  Template: {config['template_name']}", file=sys.stderr)
    print(f"  Brand colors: {config['brand']['primary_color']}, {config['brand']['secondary_color']}", file=sys.stderr)
    print(f"  Fonts: {config['brand']['header_font']}, {config['brand']['body_font']}", file=sys.stderr)
    print(f"  Layout mappings: {len(config['layout_mappings'])}", file=sys.stderr)
    print(f"  Content type routes: {len(config['content_type_routing'])}", file=sys.stderr)


if __name__ == "__main__":
    main()
