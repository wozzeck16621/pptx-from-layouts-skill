#!/usr/bin/env python3
"""
generate_pptx.py: Generate PowerPoint presentations from layout plan JSON files.

Usage:
    python generate_pptx.py layout_plan.json --template template.pptx --output output.pptx

This script handles:
- Routes slides based on content_type field (primary routing)
- Discovers placeholders from slides at runtime
- Layouts without TITLE placeholder (uses largest BODY for title)
- Structured columns with headers and body content
- Multiple tables per slide with flexible schema
- Timeline visualization
- Deliverables with title + description (multiple input patterns)
- Image placeholders (renders as labeled boxes)
- Footnotes, callouts, and quotes
- Subsection content blocks
- Branding slides with no placeholders

Version: 2.2 - Self-contained skill
"""

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Set up paths for skill's local modules
_SCRIPT_DIR = Path(__file__).resolve().parent
_SKILL_DIR = _SCRIPT_DIR.parent
_LIB_DIR = _SKILL_DIR / "lib"
_SCHEMAS_DIR = _SKILL_DIR / "schemas"
_PROJECT_ROOT = _SKILL_DIR.parents[2]

# Add skill directories to path
for _path in [str(_SKILL_DIR), str(_LIB_DIR), str(_SCHEMAS_DIR), str(_SCRIPT_DIR)]:
    if _path not in sys.path:
        sys.path.insert(0, _path)

# Schema validation
from schemas import LayoutPlan, GenerationResult, SlideResult, SlideError, SlideStatus
from schemas.layout_plan import VALID_CONTENT_TYPES
from pydantic import ValidationError

# Template config support (optional - falls back to IC defaults)
try:
    from schemas.template_config import TemplateConfig
    _HAS_TEMPLATE_CONFIG = True
except ImportError:
    _HAS_TEMPLATE_CONFIG = False
    TemplateConfig = None

# Performance instrumentation (optional)
try:
    from performance import PerfContext, PerfTimer, identify_bottlenecks
    _HAS_PERFORMANCE = True
except ImportError:
    _HAS_PERFORMANCE = False
    PerfContext = None
    PerfTimer = None
    identify_bottlenecks = None

# Margin enforcement for shape positioning
from margins import get_safe_area

# Global template config (set in main() when --config is used)
_template_config: "TemplateConfig | None" = None


def get_template_config() -> "TemplateConfig | None":
    """Get the current template config."""
    return _template_config


def set_template_config(config: "TemplateConfig | None"):
    """Set the template config."""
    global _template_config
    _template_config = config


def get_color_map() -> dict[str, str]:
    """Get color map from template config or defaults."""
    config = get_template_config()
    if config:
        return config.get_color_map()
    # IC defaults
    return {
        'blue': '0196FF',
        'black': '000000',
        'gray': '595959',
    }


def get_preset_styles() -> dict[str, dict]:
    """Get preset styles from template config or defaults."""
    config = get_template_config()
    if config:
        return config.get_preset_styles()
    # IC defaults
    return {
        'signpost': {'color': '0196FF', 'size': 14},
        'question': {'color': '0196FF', 'italic': True},
    }


def get_brand_font(font_type: str = 'body') -> str:
    """Get brand font from template config or default, with fallback chain.

    Applies the font fallback chain (Aptos → Calibri → Arial → sans-serif)
    to ensure the returned font is available on the system for consistent
    measurement. The actual font name is still used in PPTX output, but
    this function returns a font that exists for measurement purposes.

    Args:
        font_type: 'header' or 'body'

    Returns:
        Font name string (may be a fallback if preferred font unavailable)
    """
    # Import font fallback utilities (lazy import to avoid circular deps)
    try:
        from font_fallback import resolve_font_for_pptx
    except ImportError:
        # Fallback if font_fallback module not available
        def resolve_font_for_pptx(font_name):
            return font_name

    config = get_template_config()
    if config:
        if font_type == 'header':
            preferred = config.brand.header_font
        else:
            preferred = config.brand.body_font
    else:
        preferred = 'Aptos'  # IC default

    return resolve_font_for_pptx(preferred)

from pptx_compat import (
    get_run_properties,
    get_paragraph_properties,
    get_table_from_shape,
    remove_shape,
    delete_slide as compat_delete_slide,
    qn,
    make_sub_element,
    find_element,
    find_all_elements,
    etree,
)


# Phase 4: Visual validation (optional dependency)
try:
    from visual_validator import check_text_overflow
    _HAS_VISUAL_VALIDATOR = True
except ImportError:
    _HAS_VISUAL_VALIDATOR = False

# Phase 5: Content fitting (optional dependency)
try:
    from content_fitter import (
        calculate_fit_font_size,
        get_placeholder_font_info,
        apply_uniform_font_size,
        calculate_table_fit,
        get_placeholder_dimensions,
        calculate_column_widths,
        shrink_textbox_to_fit,
        shrink_placeholder_to_fit,
    )
    _HAS_CONTENT_FITTER = True
except ImportError:
    _HAS_CONTENT_FITTER = False

# Phase 6: Graceful degradation framework
try:
    from graceful_degradation import (
        Severity,
        DegradationContext,
        DegradationResult,
    )
    _HAS_GRACEFUL_DEGRADATION = True
except ImportError:
    _HAS_GRACEFUL_DEGRADATION = False
    DegradationContext = None  # noqa: N806
    DegradationResult = None  # noqa: N806

    # Minimal fallback definitions
    class Severity:  # noqa: D101
        FATAL = "fatal"
        ERROR = "error"
        WARN = "warn"

# Phase 7: Content overflow splitting (optional dependency)
try:
    from content_splitter import (
        detect_content_overflow,
        create_continuation_slides,
    )
    _HAS_CONTENT_SPLITTER = True
except ImportError:
    _HAS_CONTENT_SPLITTER = False


# Per-slide warning collection (reset before each slide's populate call)
_current_slide_warnings: list[str] = []

# Phase 5: Content shrink warnings (reset per slide in generate_presentation)
_shrink_warnings: list[str] = []

# Phase 6: Current degradation context (set during generation)
_degradation_ctx: "DegradationContext | None" = None


def _collect_warning(msg: str, category: str = "general", severity: str = "warn"):
    """Collect a warning for the current slide being processed.

    Args:
        msg: Warning message
        category: Issue category (image, layout, typography, table, etc.)
        severity: Severity level (warn, error, fatal)
    """
    _current_slide_warnings.append(msg)

    # Also record in degradation context if available
    if _HAS_GRACEFUL_DEGRADATION and _degradation_ctx is not None:
        sev = {
            "warn": Severity.WARN,
            "error": Severity.ERROR,
            "fatal": Severity.FATAL,
        }.get(severity, Severity.WARN)
        _degradation_ctx.add_issue(
            severity=sev,
            category=category,
            message=msg,
        )


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

class LayoutPlanError(Exception):
    """Raised when layout plan validation fails."""
    pass


def load_layout_plan(path: str) -> dict:
    """Load and validate layout plan JSON using Pydantic schema.

    Raises:
        LayoutPlanError: If the layout plan fails schema validation
        FileNotFoundError: If the file doesn't exist
    """
    plan_path = Path(path)
    if not plan_path.exists():
        raise FileNotFoundError(f"Layout plan not found: {path}")

    try:
        # Validate using Pydantic schema
        validated_plan = LayoutPlan.model_validate_json(plan_path.read_text(encoding='utf-8'))
        # Return as dict for compatibility with existing code
        return validated_plan.model_dump(by_alias=True)
    except ValidationError as e:
        # Format Pydantic errors into human-readable message
        error_lines = ["Layout plan validation failed:"]
        for error in e.errors():
            loc = " -> ".join(str(x) for x in error["loc"])
            msg = error["msg"]
            error_lines.append(f"  {loc}: {msg}")
        raise LayoutPlanError("\n".join(error_lines))


def clean_markdown(text: str) -> str:
    """Remove markdown and style marker formatting from text."""
    if not text:
        return ''
    if not isinstance(text, str):
        return str(text)
    # Remove markdown bold/italic
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    # Remove color markers {blue}...{/blue}, {italic}...{/italic}, etc.
    text = re.sub(r'\{(blue|italic|bold|red|green)\}', '', text)
    text = re.sub(r'\{/(blue|italic|bold|red|green)\}', '', text)
    return text.strip()


def parse_markdown_formatting(text: str) -> list:
    """Parse markdown text into a list of formatted segments.

    Returns list of dicts with:
    - text: the text content
    - bold: True if wrapped in **
    - italic: True if wrapped in *

    Example: "**Consumer:** Premium-forward" ->
    [{'text': 'Consumer:', 'bold': True}, {'text': ' Premium-forward', 'bold': False}]
    """
    if not text:
        return []
    if not isinstance(text, str):
        return [{'text': str(text), 'bold': False, 'italic': False}]

    segments = []
    # Pattern to match **bold**, *italic*, or plain text
    # Order matters: check ** before *
    pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|[^*]+)'
    parts = re.findall(pattern, text)

    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            # Bold text
            segments.append({
                'text': part[2:-2],
                'bold': True,
                'italic': False
            })
        elif part.startswith('*') and part.endswith('*'):
            # Italic text
            segments.append({
                'text': part[1:-1],
                'bold': False,
                'italic': True
            })
        else:
            # Plain text
            segments.append({
                'text': part,
                'bold': False,
                'italic': False
            })

    return segments


def set_text_with_formatting(text_frame, text: str, base_bold: bool = False, base_italic: bool = False):
    """Set text in a text_frame, preserving markdown formatting.

    Parses **bold** and *italic* markers and applies run-level formatting.

    Args:
        text_frame: PowerPoint text frame object
        text: Text with optional markdown formatting
        base_bold: Default bold state for non-marked text
        base_italic: Default italic state for non-marked text
    """
    if not text:
        return

    p = text_frame.paragraphs[0]
    p.clear()

    segments = parse_markdown_formatting(text)

    for seg in segments:
        run = p.add_run()
        run.text = seg['text']
        run.font.bold = seg['bold'] or base_bold
        run.font.italic = seg['italic'] or base_italic


def apply_style_to_run(run, style: dict):
    """Apply a style object to a PowerPoint run.

    Style object can contain:
    - color: hex color code without # (e.g., "0196FF")
    - size: font size in points (e.g., 14)
    - font: font family name (e.g., "Aptos")
    - italic: boolean
    - bold: boolean
    - underline: boolean
    - strike: boolean (strikethrough)
    - superscript: boolean
    - subscript: boolean
    - caps: boolean (all caps)

    Args:
        run: PowerPoint run object
        style: Dict with style properties
    """
    if not style:
        return

    # Color
    if style.get('color'):
        color_hex = style['color'].lstrip('#')
        if len(color_hex) != 6:
            _collect_warning(f"Invalid hex color format '{style['color']}' (expected 6 digits)")
        else:
            try:
                r = int(color_hex[0:2], 16)
                g = int(color_hex[2:4], 16)
                b = int(color_hex[4:6], 16)
                run.font.color.rgb = RGBColor(r, g, b)
            except ValueError:
                _collect_warning(f"Invalid hex color value '{style['color']}'")

    # Size
    if style.get('size'):
        try:
            size_val = int(style['size'])
            if size_val <= 0 or size_val > 1000:
                _collect_warning(f"Font size {size_val} out of range (1-1000)")
            else:
                run.font.size = Pt(size_val)
        except (ValueError, TypeError):
            _collect_warning(f"Invalid font size '{style['size']}'")

    # Font family (apply fallback chain for user-specified fonts)
    if style.get('font'):
        try:
            from font_fallback import resolve_font_for_pptx
            run.font.name = resolve_font_for_pptx(style['font'])
        except ImportError:
            run.font.name = style['font']

    # Bold
    if style.get('bold'):
        run.font.bold = True

    # Italic
    if style.get('italic'):
        run.font.italic = True

    # Underline
    if style.get('underline'):
        run.font.underline = True

    # Strikethrough
    if style.get('strike'):
        # python-pptx doesn't have direct strikethrough; use XML
        rPr = get_run_properties(run)
        strike_el = rPr.find(qn('a:strike'))
        if strike_el is None:
            strike_el = etree.SubElement(rPr, qn('a:strike'))
        strike_el.set('val', 'sngStrike')

    # Superscript (baseline offset +30% and reduced font size)
    if style.get('superscript'):
        rPr = get_run_properties(run)
        rPr.set('baseline', '30000')  # 30% above baseline
        # Reduce font size to ~60% for readability
        # If run already has explicit size, reduce it; otherwise use reduced default
        current_size = run.font.size
        if current_size:
            reduced_size = int(current_size.pt * 0.60)
            run.font.size = Pt(max(reduced_size, 6))  # Minimum 6pt
        else:
            # Default body font is ~16pt, so super/subscript should be ~10pt
            run.font.size = Pt(10)

    # Subscript (baseline offset -30% and reduced font size)
    if style.get('subscript'):
        rPr = get_run_properties(run)
        rPr.set('baseline', '-30000')  # 30% below baseline
        # Reduce font size to ~60% for readability
        current_size = run.font.size
        if current_size:
            reduced_size = int(current_size.pt * 0.60)
            run.font.size = Pt(max(reduced_size, 6))  # Minimum 6pt
        else:
            # Default body font is ~16pt, so super/subscript should be ~10pt
            run.font.size = Pt(10)

    # All caps
    if style.get('caps'):
        rPr = get_run_properties(run)
        caps_el = rPr.find(qn('a:caps'))
        if caps_el is None:
            caps_el = etree.SubElement(rPr, qn('a:caps'))
        caps_el.set('val', '1')


def apply_paragraph_format(paragraph, para_format: dict):
    """Apply paragraph-level formatting to a PowerPoint paragraph.

    Supports:
    - bullet: Bullet character ('-', '•', etc.), 'numbered', or None (no bullet)
    - level: Indent level (0-4)
    - space_before: Space before in points
    - space_after: Space after in points

    Args:
        paragraph: PowerPoint paragraph object
        para_format: Dict with paragraph formatting options
    """
    if not para_format:
        return

    # Get or create pPr (paragraph properties)
    pPr = get_paragraph_properties(paragraph)

    # Bullet character
    bullet = para_format.get('bullet')
    if bullet is not None:
        if bullet is None or bullet == 'none':
            # Remove bullet - set buNone
            # Remove any existing bullet elements
            for bu_el in pPr.findall(qn('a:buChar')) + pPr.findall(qn('a:buAutoNum')):
                pPr.remove(bu_el)
            buNone = pPr.find(qn('a:buNone'))
            if buNone is None:
                buNone = etree.SubElement(pPr, qn('a:buNone'))
        elif bullet == 'numbered':
            # Numbered list - use buAutoNum
            for bu_el in pPr.findall(qn('a:buChar')) + pPr.findall(qn('a:buNone')):
                pPr.remove(bu_el)
            buAutoNum = pPr.find(qn('a:buAutoNum'))
            if buAutoNum is None:
                buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
            buAutoNum.set('type', 'arabicPeriod')  # 1. 2. 3. etc.
        else:
            # Custom bullet character (-, •, ▸, etc.)
            for bu_el in pPr.findall(qn('a:buAutoNum')) + pPr.findall(qn('a:buNone')):
                pPr.remove(bu_el)
            buChar = pPr.find(qn('a:buChar'))
            if buChar is None:
                buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', bullet)

    # Indent level
    if 'level' in para_format:
        paragraph.level = min(4, max(0, int(para_format['level'])))

    # Space before (in points -> EMUs: 1pt = 12700 EMUs)
    if para_format.get('space_before') is not None:
        spcBef = pPr.find(qn('a:spcBef'))
        if spcBef is None:
            spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        # Clear existing children
        for child in list(spcBef):
            spcBef.remove(child)
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        # spcPts val is in hundredths of a point
        spcPts.set('val', str(int(para_format['space_before']) * 100))

    # Space after (in points -> hundredths of a point)
    if para_format.get('space_after') is not None:
        spcAft = pPr.find(qn('a:spcAft'))
        if spcAft is None:
            spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        # Clear existing children
        for child in list(spcAft):
            spcAft.remove(child)
        spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', str(int(para_format['space_after']) * 100))


def resolve_formatting(template_format: dict, layout_plan_format: dict | None, slide_num: int = 0) -> tuple[dict, list[str]]:
    """Merge template defaults with layout plan overrides.

    Layers layout plan formatting on top of template defaults, issuing
    consistency warnings when overrides change existing template values.
    Also caps bullet depth at 3 levels (0-indexed max = 2).

    Args:
        template_format: Base formatting from template (e.g. {'size': 16, 'bold': True})
        layout_plan_format: Overrides from layout plan (or None for no overrides)
        slide_num: Slide number for warning context

    Returns:
        Tuple of (merged_format_dict, list_of_warning_strings)
    """
    warnings = []
    merged = dict(template_format)

    if layout_plan_format:
        for key, plan_val in layout_plan_format.items():
            if key in merged and merged[key] != plan_val:
                template_val = merged[key]
                warnings.append(
                    f"Slide {slide_num}: Layout plan overrides template {key} ({template_val} -> {plan_val})"
                )
            merged[key] = plan_val

    # Cap bullet depth at 3 levels (0-indexed: 0, 1, 2)
    if 'level' in merged:
        level = int(merged['level'])
        if level > 2:
            warnings.append(
                f"Slide {slide_num}: Bullet level {level} capped at 2 (max 3 levels)"
            )
            merged['level'] = 2

    return merged, warnings


def route_content_to_placeholder(placeholders: list, content_items: list[dict], slide_num: int = 0) -> tuple[list[dict], list[str]]:
    """Map content items to discovered placeholders by type and index.

    Routes each content item to a placeholder based on its 'type' field
    (e.g. 'TITLE', 'BODY') and optional 'index' field. Produces warnings
    for unmatched content items and groups multiple items targeting the
    same placeholder.

    Args:
        placeholders: List of placeholder dicts from discover_placeholders()
            Each has 'type' (or 'ph_type'), 'idx', 'shape', 'position' keys
        content_items: List of content dicts with at minimum {'type': str, 'text': str}
            Optional 'index' key specifies target placeholder index
        slide_num: Slide number for warning context

    Returns:
        Tuple of (routed_list, warnings) where routed_list entries have:
            - placeholder: the placeholder dict (or None if unmatched)
            - content_items: list of content dicts routed to this placeholder
            - matched: bool indicating if a placeholder was found
    """
    warnings = []

    # Build lookup: group placeholders by type
    ph_by_type: dict[str, list[dict]] = {}
    for ph in placeholders:
        # Support both 'ph_type' (test mocks) and 'type' (discover_placeholders output)
        ph_type = ph.get('ph_type') or ph.get('type', 'OTHER')
        if ph_type not in ph_by_type:
            ph_by_type[ph_type] = []
        ph_by_type[ph_type].append(ph)

    # Route each content item to a placeholder
    # Group by (type, index) to handle multiple items for same placeholder
    routing_groups: dict[tuple[str, int], dict] = {}

    for item in content_items:
        target_type = item.get('type', '')
        target_index = item.get('index', 0)
        group_key = (target_type, target_index)

        if group_key not in routing_groups:
            # Find matching placeholder
            matched_ph = None
            type_phs = ph_by_type.get(target_type, [])

            for ph in type_phs:
                ph_idx = ph.get('idx', 0)
                if ph_idx == target_index:
                    matched_ph = ph
                    break

            # If exact index not found but there are placeholders of this type,
            # try the first one (index 0 is common default)
            if matched_ph is None and type_phs and target_index == 0:
                matched_ph = type_phs[0]

            if matched_ph is None:
                warnings.append(
                    f"Slide {slide_num}: No placeholder of type '{target_type}' index {target_index} found, content skipped"
                )

            routing_groups[group_key] = {
                'placeholder': matched_ph,
                'content_items': [],
                'matched': matched_ph is not None,
            }

        routing_groups[group_key]['content_items'].append(item)

    return list(routing_groups.values()), warnings


def render_styled_runs(paragraph, styled_runs: list, base_bold: bool = False, base_italic: bool = False):
    """Render a list of styled runs to a PowerPoint paragraph.

    Styled runs format (from slide-outline-to-layout):
    [
        {"text": "The opportunity: ", "style": null},
        {"text": "premium indulgence", "style": {"color": "0196FF"}},
        {"text": " in solo moments.", "style": null}
    ]

    Also parses markdown bold/italic (**bold**, *italic*) within each run's text.

    Args:
        paragraph: PowerPoint paragraph object (should be cleared first)
        styled_runs: List of {text, style} dicts
        base_bold: Default bold for unstyled text
        base_italic: Default italic for unstyled text
    """
    if not styled_runs:
        return

    for run_data in styled_runs:
        text = run_data.get('text', '')
        style = run_data.get('style')

        # Parse markdown formatting within this run's text
        segments = parse_markdown_formatting(text)

        for seg in segments:
            run = paragraph.add_run()
            run.text = seg['text']

            # Always set font name to ensure consistency after tf.clear()
            run.font.name = get_brand_font('body')

            # Apply base formatting (from function args)
            if base_bold or seg['bold']:
                run.font.bold = True
            if base_italic or seg['italic']:
                run.font.italic = True

            # Apply style overrides (color, size, etc. from styled run)
            if style:
                apply_style_to_run(run, style)


def add_styled_paragraph(text_frame, text: str, styled_runs: list = None, level: int = 0,
                         base_bold: bool = False, base_italic: bool = False, first_para: bool = False,
                         para_format: dict = None):
    """Add a paragraph with styled runs or fallback to markdown formatting.

    If styled_runs is provided, uses those for formatting.
    Otherwise, falls back to markdown parsing (**bold**, *italic*).

    Args:
        text_frame: PowerPoint text frame
        text: Plain text (fallback if styled_runs is None)
        styled_runs: Optional list of {text, style} dicts
        level: Paragraph indent level
        base_bold: Default bold for unstyled text
        base_italic: Default italic for unstyled text
        first_para: If True, use existing first paragraph
        para_format: Optional dict with paragraph formatting (bullet, level, space_before, space_after)

    Returns:
        The paragraph object created/used
    """
    if first_para:
        p = text_frame.paragraphs[0]
        p.clear()
    else:
        p = text_frame.add_paragraph()

    # Extract paragraph-level markers from text (e.g., {level:N}, {bullet:X})
    # This strips markers like {level:0} from the text and extracts them as formatting
    clean_text = text
    extracted_para_format = {}
    if text and not styled_runs:
        clean_text, extracted_para_format = extract_paragraph_markers(text)

    # Merge extracted paragraph format with passed para_format (passed takes precedence)
    merged_para_format = {**extracted_para_format, **(para_format or {})}

    # Apply paragraph-level formatting
    if merged_para_format:
        # Level from para_format overrides the level argument
        if 'level' in merged_para_format:
            p.level = min(4, max(0, int(merged_para_format['level'])))
        else:
            p.level = level
        # Apply bullet and spacing
        apply_paragraph_format(p, merged_para_format)
    else:
        p.level = level

    if styled_runs:
        # Use styled runs
        render_styled_runs(p, styled_runs, base_bold, base_italic)
    elif clean_text and has_typography_markers(clean_text):
        # Fall back to typography marker parsing for {bold}, {blue}, etc.
        add_styled_paragraph_with_markers(p, clean_text, base_bold, base_italic)
    else:
        # Fall back to markdown parsing
        segments = parse_markdown_formatting(clean_text)
        brand_font = get_brand_font('body')
        for seg in segments:
            run = p.add_run()
            run.text = seg['text']
            # Always set font name to ensure consistency after tf.clear()
            run.font.name = brand_font
            # Only set bold/italic if True - leave as None (inherit) otherwise
            if seg['bold'] or base_bold:
                run.font.bold = True
            if seg['italic'] or base_italic:
                run.font.italic = True

    return p


def parse_nested_typography_markers(text: str, inherited_style: dict = None) -> list:
    """Recursively parse nested typography markers into flat list of styled parts.

    Handles nested markers like {bold}{blue}text{/blue}{/bold} by recursively
    parsing content and merging styles from outer and inner markers.

    Args:
        text: Text with inline style markers
        inherited_style: Style dict inherited from outer markers

    Returns:
        List of dicts with 'text' and 'style' keys, where text is clean
        and style is a merged dict of all applicable styles.
    """
    if not text:
        return []

    inherited_style = inherited_style or {}

    # Get color map and preset styles
    brand_colors = get_color_map()
    color_map = {
        'blue': brand_colors.get('blue', '0196FF'),
        'red': 'FF0000',
        'green': '00AA00',
    }
    preset_styles = get_preset_styles()

    # Pattern for basic markers: {style}content{/style}
    basic_pattern = r'\{(blue|red|green|italic|bold|underline|strike|signpost|question|caps|super|sub)\}(.*?)\{/\1\}'
    # Pattern for custom color: {color:#RRGGBB}content{/color}
    color_pattern = r'\{color:#([0-9A-Fa-f]{6})\}(.*?)\{/color\}'
    # Pattern for size: {size:N}content{/size}
    size_pattern = r'\{size:(\d+)\}(.*?)\{/size\}'
    # Pattern for font: {font:name}content{/font}
    font_pattern = r'\{font:([^}]+)\}(.*?)\{/font\}'

    # Find all matches
    all_matches = []

    for match in re.finditer(basic_pattern, text):
        all_matches.append({
            'start': match.start(),
            'end': match.end(),
            'marker': match.group(1),
            'content': match.group(2),
            'type': 'basic'
        })

    for match in re.finditer(color_pattern, text):
        all_matches.append({
            'start': match.start(),
            'end': match.end(),
            'marker': match.group(1),  # Hex color
            'content': match.group(2),
            'type': 'custom_color'
        })

    for match in re.finditer(size_pattern, text):
        all_matches.append({
            'start': match.start(),
            'end': match.end(),
            'marker': match.group(1),  # Size in points
            'content': match.group(2),
            'type': 'size'
        })

    for match in re.finditer(font_pattern, text):
        all_matches.append({
            'start': match.start(),
            'end': match.end(),
            'marker': match.group(1),  # Font name
            'content': match.group(2),
            'type': 'font'
        })

    # Sort by start position
    all_matches.sort(key=lambda x: x['start'])

    # If no matches, return text with inherited style
    if not all_matches:
        clean_text = clean_markdown(text)
        if clean_text:
            return [{'text': clean_text, 'style': inherited_style if inherited_style else None}]
        return []

    # Build parts list with recursive parsing
    parts = []
    last_end = 0

    for match in all_matches:
        # Skip overlapping matches (can happen with nested patterns)
        if match['start'] < last_end:
            continue

        # Add text before this marker (with inherited style)
        # Note: Don't strip intermediate text - preserve whitespace between parts
        if match['start'] > last_end:
            before_text = text[last_end:match['start']]
            # Only remove markdown formatting, don't strip whitespace
            clean_before = re.sub(r'\*\*([^*]+)\*\*', r'\1', before_text)
            clean_before = re.sub(r'\*([^*]+)\*', r'\1', clean_before)
            if clean_before:
                parts.append({'text': clean_before, 'style': inherited_style if inherited_style else None})

        # Build style for this marker
        marker_style = dict(inherited_style) if inherited_style else {}

        if match['type'] == 'custom_color':
            marker_style['color'] = match['marker']
        elif match['type'] == 'size':
            marker_style['size'] = int(match['marker'])
        elif match['type'] == 'font':
            marker_style['font'] = match['marker']
        elif match['marker'] in color_map:
            marker_style['color'] = color_map[match['marker']]
        elif match['marker'] in preset_styles:
            marker_style.update(preset_styles[match['marker']])
        elif match['marker'] == 'italic':
            marker_style['italic'] = True
        elif match['marker'] == 'bold':
            marker_style['bold'] = True
        elif match['marker'] == 'underline':
            marker_style['underline'] = True
        elif match['marker'] == 'strike':
            marker_style['strike'] = True
        elif match['marker'] == 'caps':
            marker_style['caps'] = True
        elif match['marker'] == 'super':
            marker_style['superscript'] = True
        elif match['marker'] == 'sub':
            marker_style['subscript'] = True

        # Recursively parse the content with merged style
        content_parts = parse_nested_typography_markers(match['content'], marker_style)
        parts.extend(content_parts)

        last_end = match['end']

    # Add remaining text after last match (with inherited style)
    # Strip trailing whitespace only for the final segment
    if last_end < len(text):
        after_text = text[last_end:]
        # Only remove markdown formatting, then strip trailing whitespace
        clean_after = re.sub(r'\*\*([^*]+)\*\*', r'\1', after_text)
        clean_after = re.sub(r'\*([^*]+)\*', r'\1', clean_after).rstrip()
        if clean_after:
            parts.append({'text': clean_after, 'style': inherited_style if inherited_style else None})

    return parts


def add_styled_paragraph_with_markers(paragraph, text: str, base_bold: bool = False, base_italic: bool = False):
    """Parse inline style markers and render with formatting.

    FALLBACK FUNCTION: Used when pre-parsed styled runs aren't available.
    For new outlines, the ingest.py parser creates _styled fields with full
    typography support. This function provides backwards compatibility.

    Supported markers:
    - Color: {blue}text{/blue}, {red}text{/red}, {green}text{/green}
    - Custom color: {color:#RRGGBB}text{/color}
    - Style: {bold}text{/bold}, {italic}text{/italic}
    - Other: {underline}text{/underline}, {strike}text{/strike}
    - Preset: {signpost}text{/signpost}, {question}text{/question}
    - Nested: {bold}{blue}text{/blue}{/bold} - markers can be nested

    Args:
        paragraph: PowerPoint paragraph object
        text: Text with inline style markers
        base_bold: Default bold state
        base_italic: Default italic state
    """
    # Parse nested markers recursively
    parts = parse_nested_typography_markers(text)

    # If no parts (empty text), create one part with the original text
    if not parts:
        parts = [{'text': clean_markdown(text), 'style': None}]

    # Render parts
    brand_font = get_brand_font('body')
    for part in parts:
        if not part['text']:
            continue
        run = paragraph.add_run()
        run.text = part['text']

        # Always set font name to ensure consistency after tf.clear()
        run.font.name = brand_font

        # Apply base styles
        if base_bold:
            run.font.bold = True
        if base_italic:
            run.font.italic = True

        # Apply part-specific styles
        if part['style']:
            apply_style_to_run(run, part['style'])


def is_deliverable_object(item) -> bool:
    """Check if an item is a deliverable object (has title/description keys)."""
    if not isinstance(item, dict):
        return False
    return 'title' in item or 'description' in item


def extract_deliverables(slide_data: dict) -> list:
    """Extract deliverables from any supported location in slide data.
    
    Checks (in order):
    1. slide_data['deliverables']
    2. slide_data['content']['deliverables']
    3. slide_data['content']['body'] if it contains objects with title/description
    4. slide_data['cards'] (converts card format to deliverable format)
    """
    # Pattern A: Root-level deliverables array
    if slide_data.get('deliverables'):
        return slide_data['deliverables']
    
    content = slide_data.get('content', {})
    
    # Pattern B: Inside content
    if content.get('deliverables'):
        return content['deliverables']
    
    # Pattern C: Body contains objects with title/description
    body = content.get('body', [])
    if body and isinstance(body, list) and len(body) > 0:
        if is_deliverable_object(body[0]):
            return body
    
    # Pattern D: Cards array (convert to deliverable format)
    cards = slide_data.get('cards', [])
    if cards:
        deliverables = []
        for card in cards:
            deliverables.append({
                'number': card.get('number'),
                'title': card.get('title', ''),
                'description': card.get('body', '')  # cards use 'body', deliverables use 'description'
            })
        return deliverables
    
    return []


def extract_timeline(slide_data: dict) -> list:
    """Extract timeline from any supported location."""
    # Root-level
    if slide_data.get('timeline'):
        return slide_data['timeline']
    
    # Inside content
    content = slide_data.get('content', {})
    if content.get('timeline'):
        return content['timeline']
    
    # Inside extras
    extras = slide_data.get('extras', {})
    if extras.get('timeline'):
        return extras['timeline']
    
    return []


def extract_tables(slide_data: dict) -> list:
    """Extract tables from any supported location.
    
    Normalizes table structure to always have 'headers' (plural) key.
    """
    tables = []
    
    # Check all possible locations
    raw_tables = []
    
    # Root-level tables array
    if slide_data.get('tables'):
        raw_tables.extend(slide_data['tables'])
    
    content = slide_data.get('content', {})
    
    # Content-level tables array
    if content.get('tables'):
        raw_tables.extend(content['tables'])
    
    # Single table in content
    if content.get('table'):
        raw_tables.append(content['table'])
    
    # Extras table
    extras = slide_data.get('extras', {})
    if extras.get('table'):
        raw_tables.append(extras['table'])
    
    # Normalize table structure
    for table in raw_tables:
        normalized = {
            'headers': [],
            'rows': table.get('rows', []),
            'section_header': None
        }
        
        # Handle 'headers' (plural) - column headers
        if table.get('headers'):
            normalized['headers'] = table['headers']
        
        # Handle 'header' (singular) - section title, not column header
        if table.get('header') and not table.get('headers'):
            normalized['section_header'] = table['header']
            # If no column headers, infer from first row length
            if normalized['rows'] and len(normalized['rows']) > 0:
                # Use generic headers or leave empty for 2-column key-value style
                pass
        
        tables.append(normalized)
    
    return tables


def extract_columns(slide_data: dict) -> list:
    """Extract columns from slide data.

    Also converts cards to column format if columns aren't present.
    Cards have {number, title, body} -> Columns have {header, body}

    Fallback: if no columns or cards, check extras.table - this handles the case
    where the parser didn't convert a table to columns for a column-based layout.
    """
    columns = slide_data.get('columns', [])
    if columns:
        return columns

    # Convert cards to column format
    cards = slide_data.get('cards', [])
    if cards:
        return [
            {
                'header': card.get('title', ''),
                'body': [card.get('body', '')] if card.get('body') else [],
                'number': card.get('number')
            }
            for card in cards
        ]

    # Fallback: convert extras.table to columns if present
    # This provides defense-in-depth for column-based visual types with table content
    extras = slide_data.get('extras', {})
    table = extras.get('table')
    if table and table.get('headers'):
        return _table_to_columns(table)

    return []


def _table_to_columns(table: dict) -> list:
    """Convert a table structure to columns format.

    Used as fallback when parser didn't convert table to columns.
    """
    columns = []
    headers = table.get('headers', [])
    rows = table.get('rows', [])

    for col_idx, header in enumerate(headers):
        column = {
            'header': header.strip() if header else '',
            'body': [],
            'number': col_idx + 1
        }

        # Collect all row data for this column
        for row in rows:
            # Handle section header rows (dict format)
            if isinstance(row, dict):
                if col_idx == 0 and row.get('type') == 'section_header':
                    column['body'].append(row.get('text', ''))
                continue

            # Regular row (list format)
            if isinstance(row, list) and col_idx < len(row):
                cell_value = row[col_idx]
                if cell_value and cell_value.strip():
                    column['body'].append(cell_value.strip())

        columns.append(column)

    return columns


# =============================================================================
# PLACEHOLDER DISCOVERY
# =============================================================================

def discover_placeholders(slide) -> list:
    """Discover all placeholders in a slide at runtime.

    Returns list of dicts with placeholder info:
    - idx: placeholder index
    - type: placeholder type name (TITLE, BODY, PICTURE, etc.)
    - position: dict with left, top, width, height in inches
    """
    placeholders = []

    type_names = {
        PP_PLACEHOLDER.TITLE: 'TITLE',
        PP_PLACEHOLDER.CENTER_TITLE: 'CENTER_TITLE',
        PP_PLACEHOLDER.SUBTITLE: 'SUBTITLE',
        PP_PLACEHOLDER.BODY: 'BODY',
        PP_PLACEHOLDER.PICTURE: 'PICTURE',
        PP_PLACEHOLDER.SLIDE_NUMBER: 'SLIDE_NUMBER',
        PP_PLACEHOLDER.FOOTER: 'FOOTER',
        PP_PLACEHOLDER.DATE: 'DATE',
        PP_PLACEHOLDER.OBJECT: 'OBJECT',
        PP_PLACEHOLDER.CHART: 'CHART',
        PP_PLACEHOLDER.TABLE: 'TABLE',
    }

    for shape in slide.placeholders:
        try:
            ph_format = shape.placeholder_format
            if ph_format is None:
                _collect_warning("Placeholder shape has no format, skipping")
                continue

            ph_type = type_names.get(ph_format.type, 'OTHER')
            idx = getattr(ph_format, 'idx', None)
            if idx is None:
                _collect_warning("Placeholder has no idx attribute, skipping")
                continue

            placeholders.append({
                'idx': idx,
                'type': ph_type,
                'shape': shape,
                'position': {
                    'left_inches': shape.left / Emu(914400) if shape.left else 0,
                    'top_inches': shape.top / Emu(914400) if shape.top else 0,
                    'width_inches': shape.width / Emu(914400) if shape.width else 0,
                    'height_inches': shape.height / Emu(914400) if shape.height else 0,
                }
            })
        except AttributeError as e:
            _collect_warning(f"Error accessing placeholder attributes: {e}")
            continue

    return placeholders


def get_placeholder_by_idx(slide, idx: int, slide_num: int = 0, ph_type: str = None, warn_on_missing: bool = True):
    """Get placeholder shape by idx with optional warning on missing.

    Args:
        slide: PowerPoint slide object
        idx: Placeholder index to find
        slide_num: Slide number for warning messages (0 means don't include in message)
        ph_type: Optional placeholder type name for more descriptive warnings
        warn_on_missing: Whether to emit a warning if placeholder not found

    Returns:
        Placeholder shape if found, None otherwise
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape

    # Placeholder not found - emit warning if requested
    if warn_on_missing:
        type_info = f" ({ph_type})" if ph_type else ""
        slide_info = f"Slide {slide_num}: " if slide_num else ""
        _collect_warning(
            f"{slide_info}Placeholder idx {idx}{type_info} not found in layout - content skipped"
        )
    return None


def sort_placeholders_by_position(placeholders: list) -> list:
    """Sort placeholders by position: top-to-bottom, left-to-right."""
    return sorted(placeholders, key=lambda p: (
        p.get('position', {}).get('top_inches', 0),
        p.get('position', {}).get('left_inches', 0)
    ))


def get_largest_body_placeholder(placeholders: list) -> dict | None:
    """Find the largest BODY placeholder by area."""
    body_phs = [p for p in placeholders if p.get('type') == 'BODY']
    if not body_phs:
        return None
    
    def area(p):
        pos = p.get('position', {})
        return pos.get('width_inches', 0) * pos.get('height_inches', 0)
    
    return max(body_phs, key=area)


def get_title_placeholder(placeholders: list) -> dict | None:
    """Get TITLE or CENTER_TITLE placeholder."""
    for p in placeholders:
        if p.get('type') in ('TITLE', 'CENTER_TITLE'):
            return p
    return None


def get_body_placeholders(placeholders: list, sort_horizontal: bool = False) -> list:
    """Get all BODY placeholders, optionally sorted left-to-right."""
    body_phs = [p for p in placeholders if p.get('type') == 'BODY']
    if sort_horizontal:
        body_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))
    return body_phs


def set_slide_number(slide, slide_num: int) -> bool:
    """Set slide number by adding a text box at the SLIDE_NUMBER placeholder position.

    Looks up the SLIDE_NUMBER placeholder position from the slide's layout and
    creates a text box at that position with the slide number.

    Args:
        slide: PowerPoint slide object
        slide_num: The slide number to display

    Returns:
        True if slide number was set, False if no placeholder position found
    """
    # Find SLIDE_NUMBER placeholder in the slide's layout
    sn_ph = None
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            sn_ph = ph
            break

    if not sn_ph:
        return False

    # Create text box at the placeholder position
    textbox = slide.shapes.add_textbox(sn_ph.left, sn_ph.top, sn_ph.width, sn_ph.height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(slide_num)
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x59, 0x59, 0x59)  # Gray to match typical slide numbers

    return True


# =============================================================================
# TEXT SETTING FUNCTIONS
# =============================================================================

def set_text_in_placeholder(placeholder, text: str, bold: bool = False, font_size: int = None,
                            styled_runs: list = None, slide_num: int = 0,
                            style_overrides: dict = None):
    """Set text in a placeholder with optional formatting and shrink-to-fit.

    Args:
        placeholder: PowerPoint placeholder to fill
        text: Plain text to display (fallback if no styled_runs)
        bold: Whether to make text bold
        font_size: Optional font size in points
        styled_runs: Optional list of styled runs [{text, style}, ...]
                     where style can have 'color' (hex), 'bold', 'italic', 'size'
        slide_num: Slide number for warning context (Phase 5)
        style_overrides: Optional layout plan style overrides for resolve_formatting
    """
    if placeholder is None or (not text and not styled_runs):
        return False

    # Phase 5: Read font info BEFORE clearing text frame (tf.clear() erases defaults)
    font_info = None
    if _HAS_CONTENT_FITTER:
        try:
            font_info = get_placeholder_font_info(placeholder)
        except Exception:
            pass

    tf = placeholder.text_frame
    tf.clear()

    # Phase 5: Resolve formatting — merge template defaults with layout plan overrides
    template_format = {}
    if font_info:
        template_format = {
            'size': font_info.get('font_size_pt'),
            'font_name': font_info.get('font_name', get_brand_font('header')),
        }
        if bold:
            template_format['bold'] = True
    merged_format, format_warnings = resolve_formatting(template_format, style_overrides, slide_num)
    if format_warnings:
        _shrink_warnings.extend(format_warnings)

    p = tf.paragraphs[0]

    # Use styled runs if provided
    if styled_runs:
        for i, run_info in enumerate(styled_runs):
            # Add soft return (line break) between styled segments for multi-line titles
            if i > 0:
                p.add_line_break()

            run = p.add_run()
            run_text = run_info.get('text', '')
            run.text = clean_markdown(run_text)

            style = run_info.get('style') or {}

            # Apply style
            if style.get('color'):
                run.font.color.rgb = RGBColor.from_string(style['color'])
            if style.get('bold'):
                run.font.bold = True
            if style.get('italic'):
                run.font.italic = True
            if style.get('size'):
                run.font.size = Pt(style['size'])

            # Apply base formatting
            if bold and not style.get('bold'):
                run.font.bold = True
            if font_size and not style.get('size'):
                run.font.size = Pt(font_size)
    else:
        run = p.add_run()
        run.text = clean_markdown(text)

        if bold:
            run.font.bold = True
        if font_size:
            run.font.size = Pt(font_size)

    # Phase 5: Shrink-to-fit — measure placed text and shrink if overflow
    if _HAS_CONTENT_FITTER and font_info:
        try:
            dims = get_placeholder_dimensions(placeholder)
            intended_size = font_size or font_info.get('font_size_pt', 16)
            font_name = font_info.get('font_name', get_brand_font('body'))

            # Collect all placed text
            placed_text = '\n'.join(
                run.text for para in tf.paragraphs for run in para.runs
            )

            if placed_text.strip():
                fit_result = calculate_fit_font_size(
                    text=placed_text,
                    font_name=font_name,
                    intended_size_pt=intended_size,
                    width_inches=dims['usable_width_inches'],
                    height_inches=dims['usable_height_inches'],
                )

                if fit_result.get('shrunk'):
                    shrunk_size = fit_result['size_pt']
                    apply_uniform_font_size(tf, shrunk_size)

                    # Cap any styled run sizes at shrunk size
                    if styled_runs:
                        shrunk_pt = Pt(shrunk_size)
                        for para in tf.paragraphs:
                            for r in para.runs:
                                if r.font.size and r.font.size > shrunk_pt:
                                    r.font.size = shrunk_pt

                    # Collect warning
                    warning_text = (
                        f"Slide {slide_num}: Text shrunk from {intended_size}pt to "
                        f"{shrunk_size}pt in placeholder"
                    )
                    _shrink_warnings.append(warning_text)
                    print(f"  [shrink] {warning_text}", file=sys.stderr)
        except Exception as e:
            print(f"Content fitter: could not measure text for shrink: {e}", file=sys.stderr)

    return True


def has_typography_markers(text: str) -> bool:
    """Check if text contains typography markers that need special handling.

    Detects inline markers like {blue}, {signpost}, {question}, {bold}, etc.
    and paragraph-level markers like {bullet:-}, {bullet:1}, {level:2}.
    """
    # Inline style markers
    inline_pattern = r'\{(blue|red|green|italic|bold|underline|strike|signpost|question|caps|super|sub|font:[^}]+|size:\d+)\}'
    # Custom color marker
    color_pattern = r'\{color:#[0-9A-Fa-f]{6}\}'
    # Paragraph-level markers
    para_pattern = r'\{(bullet:[^}]+|level:\d+|space:(before|after):[^}]+)\}'

    return bool(re.search(inline_pattern, text) or
                re.search(color_pattern, text) or
                re.search(para_pattern, text))


def extract_paragraph_markers(text: str) -> tuple:
    """Extract paragraph-level markers from text and return (clean_text, para_format).

    Handles:
    - {bullet:-} -> dash bullet
    - {bullet:•} -> round bullet
    - {bullet:none} -> no bullet
    - {bullet:1} or {bullet:numbered} -> numbered list
    - {level:N} -> indent level
    - {space:before:Npt} -> space before
    - {space:after:Npt} -> space after

    Returns:
        tuple: (text with markers removed, dict of paragraph formatting)
    """
    para_format = {}
    clean_text = text

    # Extract bullet marker
    bullet_match = re.search(r'\{bullet:([^}]+)\}', clean_text)
    if bullet_match:
        bullet_val = bullet_match.group(1)
        if bullet_val == 'none':
            para_format['bullet'] = None
        elif bullet_val in ('1', 'numbered'):
            para_format['bullet'] = 'numbered'
        else:
            # Character like '-', '•', '▸'
            # Convert hyphen to en-dash for proper bullet rendering
            if bullet_val == '-':
                bullet_val = '–'  # en-dash
            para_format['bullet'] = bullet_val
        clean_text = clean_text.replace(bullet_match.group(0), '')

    # Extract level marker
    level_match = re.search(r'\{level:(\d+)\}', clean_text)
    if level_match:
        para_format['level'] = int(level_match.group(1))
        clean_text = clean_text.replace(level_match.group(0), '')

    # Extract space:before marker
    space_before_match = re.search(r'\{space:before:(\d+)pt?\}', clean_text)
    if space_before_match:
        para_format['space_before'] = int(space_before_match.group(1))
        clean_text = clean_text.replace(space_before_match.group(0), '')

    # Extract space:after marker
    space_after_match = re.search(r'\{space:after:(\d+)pt?\}', clean_text)
    if space_after_match:
        para_format['space_after'] = int(space_after_match.group(1))
        clean_text = clean_text.replace(space_after_match.group(0), '')

    return clean_text.strip(), para_format


def add_formatted_paragraph(text_frame, text: str, level: int = 0, base_bold: bool = False, base_italic: bool = False, first_para: bool = False):
    """Add a paragraph with markdown formatting preserved.

    Args:
        text_frame: PowerPoint text frame
        text: Text with optional **bold** and *italic* markers, or typography markers
        level: Paragraph indent level (0 = no indent)
        base_bold: Default bold for unmarked text
        base_italic: Default italic for unmarked text
        first_para: If True, use existing first paragraph instead of adding new

    Returns:
        The paragraph object created/used
    """
    if first_para:
        p = text_frame.paragraphs[0]
        p.clear()
    else:
        p = text_frame.add_paragraph()

    p.level = level

    # Extract paragraph-level markers (bullet style, indent level, spacing)
    clean_text, para_format = extract_paragraph_markers(text)

    # Override level if specified in markers
    if 'level' in para_format:
        p.level = para_format['level']

    # Check if text has typography markers that need special handling
    if has_typography_markers(clean_text):
        # Use the comprehensive marker parser
        add_styled_paragraph_with_markers(p, clean_text, base_bold, base_italic)
    else:
        # Standard markdown-only parsing (faster path for simple text)
        segments = parse_markdown_formatting(clean_text)
        brand_font = get_brand_font('body')
        for seg in segments:
            run = p.add_run()
            run.text = seg['text']
            # Explicitly set font to ensure template consistency after tf.clear()
            run.font.name = brand_font
            # Only set bold/italic if True - leave as None (inherit) otherwise
            # This matches source PPTX behavior where None means "inherit from theme"
            if seg['bold'] or base_bold:
                run.font.bold = True
            if seg['italic'] or base_italic:
                run.font.italic = True

    # Apply paragraph-level formatting (bullet character, spacing)
    if para_format:
        apply_paragraph_format(p, para_format)

    return p


def set_header_and_bullets(placeholder, header: str, bullets: list, subheader: str = None,
                           intro: list = None, key_values: dict = None,
                           header_styled: list = None, body_styled: list = None,
                           header_bold: bool = False, add_blank_after_header: bool = False,
                           add_blank_between_items: bool = False):
    """Set header, optional subheader (italic), optional intro paragraphs, and bullet points in a single placeholder.

    Preserves markdown formatting (**bold** and *italic*) in bullet text.
    Supports styled runs from slide-outline-to-layout parser.

    Args:
        placeholder: PowerPoint placeholder to fill
        header: Header text (fallback if header_styled not provided)
        bullets: List of bullet point strings (fallback if body_styled not provided)
        subheader: Optional italic subheader text
        intro: Optional list of plain text paragraphs before bullets
        key_values: Optional dict of bold key: value pairs to render as bullets
        header_styled: Optional styled runs for header [{text, style}, ...]
        body_styled: Optional styled runs for body [{"text": ..., "runs": [...]}, ...]
        header_bold: If True, explicitly set header to bold. If False, inherit from theme.
        add_blank_after_header: If True, add an empty paragraph after the header.
        add_blank_between_items: If True, add empty paragraphs between body items.
    """
    if placeholder is None:
        return False

    tf = placeholder.text_frame
    tf.clear()

    first_para = True

    # Add header if present (use styled runs if available)
    if header or header_styled:
        add_styled_paragraph(tf, header, header_styled, level=0, base_bold=header_bold, first_para=first_para)
        first_para = False

        # Add blank paragraph after header if requested
        if add_blank_after_header:
            p = tf.add_paragraph()
            p.level = 0
            # Clear any default text
            for run in p.runs:
                run.text = ''

    # Add subheader if present
    if subheader:
        add_formatted_paragraph(tf, subheader, level=0, base_italic=True, first_para=first_para)
        first_para = False

    # Add intro paragraphs if present (plain text before bullets)
    if intro:
        for intro_line in intro:
            if isinstance(intro_line, str) and intro_line:
                add_formatted_paragraph(tf, intro_line, level=0, first_para=first_para)
                first_para = False

    # Add key_values with bold keys (e.g., **Consumer:** Premium-forward)
    if key_values:
        for key, value in key_values.items():
            if first_para:
                p = tf.paragraphs[0]
                p.clear()
                first_para = False
            else:
                p = tf.add_paragraph()
            p.level = 0

            # Bold key
            run_key = p.add_run()
            run_key.text = f"{key}: "
            run_key.font.bold = True

            # Normal value (but parse for any markdown formatting)
            value_segments = parse_markdown_formatting(value)
            for seg in value_segments:
                run = p.add_run()
                run.text = seg['text']
                run.font.bold = seg['bold']
                run.font.italic = seg['italic']

    # Add bullets - use styled runs if available, otherwise parse markdown
    if body_styled:
        for i, item in enumerate(body_styled):
            # body_styled format: [{"text": "...", "runs": [{text, style}, ...], "paragraph": {...}}, ...]
            styled_runs = item.get('runs') if isinstance(item, dict) else None
            plain_text = item.get('text', '') if isinstance(item, dict) else str(item)
            para_format = item.get('paragraph') if isinstance(item, dict) else None
            add_styled_paragraph(tf, plain_text, styled_runs, level=1, first_para=first_para,
                                 para_format=para_format)
            first_para = False
            # Add blank line after each item if requested (including last for trailing blank)
            if add_blank_between_items:
                p = tf.add_paragraph()
                p.level = 1
    else:
        # Fallback: iterate bullets with markdown parsing
        # Supports both flat strings (level=1) and dicts with 'text'/'level' keys
        for i, bullet in enumerate(bullets):
            if isinstance(bullet, str):
                add_formatted_paragraph(tf, bullet, level=1, first_para=first_para)
                first_para = False
                # Add blank line after each item if requested (including last for trailing blank)
                if add_blank_between_items:
                    p = tf.add_paragraph()
                    p.level = 1
            elif isinstance(bullet, dict):
                # Dict format: {'text': '...', 'level': N}
                # Level 0 in body means first nested level (rendered as level 1)
                # Level 1+ means deeper nesting (rendered as level 2+)
                text = bullet.get('text', '')
                item_level = bullet.get('level', 0)
                # Body items start at paragraph level 1, nested items go deeper
                render_level = item_level + 1
                add_formatted_paragraph(tf, text, level=render_level, first_para=first_para)
                first_para = False
                if add_blank_between_items:
                    p = tf.add_paragraph()
                    p.level = render_level

    # Auto-shrink font if content overflows placeholder bounds
    if _HAS_CONTENT_FITTER:
        try:
            fit_result = shrink_placeholder_to_fit(
                placeholder,
                font_name=get_brand_font('body'),
            )
            if fit_result.get('shrunk'):
                _shrink_warnings.append(
                    f"Column content shrunk from {fit_result['original_size_pt']}pt to "
                    f"{fit_result['final_size_pt']}pt"
                )
        except Exception:
            pass  # Graceful degradation - keep original sizes on error

    return True


def set_deliverable_content(placeholder, title: str, description: str):
    """Set deliverable with title (bold) and description."""
    if placeholder is None:
        return False

    tf = placeholder.text_frame
    tf.clear()
    brand_font = get_brand_font('body')

    # Title paragraph
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_markdown(title)
    run.font.name = brand_font
    run.font.bold = True
    p.level = 0

    # Description paragraph
    if description:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = clean_markdown(description)
        run.font.name = brand_font
        p.level = 0

    # Auto-shrink font if content overflows placeholder bounds
    if _HAS_CONTENT_FITTER:
        try:
            shrink_placeholder_to_fit(placeholder, font_name=brand_font)
        except Exception:
            pass  # Graceful degradation

    return True


# =============================================================================
# TABLE CREATION
# =============================================================================

def set_cell_text_with_typography(cell, text: str, font_size: int = 11, font_color: RGBColor = None, bold: bool = False):
    """Set cell text with typography marker support.

    Handles inline typography markers like {bold}, {blue}, {italic} in cell text.
    Falls back to plain text if no markers are present.

    Args:
        cell: PowerPoint table cell
        text: Cell text (may contain typography markers)
        font_size: Font size in points
        font_color: Font color (None = inherit/black)
        bold: Base bold state
    """
    if not text:
        cell.text = ''
        return

    text_str = str(text)

    # Check for typography markers
    if has_typography_markers(text_str):
        # Clear cell and use typography rendering
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        render_text_with_typography(
            p,
            text=text_str,
            base_bold=bold,
            font_size=font_size,
            font_color=font_color
        )
    else:
        # Simple text assignment
        cell.text = clean_markdown(text_str)
        # Apply formatting to all runs
        brand_font = get_brand_font('body')
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = brand_font
                if bold:
                    run.font.bold = True
                if font_size:
                    run.font.size = Pt(font_size)
                if font_color:
                    run.font.color.rgb = font_color


def create_table_in_placeholder(slide, placeholder, table_data: dict):
    """Create a table in a placeholder area with IC brand styling.
    
    IC Table Style:
    - Header row: Black background, white text, bold
    - Data rows: White background, black text
    
    Handles:
    - tables with 'headers' (column headers)
    - tables with 'section_header' (section title above table)
    - tables without headers (key-value style)
    """
    if placeholder is None:
        return False
    
    headers = table_data.get('headers', [])
    rows = table_data.get('rows', [])
    section_header = table_data.get('section_header')
    
    if not rows:
        return False
    
    # Determine column count
    col_count = len(headers) if headers else len(rows[0]) if rows else 0
    if col_count == 0:
        return False
    
    # Row count: headers row (if any) + data rows
    row_count = (1 if headers else 0) + len(rows)
    data_row_count = len(rows)

    # Get placeholder position
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # Move placeholder off-slide to hide it
    placeholder.left = Inches(-100)
    placeholder.top = Inches(-100)

    # If there's a section header, add it as a text box above the table
    table_top = top
    if section_header:
        sec_header_height = Inches(0.4)
        textbox = slide.shapes.add_textbox(left, top, width, sec_header_height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = clean_markdown(section_header)
        run.font.bold = True
        run.font.size = Pt(11)
        table_top = top + sec_header_height
        height = height - sec_header_height

    # Phase 5: Calculate table fit — determine font size and row heights
    intended_data_font = 10
    table_font_size = intended_data_font
    table_row_height_emu = None
    table_header_height_emu = None

    if _HAS_CONTENT_FITTER:
        try:
            available_height_inches = height / 914400  # EMU to inches
            fit_result = calculate_table_fit(
                row_count=data_row_count,
                col_count=col_count,
                available_height_inches=available_height_inches,
                intended_font_size_pt=intended_data_font,
            )
            table_font_size = fit_result['font_size_pt']
            table_row_height_emu = int(fit_result['row_height_inches'] * 914400)
            table_header_height_emu = int(fit_result['header_height_inches'] * 914400)
            if not fit_result['fits']:
                _shrink_warnings.append(
                    f"Table with {data_row_count} rows may overflow at minimum font size {table_font_size}pt"
                )
                print(f"  [shrink] Table with {data_row_count} rows may overflow at minimum font size {table_font_size}pt", file=sys.stderr)
        except Exception as e:
            print(f"Content fitter: could not calculate table fit: {e}", file=sys.stderr)

    # Add table shape
    table_shape = slide.shapes.add_table(row_count, col_count, left, table_top, width, height)
    table = table_shape.table

    # Disable table banding/alternating row styling
    # This prevents PowerPoint's default colored row banding
    tbl = get_table_from_shape(table_shape)
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = tbl.makeelement(qn('a:tblPr'))
        tbl.insert(0, tblPr)
    tblPr.set('bandRow', '0')
    tblPr.set('bandCol', '0')
    tblPr.set('firstRow', '0')
    tblPr.set('lastRow', '0')
    tblPr.set('firstCol', '0')
    tblPr.set('lastCol', '0')

    # Calculate column widths based on content (or fall back to equal widths)
    header_font_size = min(table_font_size + 1, intended_data_font)
    if _HAS_CONTENT_FITTER:
        try:
            col_widths = calculate_column_widths(
                headers=headers,
                rows=rows,
                total_width_emu=width,
                font_size_pt=table_font_size,
                header_font_size_pt=header_font_size,
            )
            for i, col_w in enumerate(col_widths):
                table.columns[i].width = col_w
        except Exception as e:
            # Fallback to equal widths on error
            print(f"Content fitter: column width calculation failed: {e}", file=sys.stderr)
            col_width = width // col_count
            for i in range(col_count):
                table.columns[i].width = col_width
    else:
        # Fallback: equal column widths
        col_width = width // col_count
        for i in range(col_count):
            table.columns[i].width = col_width

    # Phase 5: Set row heights from calculated fit
    if table_header_height_emu and headers:
        table.rows[0].height = table_header_height_emu
    if table_row_height_emu:
        start_row = 1 if headers else 0
        for r_idx in range(start_row, row_count):
            table.rows[r_idx].height = table_row_height_emu

    # Fill header row with IC styling (black bg, white text)
    row_idx = 0
    if headers:
        for col_idx, header in enumerate(headers):
            if col_idx < col_count:
                cell = table.cell(row_idx, col_idx)

                # Black background
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 0, 0)

                # White bold text with typography support
                set_cell_text_with_typography(
                    cell, str(header),
                    font_size=header_font_size,
                    font_color=RGBColor(255, 255, 255),
                    bold=True
                )
        row_idx += 1

    # Fill data rows with IC styling (white bg, black text)
    for row_data in rows:
        # Handle section header rows (dict with 'type': 'section_header')
        if isinstance(row_data, dict) and row_data.get('type') == 'section_header':
            # Section header spans all columns
            cell = table.cell(row_idx, 0)
            cell.text = clean_markdown(str(row_data.get('text', '')))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray for section headers
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.size = Pt(table_font_size)
            # Fill remaining cells with empty text (same row)
            for col_idx in range(1, col_count):
                cell = table.cell(row_idx, col_idx)
                cell.text = ''
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            row_idx += 1
            continue

        # Regular row - pad with empty strings if needed
        padded_row = list(row_data) + [''] * (col_count - len(row_data)) if len(row_data) < col_count else row_data

        for col_idx in range(col_count):
            cell_value = padded_row[col_idx] if col_idx < len(padded_row) else ''
            cell = table.cell(row_idx, col_idx)

            # Set cell text with typography support
            set_cell_text_with_typography(
                cell, cell_value,
                font_size=table_font_size,
                font_color=RGBColor(0, 0, 0)
            )

            # White background
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        row_idx += 1

    return True


def create_table_at_position(slide, left, top, width, height, table_data: dict):
    """Create a table at a specific position with IC brand styling.

    Used for inline tables within columns or other constrained areas.

    Args:
        slide: PowerPoint slide object
        left, top, width, height: Position and size in EMUs or Inches
        table_data: Dict with 'headers' and 'rows' keys
    """
    headers = table_data.get('headers', [])
    rows = table_data.get('rows', [])

    if not headers and not rows:
        return None

    col_count = len(headers) if headers else (len(rows[0]) if rows else 0)
    if col_count == 0:
        return None

    row_count = (1 if headers else 0) + len(rows)
    data_row_count = len(rows)

    # Phase 5: Calculate table fit for inline tables (smaller intended fonts)
    intended_data_font = 8  # Inline tables use smaller fonts
    table_font_size = intended_data_font
    table_row_height_emu = None
    table_header_height_emu = None

    if _HAS_CONTENT_FITTER:
        try:
            # height may be EMU or Inches object — convert to inches
            if hasattr(height, 'inches'):
                available_height_inches = height.inches
            else:
                available_height_inches = height / 914400
            fit_result = calculate_table_fit(
                row_count=data_row_count,
                col_count=col_count,
                available_height_inches=available_height_inches,
                intended_font_size_pt=intended_data_font,
            )
            table_font_size = fit_result['font_size_pt']
            table_row_height_emu = int(fit_result['row_height_inches'] * 914400)
            table_header_height_emu = int(fit_result['header_height_inches'] * 914400)
            if not fit_result['fits']:
                _shrink_warnings.append(
                    f"Inline table with {data_row_count} rows may overflow at minimum font size {table_font_size}pt"
                )
                print(f"  [shrink] Inline table with {data_row_count} rows may overflow at minimum font size {table_font_size}pt", file=sys.stderr)
        except Exception as e:
            print(f"Content fitter: could not calculate table fit: {e}", file=sys.stderr)

    # Create table
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = table_shape.table

    # Disable table banding
    tbl = get_table_from_shape(table_shape)
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = tbl.makeelement(qn('a:tblPr'))
        tbl.insert(0, tblPr)
    tblPr.set('bandRow', '0')
    tblPr.set('bandCol', '0')
    tblPr.set('firstRow', '0')
    tblPr.set('lastRow', '0')
    tblPr.set('firstCol', '0')
    tblPr.set('lastCol', '0')

    # Calculate column widths based on content (or fall back to equal widths)
    header_font_size = min(table_font_size + 1, intended_data_font)
    # Convert width to EMU if it's an Inches object
    width_emu = int(width) if not hasattr(width, 'emu') else width.emu
    if _HAS_CONTENT_FITTER:
        try:
            col_widths = calculate_column_widths(
                headers=headers,
                rows=rows,
                total_width_emu=width_emu,
                font_size_pt=table_font_size,
                header_font_size_pt=header_font_size,
            )
            for i, col_w in enumerate(col_widths):
                table.columns[i].width = col_w
        except Exception as e:
            # Fallback to equal widths on error
            print(f"Content fitter: column width calculation failed: {e}", file=sys.stderr)
            col_width = width_emu // col_count
            for i in range(col_count):
                table.columns[i].width = col_width
    else:
        # Fallback: equal column widths
        col_width = width_emu // col_count
        for i in range(col_count):
            table.columns[i].width = col_width

    # Phase 5: Set row heights from calculated fit
    if table_header_height_emu and headers:
        table.rows[0].height = table_header_height_emu
    if table_row_height_emu:
        start_row = 1 if headers else 0
        for r_idx in range(start_row, row_count):
            table.rows[r_idx].height = table_row_height_emu

    # Fill header row (black bg, white text)
    row_idx = 0
    if headers:
        for col_idx, header in enumerate(headers):
            if col_idx < col_count:
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
                # White bold text with typography support
                set_cell_text_with_typography(
                    cell, str(header),
                    font_size=header_font_size,
                    font_color=RGBColor(255, 255, 255),
                    bold=True
                )
        row_idx += 1

    # Fill data rows (white bg, black text)
    for row_data in rows:
        # Handle section header rows (dict with 'type': 'section_header')
        if isinstance(row_data, dict) and row_data.get('type') == 'section_header':
            cell = table.cell(row_idx, 0)
            cell.text = clean_markdown(str(row_data.get('text', '')))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.size = Pt(table_font_size)
            for col_idx in range(1, col_count):
                cell = table.cell(row_idx, col_idx)
                cell.text = ''
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            row_idx += 1
            continue

        # Regular row - pad with empty strings if needed
        padded_row = list(row_data) + [''] * (col_count - len(row_data)) if len(row_data) < col_count else row_data

        for col_idx in range(col_count):
            cell_value = padded_row[col_idx] if col_idx < len(padded_row) else ''
            cell = table.cell(row_idx, col_idx)

            # Set cell text with typography support
            set_cell_text_with_typography(
                cell, cell_value,
                font_size=table_font_size,
                font_color=RGBColor(0, 0, 0)
            )

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        row_idx += 1

    return table_shape


# =============================================================================
# IMAGE PLACEHOLDER CREATION
# =============================================================================

def calculate_image_position_above_body(body_ph_position: dict, slide_width: float = 13.333) -> dict:
    """Calculate position for an image placeholder box above a body placeholder.

    Used when a grid visual type is selected but the actual template layout doesn't
    have PICTURE placeholders. Creates an image area above the body placeholder.

    Args:
        body_ph_position: Dict with left_inches, top_inches, width_inches, height_inches
        slide_width: Total slide width in inches (default 13.333 for widescreen)

    Returns:
        Dict with left, top, width, height in inches for the image placeholder
    """
    left = body_ph_position.get('left_inches', 0)
    body_top = body_ph_position.get('top_inches', 0)
    width = body_ph_position.get('width_inches', 3)

    # Image goes above body, starting from title area bottom (~1.3") to body top
    # with some padding
    image_top = 1.5  # Below title area
    image_height = max(1.5, body_top - image_top - 0.2)  # Leave gap before body

    return {
        'left_inches': left,
        'top_inches': image_top,
        'width_inches': width,
        'height_inches': image_height
    }


def create_image_placeholder_box(slide, left, top, width, height, description: str = "Insert Picture", send_to_back: bool = False):
    """Create a dashed box with label indicating where an image should go.

    Creates a rectangular shape with:
    - No fill (transparent)
    - Dashed gray border
    - Centered italic text with the description

    Args:
        slide: PowerPoint slide object
        left: Left position (EMU or Inches)
        top: Top position (EMU or Inches)
        width: Width (EMU or Inches)
        height: Height (EMU or Inches)
        description: Text label for the placeholder
        send_to_back: If True, move shape to back of z-order (for story-card overlays)
    """

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )

    # Style: no fill, dashed border
    shape.fill.background()
    shape.line.color.rgb = RGBColor(128, 128, 128)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.line.width = Pt(1)

    # Add text
    tf = shape.text_frame
    tf.clear()
    # Center text vertically
    from pptx.enum.text import MSO_ANCHOR
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[Image: {description}]"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(128, 128, 128)
    run.font.italic = True

    # Send to back if needed (for story-card layouts where text overlays image)
    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(0, sp)  # Insert at beginning = back of z-order

    return shape


def calculate_aspect_ratio_mismatch(img_width: int, img_height: int, target_width: int, target_height: int) -> dict:
    """Calculate the aspect ratio mismatch between an image and target box.

    Returns information about how much the image aspect ratio differs from the
    target box, and what percentage of the image will be cropped in cover mode.

    Args:
        img_width: Original image width (in EMU or pixels)
        img_height: Original image height (in EMU or pixels)
        target_width: Target box width (in EMU or same unit as image)
        target_height: Target box height (in EMU or same unit as image)

    Returns:
        dict with:
            - image_ratio: float - image aspect ratio (width/height)
            - target_ratio: float - target box aspect ratio
            - ratio_difference_pct: float - percentage difference between ratios
            - crop_percentage: float - percentage of image that will be cropped in cover mode
            - crop_direction: str - 'width', 'height', or 'none'
            - is_significant: bool - True if mismatch > 20%
    """
    if img_height == 0 or target_height == 0:
        return {
            'image_ratio': 0,
            'target_ratio': 0,
            'ratio_difference_pct': 0,
            'crop_percentage': 0,
            'crop_direction': 'none',
            'is_significant': False,
        }

    image_ratio = img_width / img_height
    target_ratio = target_width / target_height

    # Calculate percentage difference between aspect ratios
    if target_ratio == 0:
        ratio_difference_pct = 0
    else:
        ratio_difference_pct = abs(image_ratio - target_ratio) / target_ratio * 100

    # Calculate how much would be cropped in cover mode
    scale_x = target_width / img_width
    scale_y = target_height / img_height
    scale = max(scale_x, scale_y)

    new_width = img_width * scale
    new_height = img_height * scale

    if new_width > target_width:
        crop_amount = new_width - target_width
        crop_percentage = (crop_amount / new_width) * 100
        crop_direction = 'width'
    elif new_height > target_height:
        crop_amount = new_height - target_height
        crop_percentage = (crop_amount / new_height) * 100
        crop_direction = 'height'
    else:
        crop_percentage = 0
        crop_direction = 'none'

    # Significant if more than 20% crop or 25% ratio difference
    is_significant = crop_percentage > 20 or ratio_difference_pct > 25

    return {
        'image_ratio': round(image_ratio, 3),
        'target_ratio': round(target_ratio, 3),
        'ratio_difference_pct': round(ratio_difference_pct, 1),
        'crop_percentage': round(crop_percentage, 1),
        'crop_direction': crop_direction,
        'is_significant': is_significant,
    }


def insert_image_from_path(
    slide,
    file_path: str,
    left,
    top,
    width,
    height,
    description: str = "",
    fit_mode: str = "cover",
) -> bool:
    """Insert an actual image from a file path with aspect ratio handling.

    Supports two fit modes:
    - 'cover' (default): Image is scaled to fill the entire box while
      maintaining aspect ratio. Parts that exceed the box are cropped.
    - 'contain': Image is scaled to fit entirely within the box while
      maintaining aspect ratio. May result in letterboxing (empty space).

    Emits a warning if aspect ratio mismatch would cause significant cropping
    (>20% of image cropped or >25% aspect ratio difference).

    Args:
        slide: PowerPoint slide object
        file_path: Path to image file (absolute or relative to working directory)
        left: Left position (EMU or Inches)
        top: Top position (EMU or Inches)
        width: Width of target box (EMU or Inches)
        height: Height of target box (EMU or Inches)
        description: Fallback description if image insertion fails
        fit_mode: 'cover' (fill and crop) or 'contain' (fit within, may letterbox)

    Returns:
        True if image was inserted, False if fallback placeholder was used
    """
    from pathlib import Path as PathLib
    import os

    # Validate fit_mode
    if fit_mode not in ('cover', 'contain'):
        fit_mode = 'cover'

    # Resolve file path
    image_path = PathLib(file_path)
    if not image_path.is_absolute():
        # Try relative to current working directory
        image_path = PathLib(os.getcwd()) / file_path

    # Check if file exists
    if not image_path.exists():
        _collect_warning(f"Image not found: {file_path}, using placeholder")
        create_image_placeholder_box(slide, left, top, width, height, description or file_path)
        return False

    # Check if it's a valid image format
    valid_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.emf', '.wmf'}
    if image_path.suffix.lower() not in valid_extensions:
        _collect_warning(f"Invalid image format: {image_path.suffix}, using placeholder")
        create_image_placeholder_box(slide, left, top, width, height, description or file_path)
        return False

    try:
        # Insert image at position
        # First add image to get actual dimensions
        pic = slide.shapes.add_picture(str(image_path), left, top)

        # Get original image dimensions
        img_width = pic.width
        img_height = pic.height

        # Convert target dimensions to EMU if needed
        target_width = width.emu if hasattr(width, 'emu') else width
        target_height = height.emu if hasattr(height, 'emu') else height

        # Check aspect ratio mismatch and warn if significant
        mismatch = calculate_aspect_ratio_mismatch(img_width, img_height, target_width, target_height)
        if mismatch['is_significant']:
            img_name = PathLib(file_path).name
            _collect_warning(
                f"Image aspect ratio mismatch for '{img_name}': "
                f"image ratio {mismatch['image_ratio']:.2f}, target ratio {mismatch['target_ratio']:.2f}. "
                f"{mismatch['crop_percentage']:.0f}% will be cropped from {mismatch['crop_direction']}."
            )

        if fit_mode == 'contain':
            # Contain mode: fit within box, may have letterboxing
            scale_x = target_width / img_width
            scale_y = target_height / img_height
            scale = min(scale_x, scale_y)  # Use smaller scale to fit within

            new_width = int(img_width * scale)
            new_height = int(img_height * scale)

            # Center the image within the target box
            offset_x = (target_width - new_width) // 2
            offset_y = (target_height - new_height) // 2

            # Set final dimensions and position (centered within box)
            pic.width = new_width
            pic.height = new_height
            pic.left = int(left.emu if hasattr(left, 'emu') else left) + offset_x
            pic.top = int(top.emu if hasattr(top, 'emu') else top) + offset_y

        else:
            # Cover mode (default): fill entire box, crop overflow
            scale_x = target_width / img_width
            scale_y = target_height / img_height
            scale = max(scale_x, scale_y)  # Use larger scale to fill

            # Apply scaling
            new_width = int(img_width * scale)
            new_height = int(img_height * scale)

            # Center the image in the target box (cropping excess)
            offset_x = (new_width - target_width) // 2
            offset_y = (new_height - target_height) // 2

            # Set final dimensions
            pic.width = new_width
            pic.height = new_height

            # Position to center (with negative offset to crop)
            pic.left = int(left.emu if hasattr(left, 'emu') else left) - offset_x
            pic.top = int(top.emu if hasattr(top, 'emu') else top) - offset_y

            # Apply cropping via shape XML for precise box fit
            # Get the shape's spPr (shape properties) element
            spPr = pic._element.spPr
            if spPr is not None:
                # For proper cropping, we need to use a:srcRect
                # This requires the blipFill element
                blipFill = pic._element.blipFill
                if blipFill is not None:
                    # Find or create srcRect
                    srcRect = blipFill.find(qn('a:srcRect'))
                    if srcRect is None:
                        # Insert srcRect after blip
                        blip = blipFill.find(qn('a:blip'))
                        if blip is not None:
                            srcRect = etree.Element(qn('a:srcRect'))
                            blip.addnext(srcRect)

                    if srcRect is not None:
                        # Calculate crop percentages (in 1/100,000ths)
                        if offset_x > 0:
                            crop_pct = int(offset_x / new_width * 100000)
                            srcRect.set('l', str(crop_pct))
                            srcRect.set('r', str(crop_pct))
                        if offset_y > 0:
                            crop_pct = int(offset_y / new_height * 100000)
                            srcRect.set('t', str(crop_pct))
                            srcRect.set('b', str(crop_pct))

                        # Reset position to target after setting crop
                        pic.left = int(left.emu if hasattr(left, 'emu') else left)
                        pic.top = int(top.emu if hasattr(top, 'emu') else top)
                        pic.width = int(target_width)
                        pic.height = int(target_height)

        return True

    except Exception as e:
        _collect_warning(f"Failed to insert image {file_path}: {e}, using placeholder")
        create_image_placeholder_box(slide, left, top, width, height, description or file_path)
        return False


def insert_image_or_placeholder(
    slide,
    file_path: str | None,
    left,
    top,
    width,
    height,
    description: str = "",
    fit_mode: str = "cover",
    send_to_back: bool = False,
):
    """Insert image from file_path if provided, otherwise create placeholder box.

    Convenience wrapper for handlers to use consistently.

    Args:
        slide: PowerPoint slide object
        file_path: Optional path to image file (None uses placeholder)
        left: Left position
        top: Top position
        width: Width
        height: Height
        description: Description for placeholder/fallback
        fit_mode: 'cover' (fill and crop) or 'contain' (fit within, may letterbox)
        send_to_back: If True, send placeholder to back of z-order (for story-card overlays)
    """
    if file_path:
        insert_image_from_path(slide, file_path, left, top, width, height, description, fit_mode)
    else:
        create_image_placeholder_box(slide, left, top, width, height, description, send_to_back=send_to_back)


# =============================================================================
# FOOTNOTE/CALLOUT/QUOTE CREATION
# =============================================================================

def render_text_with_typography(
    paragraph,
    text: str = None,
    styled_runs: list = None,
    base_bold: bool = False,
    base_italic: bool = False,
    font_size: int = None,
    font_color: RGBColor = None
):
    """Unified typography renderer for all text content.

    Priority order:
    1. styled_runs (if provided) - pre-parsed rich text formatting
    2. inline markers in text - {blue}, {question}, {signpost}, etc.
    3. markdown in text - **bold**, *italic*
    4. base formatting - base_bold, base_italic

    Args:
        paragraph: PowerPoint paragraph object
        text: Plain text (used if styled_runs not provided)
        styled_runs: Pre-parsed list of {text, style} dicts
        base_bold: Default bold state for unstyled text
        base_italic: Default italic state for unstyled text
        font_size: Font size in points (None = inherit)
        font_color: Font color (None = inherit)
    """
    if styled_runs:
        # Use pre-parsed styled runs
        render_styled_runs(paragraph, styled_runs, base_bold=base_bold, base_italic=base_italic)
        # Apply size/color overrides to all runs if specified
        if font_size or font_color:
            for run in paragraph.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if font_color:
                    run.font.color.rgb = font_color
    elif text:
        # Check for inline markers
        if any(marker in text for marker in ['{blue}', '{italic}', '{bold}', '{signpost}', '{question}', '{red}', '{green}', '{color:']):
            add_styled_paragraph_with_markers(paragraph, text, base_bold=base_bold, base_italic=base_italic)
            # Apply size overrides (always) and color overrides (only if run has no color set)
            if font_size or font_color:
                for run in paragraph.runs:
                    if font_size:
                        run.font.size = Pt(font_size)
                    if font_color:
                        # Only apply default color if run doesn't have a color already set by markers
                        try:
                            has_color = run.font.color and run.font.color.type is not None
                        except AttributeError:
                            has_color = False
                        if not has_color:
                            run.font.color.rgb = font_color
        elif '**' in text or ('*' in text and '**' not in text):
            # Markdown formatting
            segments = parse_markdown_formatting(text)
            for seg in segments:
                run = paragraph.add_run()
                run.text = seg['text']
                run.font.bold = seg['bold'] or base_bold
                run.font.italic = seg['italic'] or base_italic
                if font_size:
                    run.font.size = Pt(font_size)
                if font_color:
                    run.font.color.rgb = font_color
        else:
            # Plain text
            run = paragraph.add_run()
            run.text = clean_markdown(text)
            run.font.bold = base_bold
            run.font.italic = base_italic
            if font_size:
                run.font.size = Pt(font_size)
            if font_color:
                run.font.color.rgb = font_color


def add_footnote(slide, text: str, styled_runs: list = None):
    """Add footnote text at bottom of slide.

    Args:
        slide: PowerPoint slide object
        text: Footnote text (fallback if styled_runs not provided)
        styled_runs: Optional list of {text, style} dicts for rich formatting
    """
    if not text and not styled_runs:
        return

    # Use margin-safe positioning
    safe = get_safe_area()
    left = Inches(safe.left)
    top = Inches(6.8)  # Near bottom, but within safe area
    width = Inches(safe.width)
    height = Inches(0.4)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]

    render_text_with_typography(
        p,
        text=text,
        styled_runs=styled_runs,
        base_italic=True,
        font_size=8,
        font_color=RGBColor(100, 100, 100)
    )


def add_callout(slide, text: str, position: str = "bottom", styled_runs: list = None):
    """Add callout/highlight text (bold).

    Args:
        slide: PowerPoint slide object
        text: Callout text (fallback if styled_runs not provided)
        position: "bottom" or "top"
        styled_runs: Optional list of {text, style} dicts for rich formatting
    """
    if not text and not styled_runs:
        return

    # Use margin-safe positioning
    safe = get_safe_area()
    if position == "bottom":
        left = Inches(safe.left)
        top = Inches(6.3)
    else:
        left = Inches(safe.left)
        top = Inches(1.5)

    width = Inches(safe.width)
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]

    # Only force black color when no styled_runs - let styled colors take effect
    render_text_with_typography(
        p,
        text=text,
        styled_runs=styled_runs,
        base_bold=True,
        font_size=11,
        font_color=None if styled_runs else RGBColor(0, 0, 0)
    )


def add_callout_block(slide, header: str, body: list, position: str = "bottom",
                      header_styled: list = None, body_styled: list = None):
    """Add callout block with header (bold) and bullet items.

    Args:
        slide: PowerPoint slide object
        header: Header text (fallback if header_styled not provided)
        body: List of body items (fallback if body_styled not provided)
        position: "bottom" or "top"
        header_styled: Optional styled runs for header
        body_styled: Optional list of styled items (each item is a list of {text, style})
    """
    if not header and not body and not header_styled and not body_styled:
        return

    # Use margin-safe positioning
    safe = get_safe_area()
    if position == "bottom":
        left = Inches(safe.left)
        top = Inches(5.5)
    else:
        left = Inches(safe.left)
        top = Inches(1.5)

    width = Inches(safe.width)
    # Calculate height based on content
    line_count = (1 if header or header_styled else 0) + len(body if body else body_styled or [])
    height = Inches(0.3 * line_count)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    first_para = True

    # Add header
    if header or header_styled:
        p = tf.paragraphs[0]
        first_para = False
        render_text_with_typography(
            p,
            text=header,
            styled_runs=header_styled,
            base_bold=True,
            font_size=10,
            font_color=RGBColor(0, 0, 0)
        )

    # Add body items
    items = body_styled if body_styled else body
    for i, item in enumerate(items or []):
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        brand_font = get_brand_font('body')
        if body_styled and isinstance(item, dict) and 'runs' in item:
            # Styled body item
            run = p.add_run()
            run.text = "• "
            run.font.name = brand_font
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(60, 60, 60)
            render_styled_runs(p, item['runs'])
            # Apply size/color to rendered runs (font.name already set by render_styled_runs)
            for run in p.runs[1:]:  # Skip bullet run
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(60, 60, 60)
        else:
            # Plain text item
            item_text = item if isinstance(item, str) else str(item)
            run = p.add_run()
            run.text = f"• {clean_markdown(item_text)}"
            run.font.name = brand_font
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(60, 60, 60)


def add_quote(slide, text: str, position: str = "bottom", styled_runs: list = None, attribution: str = None):
    """Add quote text (italic, with quotation marks) and optional attribution.

    Args:
        slide: PowerPoint slide object
        text: Quote text (fallback if styled_runs not provided)
        position: "bottom", "top", or "center"
        styled_runs: Optional list of {text, style} dicts for styled formatting
        attribution: Optional attribution line (e.g., "— Speaker Name, Title")
    """
    if not text and not styled_runs:
        return

    # Use margin-safe positioning
    safe = get_safe_area()
    if position == "center":
        left = Inches(safe.left)
        top = Inches(3.5)  # Vertically centered on 7.5" slide
    elif position == "bottom":
        left = Inches(safe.left)
        top = Inches(6.0)
    else:
        left = Inches(safe.left)
        top = Inches(5.5)

    width = Inches(safe.width)
    # Increased height from 0.6"/0.8" to 1.2"/1.4" to accommodate longer quotes/statements
    # Fixes overflow in hero-statement visual type (1 case with 221-char statement)
    height = Inches(1.4) if attribution else Inches(1.2)
    brand_font = get_brand_font('body')

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    if styled_runs:
        # Add opening quote
        run = p.add_run()
        run.text = '"'
        run.font.name = brand_font
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = RGBColor(80, 80, 80)

        # Render styled runs
        for run_data in styled_runs:
            run_text = run_data.get('text', '')
            style = run_data.get('style')
            run = p.add_run()
            run.text = run_text
            run.font.name = brand_font
            run.font.size = Pt(11)
            run.font.italic = True
            run.font.color.rgb = RGBColor(80, 80, 80)
            # Apply style overrides (e.g., color)
            if style:
                apply_style_to_run(run, style)

        # Add closing quote
        run = p.add_run()
        run.text = '"'
        run.font.name = brand_font
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = RGBColor(80, 80, 80)
    else:
        run = p.add_run()
        # Add quotation marks if not already present
        quote_text = clean_markdown(text)
        if not quote_text.startswith('"') and not quote_text.startswith('"'):
            quote_text = f'"{quote_text}"'
        run.text = quote_text
        run.font.name = brand_font
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = RGBColor(80, 80, 80)

    # Add attribution on new line if provided
    if attribution:
        p2 = tf.add_paragraph()
        run = p2.add_run()
        # Add em-dash if not already present
        attr_text = attribution if attribution.startswith('—') or attribution.startswith('-') else f"— {attribution}"
        run.text = attr_text
        run.font.name = brand_font
        run.font.size = Pt(10)
        run.font.italic = False
        run.font.color.rgb = RGBColor(100, 100, 100)


def add_hero_statement(slide, text: str, styled_runs: list = None):
    """Add hero statement text - large, bold, centered for maximum impact.

    Hero statements are single punchy statements (1-2 sentences) that should
    dominate the slide visually. Unlike quotes, they use:
    - Large font (28pt)
    - Regular weight (not italic)
    - Black text (not gray)
    - Center alignment (horizontal and vertical)
    - No quotation marks

    Args:
        slide: PowerPoint slide object
        text: Statement text (fallback if styled_runs not provided)
        styled_runs: Optional list of {text, style} dicts for styled formatting
    """
    if not text and not styled_runs:
        return

    # Use margin-safe positioning, vertically centered
    safe = get_safe_area()
    left = Inches(safe.left)
    # Center vertically: slide is 7.5" tall, title ~1.6" from top
    # Statement should be in the middle third, around 3.0-3.5" from top
    top = Inches(3.0)
    width = Inches(safe.width)
    height = Inches(1.5)  # Allow for 2-line statements at 28pt

    brand_font = get_brand_font('body')

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if styled_runs:
        # Render styled runs with hero statement formatting
        for run_data in styled_runs:
            run_text = run_data.get('text', '')
            style = run_data.get('style')
            run = p.add_run()
            run.text = run_text
            run.font.name = brand_font
            run.font.size = Pt(28)
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = RGBColor(0, 0, 0)
            # Apply style overrides (e.g., color, bold)
            if style:
                apply_style_to_run(run, style)
    else:
        run = p.add_run()
        run.text = clean_markdown(text)
        run.font.name = brand_font
        run.font.size = Pt(28)
        run.font.bold = False
        run.font.italic = False
        run.font.color.rgb = RGBColor(0, 0, 0)


# =============================================================================
# SLIDE POPULATION FUNCTIONS
# =============================================================================

def populate_title_cover_slide(slide, placeholders: list, slide_data: dict):
    """
    Special handling for title-cover layout (may not have TITLE placeholder).

    Maps content to BODY placeholders by position:
    - Largest BODY: title + subtitle
    - Bottom BODYs: metadata items (sorted bottom-to-top, left-to-right)
    - ALL unused BODYs: cleared to prevent "Click to add text"
    """
    content = slide_data.get('content', {})
    title = content.get('title', '')
    subtitle = content.get('subtitle', '')
    metadata = content.get('metadata', [])
    slide_num = slide_data.get('slide_number', 0)

    # Track which placeholders we've filled
    used_idxs = set()

    # Try TITLE placeholder first
    title_ph = get_title_placeholder(placeholders)
    if title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'], slide_num=slide_num, ph_type='TITLE')
        full_title = title
        if subtitle:
            full_title = f"{title}\n{subtitle}"
        if ph:
            set_text_in_placeholder(ph, full_title, bold=True)
            used_idxs.add(title_ph['idx'])
    else:
        # No TITLE placeholder - use largest BODY for title
        body_phs = get_body_placeholders(placeholders)
        largest = get_largest_body_placeholder(placeholders)

        if title and largest:
            ph = get_placeholder_by_idx(slide, largest['idx'], slide_num=slide_num, ph_type='BODY')
            if ph:
                full_title = title
                if subtitle:
                    full_title = f"{title}\n{subtitle}"
                set_text_in_placeholder(ph, full_title, bold=True)
                used_idxs.add(largest['idx'])

    # Get all BODY placeholders for metadata distribution
    body_phs = get_body_placeholders(placeholders)
    remaining_phs = [p for p in body_phs if p['idx'] not in used_idxs]

    # Title-cover layout uses specific placeholder convention:
    # - idx 15 (bottom-center): client attribution (first metadata item)
    # - idx 10 (bottom-left): date (second metadata item)
    # This matches Inner Chapter template conventions
    if len(metadata) >= 1:
        # First metadata item (e.g., "Mondelez × Inner Chapter") goes to idx 15
        ph_15 = get_placeholder_by_idx(slide, 15, slide_num=slide_num, ph_type='metadata', warn_on_missing=False)
        if ph_15:
            set_text_in_placeholder(ph_15, metadata[0])
            used_idxs.add(15)

    if len(metadata) >= 2:
        # Second metadata item (e.g., "January 2026") goes to idx 10
        ph_10 = get_placeholder_by_idx(slide, 10, slide_num=slide_num, ph_type='metadata', warn_on_missing=False)
        if ph_10:
            set_text_in_placeholder(ph_10, metadata[1])
            used_idxs.add(10)

    # CRITICAL: Clear ALL unused placeholders to prevent "Click to add text"
    all_body_phs = get_body_placeholders(placeholders)
    for body_ph in all_body_phs:
        if body_ph['idx'] not in used_idxs:
            ph = get_placeholder_by_idx(slide, body_ph['idx'])
            if ph:
                # Set empty text to clear the placeholder
                tf = ph.text_frame
                tf.clear()
                # Add an empty paragraph to fully clear
                p = tf.paragraphs[0]
                p.clear()


def populate_content_slide(slide, placeholders: list, slide_data: dict):
    """
    Standard content slide population.

    Handles:
    - TITLE placeholder for title
    - Single BODY placeholder with combined headline + bullets
    - Tables (from any supported location)
    - Footnotes, callouts, and quotes
    - Subsections
    - Image placeholders (for content-image-right layouts with PICTURE placeholders)
    - Styled runs (headline_styled, body_styled, quote_styled)
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    slide_num = slide_data.get('slide_number', 0)

    title = content.get('title', '')
    headline = content.get('headline', '')
    body = content.get('body', [])
    subsection = content.get('subsection', {})

    # Extract styled runs (from slide-outline-to-layout parser)
    headline_styled = slide_data.get('headline_styled')
    body_styled = slide_data.get('body_styled')
    quote_styled = slide_data.get('quote_styled')

    # Get placeholder references
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders)

    # Set title
    if title:
        if title_ph:
            ph = get_placeholder_by_idx(slide, title_ph['idx'], slide_num=slide_num, ph_type='TITLE')
            if ph:
                set_text_in_placeholder(ph, title)
        elif body_phs:
            # No TITLE placeholder, use largest BODY
            largest = get_largest_body_placeholder(placeholders)
            if largest:
                ph = get_placeholder_by_idx(slide, largest['idx'], slide_num=slide_num, ph_type='BODY')
                if ph:
                    set_text_in_placeholder(ph, title, bold=True)
                    body_phs = [p for p in body_phs if p['idx'] != largest['idx']]

    # Handle tables
    tables = extract_tables(slide_data)
    string_body = [b for b in body if isinstance(b, str)]

    # Decide table placement strategy based on content
    # If we have both body content AND tables, create table as free-form shape
    # If we only have tables, use body placeholders for tables
    has_body_content = bool(headline or string_body)

    if tables:
        if body_phs and not has_body_content:
            # No body content - use body placeholders for tables
            for i, table in enumerate(tables):
                if i < len(body_phs):
                    ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                    if ph:
                        create_table_in_placeholder(slide, ph, table)
            # Tables consumed body placeholders
            body_phs = []
        else:
            # Either no body placeholders OR we have body content
            # Add tables as free-form shapes
            # Position tables on the right side (for content-image-right layouts)
            # or below the body content area
            picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
            if picture_phs:
                # Use PICTURE placeholder position for table
                pic_pos = picture_phs[0].get('position', {})
                table_left = Inches(pic_pos.get('left_inches', 6.74))
                table_top = Inches(pic_pos.get('top_inches', 0.42))
                table_width = Inches(pic_pos.get('width_inches', 6.0))
                table_height = Inches(min(pic_pos.get('height_inches', 6.0), 6.0))
            else:
                # Default positioning - full width table using safe margins
                safe = get_safe_area()
                table_top = Inches(1.8)
                table_left = Inches(safe.left)
                table_width = Inches(safe.width)
                table_height = Inches(4.5)

            for i, table in enumerate(tables):
                create_table_at_position(
                    slide,
                    table_left,
                    table_top + Inches(i * 2.5),  # Stack multiple tables
                    table_width,
                    table_height,
                    table
                )

    # Set body content (headline + bullets combined)
    # Only process if body contains strings, not objects
    # (string_body already computed above for table placement decision)

    if (headline or string_body or headline_styled or body_styled) and body_phs:
        ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
        if ph:
            # Determine if headline should be bold based on slide_data flag
            # This flag is set by ingest when headline comes from a bold line
            headline_is_bold = slide_data.get('headline_bold', False)
            set_header_and_bullets(
                ph, headline, string_body,
                header_styled=headline_styled,
                body_styled=body_styled,
                header_bold=headline_is_bold
            )
            
            # Handle subsection - append to same placeholder
            if subsection and subsection.get('header'):
                tf = ph.text_frame
                # Add blank line
                p = tf.add_paragraph()
                p.add_run().text = ""
                # Add subsection header
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = clean_markdown(subsection.get('header', ''))
                run.font.bold = True
                p.level = 0
                # Add subsection items
                for item in subsection.get('items', []):
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = clean_markdown(item)
                    p.level = 1
    
    # Add extras
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)
    
    # Handle callout_header + callout_body pattern
    callout_header = extras.get('callout_header')
    callout_body = extras.get('callout_body', [])
    if callout_header or callout_body:
        add_callout_block(slide, callout_header, callout_body)
    else:
        # Fall back to simple callout
        callout = content.get('callout') or extras.get('callout')
        if callout:
            add_callout(slide, callout)
    
    quote = content.get('quote') or extras.get('quote')
    if quote or quote_styled:
        add_quote(slide, quote, styled_runs=quote_styled)

    # Handle image placeholder (for content-image-right layouts)
    image_placeholder = extras.get('image_placeholder') or content.get('image_placeholder')
    image_file_path = extras.get('image_file') or content.get('image_file')  # Path to actual image
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']

    if (image_placeholder or image_file_path) and picture_phs:
        pic_ph = picture_phs[0]
        pos = pic_ph.get('position', {})
        # Move original placeholder off-slide
        pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
        if pic_shape:
            pic_shape.left = Inches(-100)
            pic_shape.top = Inches(-100)
        # Insert image or create placeholder box
        insert_image_or_placeholder(
            slide,
            image_file_path,
            Inches(pos.get('left_inches', 0)),
            Inches(pos.get('top_inches', 0)),
            Inches(pos.get('width_inches', 4)),
            Inches(pos.get('height_inches', 4)),
            image_placeholder or "Insert image"
        )


def populate_column_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle column layouts (column-2, column-3, column-4).

    Uses structured `columns` array if present, otherwise distributes body items.
    Each column gets: header (bold) + subheader (italic) + body bullets + optional output

    Also handles:
    - headline: Bold text below title, above columns
    - leading_text: Plain text paragraph above columns
    - trailing_note: Labeled note below columns
    """
    content = slide_data.get('content', {})
    columns = extract_columns(slide_data)
    extras = slide_data.get('extras', {})
    slide_num = slide_data.get('slide_number', 0)

    title = content.get('title', '')
    headline = content.get('headline', '')
    body = content.get('body', [])
    leading_text = slide_data.get('leading_text', '')
    trailing_note = slide_data.get('trailing_note', {})

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'], slide_num=slide_num, ph_type='TITLE')
        if ph:
            set_text_in_placeholder(ph, title)

    # Add headline as callout below title if present
    if headline:
        add_callout(slide, headline, position='top')

    # If we have structured columns, use them
    if columns:
        for col_idx, column in enumerate(columns):
            if col_idx >= len(body_phs):
                break

            col_header = column.get('header', '')
            col_subheader = column.get('subheader', '')
            col_intro = column.get('intro', [])  # Plain text before bullets
            col_body = column.get('body', [])
            col_output = column.get('output', '')
            col_table = column.get('table')  # Inline table
            col_image = column.get('image_placeholder')
            col_file_path = column.get('file_path')  # Path to actual image file

            # Get placeholder info
            ph = get_placeholder_by_idx(slide, body_phs[col_idx]['idx'], slide_num=slide_num, ph_type='BODY')
            ph_pos = body_phs[col_idx].get('position', {})

            # If column has a table, render header + table instead of bullets
            if col_table and ph:
                # Get placeholder position for table placement
                ph_left = Inches(ph_pos.get('left_inches', 0))
                ph_top = Inches(ph_pos.get('top_inches', 0))
                ph_width = Inches(ph_pos.get('width_inches', 2))
                ph_height = Inches(ph_pos.get('height_inches', 3))

                # Set header in placeholder
                if col_header:
                    set_text_in_placeholder(ph, col_header, bold=True)
                    # Move placeholder to just show header (shrink height)
                    ph.height = Inches(0.4)
                    # Position table below header
                    table_top = ph_top + Inches(0.5)
                    table_height = ph_height - Inches(0.5)
                else:
                    # Hide placeholder, use full area for table
                    ph.left = Inches(-100)
                    ph.top = Inches(-100)
                    table_top = ph_top
                    table_height = ph_height

                # Create table in column area
                create_table_at_position(
                    slide,
                    ph_left, table_top,
                    ph_width, table_height,
                    col_table
                )
            else:
                # Standard column content: header + intro + key_values + bullets
                col_key_values = column.get('key_values', {})
                col_number = column.get('number')  # For cards converted to columns
                col_header_styled = column.get('header_styled')  # Typography-aware header

                # Build body content
                body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                if col_output:
                    body_items.append(f"Output: {col_output}")

                # For numbered cards, add the number prefix to the header
                header_with_number = col_header
                header_styled_with_number = col_header_styled
                if col_number is not None and col_header:
                    header_with_number = f"{col_number}. {col_header}"
                    # Also prepend number to styled header if present
                    if col_header_styled:
                        header_styled_with_number = [{'text': f"{col_number}. ", 'style': None}] + col_header_styled

                # Set column content with intro text and key_values (bold keys)
                # For numbered cards (framework cards), add blank lines between items
                is_numbered_card = col_number is not None
                if ph:
                    set_header_and_bullets(
                        ph, header_with_number, body_items,
                        subheader=col_subheader,
                        intro=col_intro,
                        key_values=col_key_values if col_key_values else None,
                        header_styled=header_styled_with_number,
                        header_bold=True,  # Column headers are bold
                        add_blank_after_header=is_numbered_card,
                        add_blank_between_items=is_numbered_card
                    )

            # Handle image placeholder or actual image
            if (col_image or col_file_path) and col_idx < len(picture_phs):
                pic_ph = picture_phs[col_idx]
                pos = pic_ph.get('position', {})
                # Move the picture placeholder off-slide
                pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                if pic_shape:
                    pic_shape.left = Inches(-100)
                    pic_shape.top = Inches(-100)
                insert_image_or_placeholder(
                    slide,
                    col_file_path,
                    Inches(pos.get('left_inches', 0)),
                    Inches(pos.get('top_inches', 0)),
                    Inches(pos.get('width_inches', 2)),
                    Inches(pos.get('height_inches', 1.5)),
                    col_image or "Insert image"
                )
    else:
        # Fallback: distribute body items evenly across columns
        # First item in each column becomes the header (bold, level 0)
        # Remaining items become bullets (level 1)
        string_body = [b for b in body if isinstance(b, str)]
        num_columns = len(body_phs)
        if num_columns > 0 and string_body:
            items_per_column = max(1, len(string_body) // num_columns)

            for col_idx, body_ph in enumerate(body_phs):
                start_idx = col_idx * items_per_column
                end_idx = start_idx + items_per_column if col_idx < num_columns - 1 else len(string_body)
                column_items = string_body[start_idx:end_idx]

                if column_items:
                    ph = get_placeholder_by_idx(slide, body_ph['idx'])
                    if ph:
                        # First item is header, rest are bullets
                        col_header = column_items[0] if column_items else None
                        col_bullets = column_items[1:] if len(column_items) > 1 else []
                        set_header_and_bullets(ph, col_header, col_bullets, header_bold=True)
    
    # Add leading_text as callout at top if present
    if leading_text:
        add_callout(slide, leading_text, position='top')

    # Add trailing_note as labeled note if present
    if trailing_note:
        note_label = trailing_note.get('label', '')
        note_text = trailing_note.get('text', '')
        if note_label and note_text:
            add_callout(slide, f"**{note_label}:** {note_text}", position='bottom')
        elif note_text:
            add_callout(slide, note_text, position='bottom')

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_data_contrast_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle data contrast visualization.

    Renders two contrasting metrics side-by-side with visual prominence.
    Metrics are displayed with large font sizes to create visual tension/impact.
    Layout: Large value on top, label below, description at bottom.
    """
    content = slide_data.get('content', {})
    data_contrast = slide_data.get('data_contrast', {})
    extras = slide_data.get('extras', {})
    slide_num = slide_data.get('slide_number', 0)

    title = content.get('title', '')
    headline = content.get('headline', '')
    headline_styled = slide_data.get('headline_styled')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'], slide_num=slide_num, ph_type='TITLE')
        if ph:
            set_text_in_placeholder(ph, title)

    # Set headline as callout text at top of slide
    if headline:
        add_callout(slide, headline, position="top", styled_runs=headline_styled)

    # Extract data contrast fields
    metrics = data_contrast.get('metrics', [])
    interpretation = data_contrast.get('interpretation', '')
    question = data_contrast.get('question', '')

    # Hide body placeholders - we'll draw custom shapes for prominent metrics
    for body_ph in body_phs:
        ph = get_placeholder_by_idx(slide, body_ph['idx'], slide_num=slide_num, ph_type='BODY', warn_on_missing=False)
        if ph:
            ph.left = Inches(-100)
            ph.top = Inches(-100)

    # Get safe area for positioning metrics
    safe = get_safe_area()
    brand_font = get_brand_font('body')

    # Calculate positions for side-by-side metrics
    num_metrics = len(metrics)
    if num_metrics == 0:
        return

    # Positioning: metrics centered in content area
    content_top = Inches(1.8)  # Below title
    content_height = Inches(3.5)  # Available height for metrics

    # Divide horizontal space evenly between metrics
    metric_width = Inches(safe.width / num_metrics)

    for i, metric in enumerate(metrics):
        label = metric.get('label', '')
        value = metric.get('value', '')
        description = metric.get('description', '')

        # Calculate horizontal position for this metric
        metric_left = Inches(safe.left) + (metric_width * i)
        metric_center = metric_left + (metric_width / 2)

        # Create text box for the large metric value
        value_box = slide.shapes.add_textbox(
            metric_left + Inches(0.2),
            content_top,
            metric_width - Inches(0.4),
            Inches(1.2)
        )
        tf = value_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = clean_markdown(value)
        run.font.name = brand_font
        run.font.size = Pt(48)  # Large, prominent font
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)

        # Create text box for the label
        label_box = slide.shapes.add_textbox(
            metric_left + Inches(0.2),
            content_top + Inches(1.2),
            metric_width - Inches(0.4),
            Inches(0.5)
        )
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = clean_markdown(label)
        run.font.name = brand_font
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(80, 80, 80)

        # Create text box for the description
        if description:
            desc_box = slide.shapes.add_textbox(
                metric_left + Inches(0.2),
                content_top + Inches(1.7),
                metric_width - Inches(0.4),
                Inches(0.8)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = clean_markdown(description)
            run.font.name = brand_font
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(100, 100, 100)

    # Add interpretation as callout if provided
    if interpretation:
        add_callout(slide, f"Interpretation: {interpretation}")

    # Add question as footnote if provided
    if question:
        add_footnote(slide, f"Core question: {question}")


def populate_timeline_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle timeline visualization.

    Routes to either:
    - Gantt table for complex timelines (4+ entries with week ranges)
    - Visual horizontal timeline for simple timelines

    Draws a visual horizontal timeline with:
    - Connecting line across the slide
    - Circle markers for each date
    - Date labels above the line
    - Activity labels below the line
    - Milestone entries highlighted in green
    """
    content = slide_data.get('content', {})
    timeline = extract_timeline(slide_data)
    extras = slide_data.get('extras', {})
    slide_num = slide_data.get('slide_number', 0)

    title = content.get('title', '')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders)

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'], slide_num=slide_num, ph_type='TITLE')
        if ph:
            set_text_in_placeholder(ph, title)

    if not timeline:
        return

    # Check if timeline should be rendered as Gantt chart
    # Import here to avoid circular imports and keep optional
    try:
        from gantt_renderer import should_render_as_gantt, build_gantt_data_from_timeline_entries, create_gantt_table
        if should_render_as_gantt(timeline):
            # Hide body placeholders
            for body_ph in body_phs:
                ph = get_placeholder_by_idx(slide, body_ph['idx'], slide_num=slide_num, ph_type='BODY', warn_on_missing=False)
                if ph:
                    ph.left = Inches(-100)
                    ph.top = Inches(-100)

            # Build Gantt data from timeline entries
            gantt_data = build_gantt_data_from_timeline_entries(timeline, title=title or "Timeline")

            # Render Gantt table with margin-safe positioning
            safe = get_safe_area()
            table_left = Inches(safe.left)
            table_top = Inches(1.5)
            table_width = Inches(safe.width)
            # Height based on number of tasks
            num_tasks = len(gantt_data.all_tasks)
            table_height = Inches(min(5.5, 0.35 * (num_tasks + 2) + 0.5))

            create_gantt_table(slide, gantt_data, table_left, table_top, table_width, table_height)

            # Add footnote if present
            footnote = extras.get('footnote') or extras.get('note') or content.get('footnote')
            if footnote:
                add_footnote(slide, footnote)
            return
    except ImportError:
        pass  # Gantt renderer not available, fall through to simple timeline

    # Hide body placeholders (we'll draw custom shapes)
    for body_ph in body_phs:
        ph = get_placeholder_by_idx(slide, body_ph['idx'], slide_num=slide_num, ph_type='BODY', warn_on_missing=False)
        if ph:
            ph.left = Inches(-100)
            ph.top = Inches(-100)

    # Timeline visual parameters
    n_items = len(timeline)
    if n_items == 0:
        return

    # Position timeline in content area with margin-safe positioning
    safe = get_safe_area()
    line_left = Inches(safe.left)
    line_right = Inches(safe.right)
    line_width = line_right - line_left
    line_y = Inches(3.8)  # Vertical center of content area
    line_height = Pt(3)
    marker_size = Pt(14)
    milestone_marker_size = Pt(18)

    # Draw connecting line
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        line_left, line_y - line_height // 2,
        line_width, line_height
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line_shape.line.fill.background()

    # Calculate spacing between markers
    if n_items > 1:
        spacing = line_width / (n_items - 1)
    else:
        spacing = 0

    # Draw markers and labels
    for i, entry in enumerate(timeline):
        date = entry.get('date', '')
        activity = entry.get('activity', '')
        is_milestone = entry.get('milestone', False) or entry.get('highlight', False) or entry.get('is_milestone', False)

        # Calculate x position
        if n_items > 1:
            x = line_left + (spacing * i)
        else:
            x = line_left + line_width / 2

        # Determine marker size and color
        if is_milestone:
            size = milestone_marker_size
            color = RGBColor(0, 128, 0)  # Green for milestones
        else:
            size = marker_size
            color = RGBColor(0, 0, 0)  # Black for regular

        # Draw circle marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - size / 2, line_y - size / 2,
            size, size
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        # Add date label above the line
        date_box_width = Inches(1.2)
        date_box_height = Inches(0.5)
        date_box = slide.shapes.add_textbox(
            x - date_box_width / 2,
            line_y - Inches(0.9),
            date_box_width,
            date_box_height
        )
        tf = date_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = clean_markdown(date)
        run.font.size = Pt(9)
        run.font.bold = True
        if is_milestone:
            run.font.color.rgb = RGBColor(0, 128, 0)

        # Add activity label below the line
        # Uses auto-font-sizing to fit text within constrained width
        activity_box_width = Inches(1.6)
        activity_box_height = Inches(1.2)
        activity_box = slide.shapes.add_textbox(
            x - activity_box_width / 2,
            line_y + Inches(0.4),
            activity_box_width,
            activity_box_height
        )
        tf = activity_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = clean_markdown(activity)
        run.font.size = Pt(8)
        if is_milestone:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 128, 0)

        # Auto-shrink font if text overflows
        if _HAS_CONTENT_FITTER:
            try:
                fit_result = shrink_textbox_to_fit(
                    activity_box,
                    font_name=get_brand_font('body'),
                    intended_size_pt=8,
                    floor_pct=0.625,  # Allow shrink to 5pt minimum (8 * 0.625 = 5)
                )
                if fit_result.get('shrunk'):
                    # Re-apply milestone styling after shrinking
                    if is_milestone:
                        for para in tf.paragraphs:
                            for r in para.runs:
                                r.font.bold = True
                                r.font.color.rgb = RGBColor(0, 128, 0)
            except Exception:
                pass  # Graceful degradation - keep original size on error

    # Add footnote/note
    footnote = extras.get('footnote') or extras.get('note') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_deliverables_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle deliverables visualization.

    Each deliverable has title + description + optional image_placeholder.
    Supports multiple input patterns (see extract_deliverables).
    """
    content = slide_data.get('content', {})
    deliverables = extract_deliverables(slide_data)
    extras = slide_data.get('extras', {})

    title = content.get('title', '')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    if not deliverables or not body_phs:
        return

    # If multiple body placeholders, distribute deliverables
    if len(body_phs) >= len(deliverables):
        for i, deliv in enumerate(deliverables):
            if i < len(body_phs):
                ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                deliv_title = deliv.get('title', '')
                deliv_desc = deliv.get('description', '')
                # Handle numbered deliverables
                if deliv.get('number'):
                    deliv_title = f"{deliv['number']}. {deliv_title}"
                set_deliverable_content(ph, deliv_title, deliv_desc)

                # Handle per-deliverable image placeholder or actual image
                deliv_image = deliv.get('image_placeholder', '')
                deliv_file_path = deliv.get('file_path')
                if (deliv_image or deliv_file_path) and i < len(picture_phs):
                    pic_ph = picture_phs[i]
                    pos = pic_ph.get('position', {})
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2.5)),
                        deliv_image or "Insert image"
                    )
    else:
        # Combine all deliverables into single body placeholder
        ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
        if ph:
            tf = ph.text_frame
            tf.clear()
            
            first_para = True
            for deliv in deliverables:
                # Title
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                
                deliv_title = deliv.get('title', '')
                if deliv.get('number'):
                    deliv_title = f"{deliv['number']}. {deliv_title}"
                
                run = p.add_run()
                run.text = clean_markdown(deliv_title)
                run.font.bold = True
                p.level = 0
                
                # Description
                desc = deliv.get('description', '')
                if desc:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = clean_markdown(desc)
                    p.level = 1
    
    # Add footnote
    footnote = extras.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_contact_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle contact/next steps slide.

    Title + action items + contact info
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})

    title = content.get('title', '')
    body = content.get('body', [])
    # Check both content.contact and extras.contact
    contact = content.get('contact', {}) or extras.get('contact', {})

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders)

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    # Build contact body content
    string_body = [b for b in body if isinstance(b, str)]
    body_with_contact = list(string_body)

    # Only add from contact dict if body is empty (avoid duplication)
    if contact and not body_with_contact:
        name = contact.get('name', '')
        ctitle = contact.get('title', '')  # renamed to avoid shadowing
        email = contact.get('email', '')
        phone = contact.get('phone', '')
        # Strip leading "- " from name/title if present (from parser)
        if name and name.startswith('- '):
            name = name[2:]
        if ctitle and ctitle.startswith('- '):
            ctitle = ctitle[2:]
        if name or email:
            if name:
                body_with_contact.append(name)
            if ctitle:
                body_with_contact.append(ctitle)
            if email:
                body_with_contact.append(email)
            if phone:
                body_with_contact.append(phone)

    # Set body content if we have anything to show
    if body_with_contact and body_phs:
        ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
        if ph:
            set_header_and_bullets(ph, None, body_with_contact)


def populate_closing_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle closing statement slides.

    These often use title-centered layouts (which have only a title placeholder).
    For hero-statement slides, the statement is added as centered text (like quotes).
    Renders: title, hero statement (centered), body bullets (if body placeholder exists), quote

    Supports styled runs for typography (headline_styled, body_styled, quote_styled).
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    visual_type = slide_data.get('visual_type', '')

    title = content.get('title', '')
    headline = content.get('headline', '')
    body = content.get('body', [])
    quote = content.get('quote') or extras.get('quote')
    quote_attribution = extras.get('quote_attribution')

    # Extract styled runs (from slide-outline-to-layout parser)
    headline_styled = slide_data.get('headline_styled')
    body_styled = slide_data.get('body_styled')
    quote_styled = slide_data.get('quote_styled')

    # For hero-statement, extract the statement text from extras.contact.name
    # (parser stores non-structured text there for closing content types)
    hero_statement = None
    if visual_type == 'hero-statement':
        contact = extras.get('contact', {})
        if contact and contact.get('name'):
            hero_statement = contact.get('name')
        elif headline:
            hero_statement = headline
            headline = ''  # Don't use headline twice

    # Skip headline if it's a visual marker (not actual content)
    if headline and headline.lower().startswith('visual:'):
        headline = ''
        headline_styled = None

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders)

    # For hero-statement, always use the title as the title (not the statement)
    # For other closing slides, headline can become title if no body placeholder
    use_headline_as_title = (headline or headline_styled) and not body_phs and visual_type != 'hero-statement'

    # Set title
    if title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        if ph:
            if use_headline_as_title and headline_styled:
                # Use styled headline as title
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                render_styled_runs(p, headline_styled, base_bold=True)
            elif use_headline_as_title and headline:
                # Use plain headline as title
                set_text_in_placeholder(ph, headline, bold=True)
            elif title:
                # Use regular title
                set_text_in_placeholder(ph, title)
    elif title and body_phs:
        # Fallback: use largest body for title if no title placeholder
        largest = get_largest_body_placeholder(placeholders)
        if largest:
            ph = get_placeholder_by_idx(slide, largest['idx'])
            set_text_in_placeholder(ph, title, bold=True)
            body_phs = [p for p in body_phs if p['idx'] != largest['idx']]

    # Set body content with styled runs (only if there's a body placeholder)
    string_body = [b for b in body if isinstance(b, str)]
    if body_phs and (headline or string_body or headline_styled or body_styled):
        # Only render headline in body if we didn't use it as title
        h = None if use_headline_as_title else headline
        h_styled = None if use_headline_as_title else headline_styled

        ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
        if ph:
            set_header_and_bullets(
                ph, h, string_body,
                header_styled=h_styled,
                body_styled=body_styled
            )

    # Add hero statement as large, centered text (distinct from quotes)
    if hero_statement and visual_type == 'hero-statement':
        add_hero_statement(slide, hero_statement, styled_runs=headline_styled)

    # Add quote with styled runs if available
    # Use center position for quote-hero slides, bottom for others
    if quote or quote_styled:
        quote_position = "center" if visual_type == 'quote-hero' else "bottom"
        add_quote(slide, quote, position=quote_position, styled_runs=quote_styled, attribution=quote_attribution)


def populate_comparison_tables_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle comparison-tables slides (two tables side-by-side).

    Uses title-centered layout with manually positioned textbox headers and tables.
    Matches reference slide 21 "Investment" pattern from chips-analysis.json.

    Positions (based on reference):
    - Left table: left=0.17, top=2.18, width=5.853
    - Right table: left=6.864, top=2.18, width=5.853
    - Headers: 0.6 inches above tables at top=1.584
    """
    content = slide_data.get('content', {})
    table_blocks = slide_data.get('table_blocks', [])
    extras = slide_data.get('extras', {})

    title = content.get('title', '')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    if not table_blocks:
        return

    # Fixed positions for two-column table layout using safe margins
    # Slide is 13.333" wide, with 0.25" (18pt) margins
    safe = get_safe_area()
    # Calculate column widths: (total_width - gap) / 2
    gap = 0.5  # Gap between tables
    col_width = (safe.width - gap) / 2
    positions = [
        {'left': safe.left, 'header_top': 1.584, 'table_top': 2.18, 'width': col_width},
        {'left': safe.left + col_width + gap, 'header_top': 1.584, 'table_top': 2.18, 'width': col_width},
    ]

    # Calculate available height (slide is 7.5" tall, leave margin at bottom)
    table_height = 4.2  # Based on reference

    # Render up to 2 table blocks side-by-side
    for i, block in enumerate(table_blocks[:2]):
        pos = positions[i]

        block_header = block.get('header', '')
        table_data = block.get('table', {})

        if not table_data:
            continue

        # Add header textbox above table
        if block_header:
            # Increased height from 0.252" to 0.4" to properly fit 11pt text
            # Fixes zero-height overflow in comparison-tables visual type (10 cases)
            header_height = Inches(0.4)
            textbox = slide.shapes.add_textbox(
                Inches(pos['left']),
                Inches(pos['header_top']),
                Inches(pos['width']),
                header_height
            )
            tf = textbox.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = clean_markdown(block_header)
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0, 0, 0)

        # Create table at fixed position
        create_table_at_position(
            slide,
            Inches(pos['left']),
            Inches(pos['table_top']),
            Inches(pos['width']),
            Inches(table_height),
            table_data
        )

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_story_card_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle story-card slides (full-bleed image with text overlay).

    Uses content-image-right-b layout with full-bleed PICTURE placeholder.
    Supports styled runs for typography (body_styled, headline_styled).
    """
    content = slide_data.get('content', {})
    background = slide_data.get('background', '')
    extras = slide_data.get('extras', {})

    title = content.get('title', '')
    headline = content.get('headline', '')
    body = content.get('body', [])
    quote = content.get('quote') or extras.get('quote')

    # Get styled runs if available
    body_styled = slide_data.get('body_styled', [])
    headline_styled = slide_data.get('headline_styled')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']

    # Set title (if there's a title placeholder)
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    # Handle background image placeholder or actual image
    # background can be: a file path (if exists), or a description
    background_file = extras.get('background_file')  # Explicit file path
    if background and picture_phs:
        pic_ph = picture_phs[0]
        pos = pic_ph.get('position', {})
        # Move original placeholder off-slide
        pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
        if pic_shape:
            pic_shape.left = Inches(-100)
            pic_shape.top = Inches(-100)
        # Check if background is a file path (has extension) or description
        from pathlib import Path as PathLib
        bg_path = background_file or background
        is_file_path = PathLib(bg_path).suffix.lower() in {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif'}
        insert_image_or_placeholder(
            slide,
            bg_path if is_file_path else None,
            Inches(pos.get('left_inches', 0)),
            Inches(pos.get('top_inches', 0)),
            Inches(pos.get('width_inches', 13.333)),
            Inches(pos.get('height_inches', 7.5)),
            background if not is_file_path else "Background image",
            send_to_back=True  # Story-card: image goes behind text overlay
        )

    # Set content (headline + body + quote)
    string_body = [b for b in body if isinstance(b, str)]
    if body_phs:
        ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
        if ph:
            tf = ph.text_frame
            tf.clear()

            first_para = True

            # Add headline - use styled runs if available (for blue color)
            if headline:
                p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                first_para = False
                p.level = 0

                # Check if headline has styled runs (color markers)
                if headline_styled:
                    render_styled_runs(p, headline_styled, base_bold=True)
                elif '{blue}' in headline or '{italic}' in headline:
                    # Parse inline style markers
                    add_styled_paragraph_with_markers(p, headline, base_bold=True)
                else:
                    run = p.add_run()
                    run.text = clean_markdown(headline)
                    run.font.bold = True

            # Add body paragraphs - use body_styled if available
            if body_styled:
                for item in body_styled:
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.level = 0

                    styled_runs = item.get('runs') if isinstance(item, dict) else None
                    plain_text = item.get('text', '') if isinstance(item, dict) else str(item)

                    if styled_runs:
                        render_styled_runs(p, styled_runs)
                    else:
                        run = p.add_run()
                        run.text = clean_markdown(plain_text)
            else:
                for item in string_body:
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.level = 0

                    # Check for inline style markers
                    if '{italic}' in item or '{blue}' in item or '{bold}' in item:
                        add_styled_paragraph_with_markers(p, item)
                    elif '**' in item:
                        # Use markdown parsing for **bold**
                        segments = parse_markdown_formatting(item)
                        for seg in segments:
                            run = p.add_run()
                            run.text = seg['text']
                            run.font.bold = seg['bold']
                            run.font.italic = seg['italic']
                    else:
                        run = p.add_run()
                        run.text = clean_markdown(item)

            # Add quote (italic)
            if quote:
                if not first_para:
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                run = p.add_run()
                quote_text = clean_markdown(quote)
                if not quote_text.startswith('"') and not quote_text.startswith('"'):
                    quote_text = f'"{quote_text}"'
                run.text = quote_text
                run.font.italic = True
                p.level = 0


def populate_grid_3x2_image_top_3_body(slide, placeholders: list, slide_data: dict):
    """
    Handle grid-3x2-image-top-3-body layout (layout 36).

    This layout has:
    - Title placeholder at top
    - 3 columns, each with an image area (top) and body text (bottom)
    - PICTURE placeholders: idx 65 (middle), 68 (right), and potentially 64 (left)
    - BODY placeholders: idx 16 (left), 48 (middle), 49 (right)

    Used for: Research objectives, deliverables with visuals, 3-item grids.

    Content sources (in priority order):
    1. slide_data['deliverables'] or slide_data['content']['deliverables']
    2. slide_data['columns'] with image_placeholder
    3. slide_data['content']['body'] containing objects with title/description
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)
    deliverables = extract_deliverables(slide_data)

    title = content.get('title', '')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    # If we have structured columns with image placeholders, use them
    if columns:
        for col_idx, column in enumerate(columns):
            if col_idx >= len(body_phs):
                break

            col_header = column.get('header', '')
            col_subheader = column.get('subheader', '')
            col_body = column.get('body', [])
            col_image = column.get('image_placeholder', '')
            col_file_path = column.get('file_path')

            # Set body content
            ph = get_placeholder_by_idx(slide, body_phs[col_idx]['idx'])
            if ph:
                body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                set_header_and_bullets(ph, col_header, body_items, subheader=col_subheader, header_bold=True)

            # Handle image placeholder or actual image
            if col_image or col_file_path:
                if col_idx < len(picture_phs):
                    # Use existing PICTURE placeholder position
                    pic_ph = picture_phs[col_idx]
                    pos = pic_ph.get('position', {})
                    # Move original placeholder off-slide
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2.5)),
                        col_image or "Insert image"
                    )
                elif col_idx < len(body_phs):
                    # Calculate position above body placeholder
                    body_pos = body_phs[col_idx].get('position', {})
                    img_pos = calculate_image_position_above_body(body_pos)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 3)),
                        Inches(img_pos.get('height_inches', 2)),
                        col_image or "Insert image"
                    )

    # If we have deliverables, distribute them across body placeholders
    elif deliverables:
        for i, deliv in enumerate(deliverables):
            if i >= len(body_phs):
                break

            ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
            deliv_title = deliv.get('title', '')
            deliv_desc = deliv.get('description', '')

            # Handle numbered deliverables
            if deliv.get('number'):
                deliv_title = f"{deliv['number']}. {deliv_title}"

            set_deliverable_content(ph, deliv_title, deliv_desc)

            # Check if deliverable has an image placeholder or file path
            deliv_image = deliv.get('image_placeholder', '')
            deliv_file_path = deliv.get('file_path')
            if deliv_image or deliv_file_path:
                if i < len(picture_phs):
                    pic_ph = picture_phs[i]
                    pos = pic_ph.get('position', {})
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2.5)),
                        deliv_image or "Insert image"
                    )
                elif i < len(body_phs):
                    # Calculate position above body placeholder
                    body_pos = body_phs[i].get('position', {})
                    img_pos = calculate_image_position_above_body(body_pos)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 3)),
                        Inches(img_pos.get('height_inches', 2)),
                        deliv_image or "Insert image"
                    )

    # Fallback: use body content if present
    else:
        body = content.get('body', [])
        headline = content.get('headline', '')

        # Check if body contains deliverable-style objects
        if body and isinstance(body, list) and len(body) > 0 and is_deliverable_object(body[0]):
            for i, item in enumerate(body):
                if i >= len(body_phs):
                    break
                ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                item_title = item.get('title', '')
                item_desc = item.get('description', '')
                set_deliverable_content(ph, item_title, item_desc)
        elif body_phs:
            # Standard content: put in first body placeholder
            string_body = [b for b in body if isinstance(b, str)]
            if headline or string_body:
                ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
                if ph:
                    set_header_and_bullets(ph, headline, string_body)

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_grid_3x2_image_top_6_body_a(slide, placeholders: list, slide_data: dict):
    """
    Handle grid-3x2-image-top-6-body-a layout (layout 37).

    This layout has:
    - Title placeholder at top
    - 3 columns, each with:
      - Image area (top) - PICTURE placeholders: idx 64 (left), 65 (middle), 68 (right)
      - Header/label area (middle) - BODY placeholders: idx 41 (left), 69 (middle), 53 (right)
      - Body content area (bottom) - BODY placeholders: idx 16 (left), 48 (middle), 49 (right)

    Used for: Process phases, approach overviews with 3 stages each having header + bullets.

    Content sources:
    1. slide_data['columns'] - each column has header (for label row) + body (for content row)
    2. slide_data['phases'] - alternative name for 3-phase process content
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)
    # Also check for 'phases' as alternative name
    if not columns:
        columns = slide_data.get('phases', [])

    title = content.get('title', '')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=False)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    # Categorize body placeholders by vertical position
    # Top row (headers) vs bottom row (content)
    if body_phs:
        # Sort by top position to identify rows
        sorted_by_top = sorted(body_phs, key=lambda p: p.get('position', {}).get('top_inches', 0))

        # Find the midpoint between top and bottom rows
        top_positions = [p.get('position', {}).get('top_inches', 0) for p in sorted_by_top]
        if len(top_positions) >= 4:
            # Find the gap between header row and content row
            mid_y = (top_positions[2] + top_positions[3]) / 2
        else:
            mid_y = 4.5  # Default midpoint

        # Split into header row (top) and content row (bottom)
        header_row_phs = [p for p in body_phs if p.get('position', {}).get('top_inches', 0) < mid_y]
        content_row_phs = [p for p in body_phs if p.get('position', {}).get('top_inches', 0) >= mid_y]

        # Sort each row left-to-right
        header_row_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))
        content_row_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))
    else:
        header_row_phs = []
        content_row_phs = []

    # Process columns
    if columns:
        for col_idx, column in enumerate(columns):
            col_header = column.get('header', '')
            col_subheader = column.get('subheader', '')
            col_intro = column.get('intro', [])  # List of intro lines (header + italic subtext)
            col_body = column.get('body', [])
            col_image = column.get('image_placeholder', '')
            col_file_path = column.get('file_path')

            # Set header in top row placeholder
            # Use intro list if available (contains header line + italic subtext line)
            if col_idx < len(header_row_phs):
                ph = get_placeholder_by_idx(slide, header_row_phs[col_idx]['idx'])
                if ph:
                    if col_intro and len(col_intro) >= 2:
                        # First line is header (bold), second line is subtext (italic)
                        header_line = col_intro[0] if col_intro else col_header
                        subtext_line = col_intro[1] if len(col_intro) > 1 else ''
                        # Strip {italic} markers - the subheader handler sets italic automatically
                        import re
                        subtext_line = re.sub(r'\{/?italic\}', '', subtext_line)
                        set_header_and_bullets(ph, header_line, [], subheader=subtext_line, header_bold=True)
                    elif col_subheader:
                        set_header_and_bullets(ph, col_header, [], subheader=col_subheader, header_bold=True)
                    elif col_header:
                        set_text_in_placeholder(ph, col_header, bold=True)

            # Set body content in bottom row placeholder
            if col_idx < len(content_row_phs):
                ph = get_placeholder_by_idx(slide, content_row_phs[col_idx]['idx'])
                if ph:
                    body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                    if body_items:
                        set_header_and_bullets(ph, None, body_items)

            # Handle image placeholder or actual image
            if col_image or col_file_path:
                if col_idx < len(picture_phs):
                    pic_ph = picture_phs[col_idx]
                    pos = pic_ph.get('position', {})
                    # Move original placeholder off-slide
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2)),
                        col_image or "Insert image"
                    )
                elif col_idx < len(header_row_phs):
                    # Calculate position above header row placeholder
                    header_pos = header_row_phs[col_idx].get('position', {})
                    img_pos = calculate_image_position_above_body(header_pos)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 3)),
                        Inches(img_pos.get('height_inches', 2)),
                        col_image or "Insert image"
                    )

    # Fallback: use content body if no columns
    else:
        body = content.get('body', [])
        headline = content.get('headline', '')

        # Put headline + body in first content placeholder
        if content_row_phs:
            string_body = [b for b in body if isinstance(b, str)]
            if headline or string_body:
                ph = get_placeholder_by_idx(slide, content_row_phs[0]['idx'])
                if ph:
                    set_header_and_bullets(ph, headline, string_body)

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_grid_2x2_image_top_2_body_a(slide, placeholders: list, slide_data: dict):
    """
    Handle grid-2x2-image-top-2-body-a layout (layout 42).

    This layout has:
    - Title placeholder at top
    - 2 columns, each with:
      - Image area (top) - PICTURE placeholders: idx 64 (left), 68 (right)
      - Body content area (bottom) - BODY placeholders: idx 16 (left), 49 (right)

    Used for: Two-city comparisons, two-segment analyses, side-by-side content with images.

    Content sources:
    1. slide_data['columns'] - each column has header + body + optional image_placeholder
    2. slide_data['content']['body'] containing objects with title/description
    3. Standard headline + body content
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)
    deliverables = extract_deliverables(slide_data)

    title = content.get('title', '')
    title_styled = slide_data.get('title_styled')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title (use styled runs if available for phase prefix coloring)
    if (title or title_styled) and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title, styled_runs=title_styled)

    # If we have structured columns, use them
    if columns:
        for col_idx, column in enumerate(columns):
            if col_idx >= len(body_phs):
                break

            col_header = column.get('header', '')
            col_subheader = column.get('subheader', '')
            col_intro = column.get('intro', [])
            col_body = column.get('body', [])
            col_image = column.get('image_placeholder', '')
            col_file_path = column.get('file_path')

            # Use intro[0] as header if it has timing info (e.g., "Shanghai IDI (90 mins)")
            # This is because ingest extracts short header from [Column: X] but full header is in bold line
            if col_intro and isinstance(col_intro, list) and col_intro:
                col_header = col_intro[0]

            # Set body content with blank line after header
            ph = get_placeholder_by_idx(slide, body_phs[col_idx]['idx'])
            if ph:
                body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                set_header_and_bullets(ph, col_header, body_items, subheader=col_subheader, header_bold=True,
                                       add_blank_after_header=True)

            # Handle image placeholder or actual image
            if col_image or col_file_path:
                if col_idx < len(picture_phs):
                    pic_ph = picture_phs[col_idx]
                    pos = pic_ph.get('position', {})
                    # Move original placeholder off-slide
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 5.8)),
                        Inches(pos.get('height_inches', 2.9)),
                        col_image or "Insert image"
                    )
                elif col_idx < len(body_phs):
                    # Calculate position above body placeholder
                    body_pos = body_phs[col_idx].get('position', {})
                    img_pos = calculate_image_position_above_body(body_pos)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 5.8)),
                        Inches(img_pos.get('height_inches', 2.5)),
                        col_image or "Insert image"
                    )

    # If we have deliverables, distribute them across body placeholders
    elif deliverables:
        for i, deliv in enumerate(deliverables):
            if i >= len(body_phs):
                break

            ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
            deliv_title = deliv.get('title', '')
            deliv_desc = deliv.get('description', '')

            # Handle numbered deliverables
            if deliv.get('number'):
                deliv_title = f"{deliv['number']}. {deliv_title}"

            set_deliverable_content(ph, deliv_title, deliv_desc)

            # Check if deliverable has an image placeholder or file path
            deliv_image = deliv.get('image_placeholder', '')
            deliv_file_path = deliv.get('file_path')
            if deliv_image or deliv_file_path:
                if i < len(picture_phs):
                    pic_ph = picture_phs[i]
                    pos = pic_ph.get('position', {})
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 5.8)),
                        Inches(pos.get('height_inches', 2.9)),
                        deliv_image
                    )
                elif i < len(body_phs):
                    # Calculate position above body placeholder
                    body_pos = body_phs[i].get('position', {})
                    img_pos = calculate_image_position_above_body(body_pos)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 5.8)),
                        Inches(img_pos.get('height_inches', 2.5)),
                        deliv_image or "Insert image"
                    )

    # Fallback: use body content if present
    else:
        body = content.get('body', [])
        headline = content.get('headline', '')

        # Check if body contains deliverable-style objects
        if body and isinstance(body, list) and len(body) > 0 and is_deliverable_object(body[0]):
            for i, item in enumerate(body):
                if i >= len(body_phs):
                    break
                ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                item_title = item.get('title', '')
                item_desc = item.get('description', '')
                set_deliverable_content(ph, item_title, item_desc)
        elif body_phs:
            # Standard content: put in first body placeholder
            string_body = [b for b in body if isinstance(b, str)]
            if headline or string_body:
                ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
                if ph:
                    set_header_and_bullets(ph, headline, string_body)

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_grid_2x2_image_top_4_body(slide, placeholders: list, slide_data: dict):
    """
    Handle grid-2x2-image-top-4-body layout (layout 43).

    This layout has:
    - Title placeholder at top (idx 0)
    - 2x2 grid with 4 body areas:
      - Top row: idx 41 (left), idx 58 (right) - smaller label areas
      - Bottom row: idx 16 (left), idx 49 (right) - larger content areas
    - 2 image areas above left/right columns: idx 64 (left), 68 (right)

    Used for: 2x2 matrix with icons, quadrant analysis, side-by-side comparisons with visuals.

    Content sources:
    1. slide_data['columns'] - 2 columns, each with header + body + optional image
    2. slide_data['quadrants'] - 4 quadrant items for true 2x2
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)

    title = content.get('title', '')
    title_styled = slide_data.get('title_styled')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title
    if (title or title_styled) and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title, styled_runs=title_styled)

    # Sort body placeholders by position (top-left, top-right, bottom-left, bottom-right)
    # Top row has smaller height, bottom row has larger height
    if len(body_phs) >= 4:
        # Sort by top position first, then by left
        body_phs.sort(key=lambda p: (p.get('position', {}).get('top_inches', 0),
                                      p.get('position', {}).get('left_inches', 0)))

    # If we have 2 columns, distribute across 2x2 (each column gets top label + bottom content)
    if columns and len(columns) >= 2:
        for col_idx, column in enumerate(columns[:2]):
            col_header = column.get('header', '')
            col_body = column.get('body', [])
            col_image = column.get('image_placeholder', '')
            col_file_path = column.get('file_path')

            # Top row placeholder (label/header)
            top_idx = col_idx  # 0 for left, 1 for right
            if top_idx < len(body_phs):
                ph = get_placeholder_by_idx(slide, body_phs[top_idx]['idx'])
                if ph:
                    set_text_in_placeholder(ph, col_header)

            # Bottom row placeholder (body content)
            bottom_idx = col_idx + 2  # 2 for left, 3 for right
            if bottom_idx < len(body_phs):
                ph = get_placeholder_by_idx(slide, body_phs[bottom_idx]['idx'])
                if ph:
                    body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                    set_header_and_bullets(ph, '', body_items)

            # Handle image
            if col_image or col_file_path:
                if col_idx < len(picture_phs):
                    pic_ph = picture_phs[col_idx]
                    pos = pic_ph.get('position', {})
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2)),
                        col_image or "Insert image"
                    )

    # Fallback: distribute body items across all 4 quadrants
    else:
        body = content.get('body', [])
        if isinstance(body, list):
            for i, item in enumerate(body[:4]):
                if i < len(body_phs):
                    ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                    if ph:
                        if isinstance(item, dict):
                            item_title = item.get('title', '')
                            item_body = item.get('body', [])
                            set_header_and_bullets(ph, item_title, item_body if isinstance(item_body, list) else [item_body])
                        else:
                            set_text_in_placeholder(ph, str(item))

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_grid_2x4(slide, placeholders: list, slide_data: dict):
    """
    Handle grid-2x4 layout (layout 13).

    This layout has:
    - Title placeholder at top (idx 0)
    - 4 body columns (idx 15, 18, 19, 17) arranged horizontally

    Used for: 4-item horizontal grids, phase summaries, category comparisons.

    Content sources:
    1. slide_data['columns'] - 4 columns with header + body
    2. slide_data['content']['body'] - list of 4 items
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)

    title = content.get('title', '')
    title_styled = slide_data.get('title_styled')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)

    # Set title
    if (title or title_styled) and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title, styled_runs=title_styled)

    # If we have columns, use them
    if columns:
        for col_idx, column in enumerate(columns[:4]):
            if col_idx >= len(body_phs):
                break

            col_header = column.get('header', '')
            col_body = column.get('body', [])

            ph = get_placeholder_by_idx(slide, body_phs[col_idx]['idx'])
            if ph:
                body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                set_header_and_bullets(ph, col_header, body_items, header_bold=True, add_blank_after_header=True)

    # Fallback: distribute body items
    else:
        body = content.get('body', [])
        if isinstance(body, list):
            for i, item in enumerate(body[:4]):
                if i < len(body_phs):
                    ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                    if ph:
                        if isinstance(item, dict):
                            item_title = item.get('title', '')
                            item_body = item.get('body', [])
                            set_header_and_bullets(ph, item_title, item_body if isinstance(item_body, list) else [item_body])
                        else:
                            set_text_in_placeholder(ph, str(item))

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_grid_3x2_text_top_4_body(slide, placeholders: list, slide_data: dict):
    """
    Handle grid-3x2-text-top-4-body layout (layout 16).

    This layout has:
    - 4 body placeholders arranged in a 3x2 pattern:
      - idx 19: top-left text area (header row)
      - idx 18: middle-left text area
      - idx 20: right column 1 (tall)
      - idx 21: right column 2 (tall)
    - 2 image placeholders: idx 22, 23 (bottom-left area)

    Used for: Complex grids with header/description on left, detail columns on right.

    Content sources:
    1. slide_data['columns'] - structured column data
    2. slide_data['content']['body'] - list items
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)

    # Get placeholders sorted by position
    body_phs = get_body_placeholders(placeholders, sort_horizontal=False)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Sort body placeholders: top-to-bottom, then left-to-right
    body_phs.sort(key=lambda p: (p.get('position', {}).get('top_inches', 0),
                                  p.get('position', {}).get('left_inches', 0)))

    # If we have columns, distribute content
    if columns:
        for col_idx, column in enumerate(columns[:4]):
            if col_idx >= len(body_phs):
                break

            col_header = column.get('header', '')
            col_body = column.get('body', [])
            col_image = column.get('image_placeholder', '')
            col_file_path = column.get('file_path')

            ph = get_placeholder_by_idx(slide, body_phs[col_idx]['idx'])
            if ph:
                body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                set_header_and_bullets(ph, col_header, body_items, header_bold=True, add_blank_after_header=True)

            # Handle images for first two columns (left side)
            if (col_image or col_file_path) and col_idx < len(picture_phs):
                pic_ph = picture_phs[col_idx]
                pos = pic_ph.get('position', {})
                pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                if pic_shape:
                    pic_shape.left = Inches(-100)
                    pic_shape.top = Inches(-100)
                insert_image_or_placeholder(
                    slide,
                    col_file_path,
                    Inches(pos.get('left_inches', 0)),
                    Inches(pos.get('top_inches', 0)),
                    Inches(pos.get('width_inches', 3)),
                    Inches(pos.get('height_inches', 3)),
                    col_image or "Insert image"
                )

    # Fallback: distribute body items
    else:
        body = content.get('body', [])
        if isinstance(body, list):
            for i, item in enumerate(body[:4]):
                if i < len(body_phs):
                    ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                    if ph:
                        if isinstance(item, str):
                            set_text_in_placeholder(ph, item)
                        elif isinstance(item, dict):
                            item_title = item.get('title', '')
                            item_body = item.get('body', [])
                            set_header_and_bullets(ph, item_title, item_body if isinstance(item_body, list) else [item_body])

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_content_image_top_4_body(slide, placeholders: list, slide_data: dict):
    """
    Handle content-image-top-4-body layout (layout 32).

    This layout has:
    - Title placeholder at top
    - 4 columns, each with:
      - Image area (top) - PICTURE placeholders: idx 64 (col 1), 65 (col 2), 66 (col 3), 68 (col 4)
      - Body content area (bottom) - BODY placeholders: idx 16 (col 1), 47 (col 2), 48 (col 3), 49 (col 4)

    Used for: 4-item deliverables with images, phase overviews with 4 stages.

    Content sources:
    1. slide_data['deliverables'] or slide_data['content']['deliverables']
    2. slide_data['columns'] with image_placeholder
    3. slide_data['cards'] (converted to deliverables format)
    4. slide_data['content']['body'] containing objects with title/description
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})
    columns = extract_columns(slide_data)
    deliverables = extract_deliverables(slide_data)

    title = content.get('title', '')

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders, sort_horizontal=True)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']
    picture_phs.sort(key=lambda p: p.get('position', {}).get('left_inches', 0))

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    # If we have structured columns with image placeholders, use them
    if columns:
        for col_idx, column in enumerate(columns):
            if col_idx >= len(body_phs):
                break

            col_header = column.get('header', '')
            col_subheader = column.get('subheader', '')
            col_body = column.get('body', [])
            col_image = column.get('image_placeholder', '')
            col_file_path = column.get('file_path')

            # Set body content
            ph = get_placeholder_by_idx(slide, body_phs[col_idx]['idx'])
            if ph:
                body_items = list(col_body) if isinstance(col_body, list) else [col_body] if col_body else []
                set_header_and_bullets(ph, col_header, body_items, subheader=col_subheader, header_bold=True)

            # Handle image placeholder or actual image
            if col_image or col_file_path:
                if col_idx < len(picture_phs):
                    pic_ph = picture_phs[col_idx]
                    pos = pic_ph.get('position', {})
                    # Move original placeholder off-slide
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2.9)),
                        col_image or "Insert image"
                    )
                elif col_idx < len(body_phs):
                    # Calculate position above body placeholder
                    body_pos = body_phs[col_idx].get('position', {})
                    img_pos = calculate_image_position_above_body(body_pos)
                    insert_image_or_placeholder(
                        slide,
                        col_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 3)),
                        Inches(img_pos.get('height_inches', 2)),
                        col_image or "Insert image"
                    )

    # If we have deliverables, distribute them across body placeholders
    elif deliverables:
        for i, deliv in enumerate(deliverables):
            if i >= len(body_phs):
                break

            ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
            deliv_title = deliv.get('title', '')
            deliv_desc = deliv.get('description', '')

            # Handle numbered deliverables
            if deliv.get('number'):
                deliv_title = f"{deliv['number']}. {deliv_title}"

            set_deliverable_content(ph, deliv_title, deliv_desc)

            # Check if deliverable has an image placeholder or file path
            deliv_image = deliv.get('image_placeholder', '')
            deliv_file_path = deliv.get('file_path')
            if deliv_image or deliv_file_path:
                if i < len(picture_phs):
                    pic_ph = picture_phs[i]
                    pos = pic_ph.get('position', {})
                    pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
                    if pic_shape:
                        pic_shape.left = Inches(-100)
                        pic_shape.top = Inches(-100)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(pos.get('left_inches', 0)),
                        Inches(pos.get('top_inches', 0)),
                        Inches(pos.get('width_inches', 3)),
                        Inches(pos.get('height_inches', 2.9)),
                        deliv_image or "Insert image"
                    )
                elif i < len(body_phs):
                    # Calculate position above body placeholder
                    body_pos = body_phs[i].get('position', {})
                    img_pos = calculate_image_position_above_body(body_pos)
                    insert_image_or_placeholder(
                        slide,
                        deliv_file_path,
                        Inches(img_pos.get('left_inches', 0)),
                        Inches(img_pos.get('top_inches', 1.5)),
                        Inches(img_pos.get('width_inches', 3)),
                        Inches(img_pos.get('height_inches', 2)),
                        deliv_image or "Insert image"
                    )

    # Fallback: use body content if present
    else:
        body = content.get('body', [])
        headline = content.get('headline', '')

        # Check if body contains deliverable-style objects
        if body and isinstance(body, list) and len(body) > 0 and is_deliverable_object(body[0]):
            for i, item in enumerate(body):
                if i >= len(body_phs):
                    break
                ph = get_placeholder_by_idx(slide, body_phs[i]['idx'])
                item_title = item.get('title', '')
                item_desc = item.get('description', '')
                set_deliverable_content(ph, item_title, item_desc)
        elif body_phs:
            # Standard content: put in first body placeholder
            string_body = [b for b in body if isinstance(b, str)]
            if headline or string_body:
                ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
                if ph:
                    set_header_and_bullets(ph, headline, string_body)

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


def populate_table_with_image_slide(slide, placeholders: list, slide_data: dict):
    """
    Handle table-with-image slides.

    Two modes:
    1. Body content on left, table on right (when body content exists)
    2. Table in body area, image placeholder on right (legacy mode)

    Uses content-image-right-a layout.
    """
    content = slide_data.get('content', {})
    extras = slide_data.get('extras', {})

    title = content.get('title', '')
    headline = content.get('headline', '')
    body = content.get('body', [])
    image_placeholder = extras.get('image_placeholder', '')
    image_file_path = extras.get('image_file')  # Path to actual image

    # Get placeholders
    title_ph = get_title_placeholder(placeholders)
    body_phs = get_body_placeholders(placeholders)
    picture_phs = [p for p in placeholders if p.get('type') == 'PICTURE']

    # Set title
    if title and title_ph:
        ph = get_placeholder_by_idx(slide, title_ph['idx'])
        set_text_in_placeholder(ph, title)

    # Handle table
    tables = extract_tables(slide_data)

    # Check if we have body content - if so, use new mode (body left, table right)
    if body and tables:
        # Mode 1: Body content in body placeholder, table in picture area
        string_body = [b for b in body if isinstance(b, str)]
        if body_phs and string_body:
            ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
            if ph:
                set_header_and_bullets(ph, headline, string_body)

        # Put table where picture placeholder would be
        if picture_phs:
            pic_ph = picture_phs[0]
            pos = pic_ph.get('position', {})
            # Move original placeholder off-slide
            pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
            if pic_shape:
                pic_shape.left = Inches(-100)
                pic_shape.top = Inches(-100)
            # Create table in picture area
            create_table_at_position(
                slide,
                Inches(pos.get('left_inches', 7.5)),
                Inches(pos.get('top_inches', 1.5)),
                Inches(pos.get('width_inches', 5.5)),
                Inches(pos.get('height_inches', 4.5)),
                tables[0]
            )
    elif tables and body_phs:
        # Mode 2 (legacy): Table in body area
        ph = get_placeholder_by_idx(slide, body_phs[0]['idx'])
        if ph:
            create_table_in_placeholder(slide, ph, tables[0])

        # Handle image placeholder or actual image
        if (image_placeholder or image_file_path) and picture_phs:
            pic_ph = picture_phs[0]
            pos = pic_ph.get('position', {})
            # Move original placeholder off-slide
            pic_shape = get_placeholder_by_idx(slide, pic_ph['idx'])
            if pic_shape:
                pic_shape.left = Inches(-100)
                pic_shape.top = Inches(-100)
            # Insert image or create placeholder box
            insert_image_or_placeholder(
                slide,
                image_file_path,
                Inches(pos.get('left_inches', 0)),
                Inches(pos.get('top_inches', 0)),
                Inches(pos.get('width_inches', 4)),
                Inches(pos.get('height_inches', 3)),
                image_placeholder or "Insert image"
            )

    # Add footnote
    footnote = extras.get('footnote') or content.get('footnote')
    if footnote:
        add_footnote(slide, footnote)


# =============================================================================
# MAIN SLIDE ROUTING
# =============================================================================

def populate_slide(slide, slide_data: dict):
    """Populate a single slide with content based on content_type field.
    
    Primary routing is based on content_type field.
    Secondary routing uses visual_type and data presence as fallback.
    """
    layout_info = slide_data.get('layout', {})
    layout_name = layout_info.get('name', '')
    content = slide_data.get('content', {})
    content_type = slide_data.get('content_type', '')
    visual_type = slide_data.get('visual_type', '')
    
    # Skip branding slides (no content)
    if content_type == 'branding' or layout_name == 'master-base':
        return
    
    # Discover placeholders from the actual slide at runtime
    placeholders = discover_placeholders(slide)
    
    # Skip if no placeholders discovered
    if not placeholders:
        return
    
    # =========================================================================
    # PRIMARY ROUTING: Based on content_type field
    # =========================================================================
    
    # Title/cover slides
    if content_type == 'title_slide':
        populate_title_cover_slide(slide, placeholders, slide_data)
        return
    
    # Contact slides
    if content_type == 'contact':
        populate_contact_slide(slide, placeholders, slide_data)
        return
    
    # Closing slides
    if content_type == 'closing':
        populate_closing_slide(slide, placeholders, slide_data)
        return

    # Quote slides (use closing handler which supports quote + attribution)
    if content_type == 'quote':
        populate_closing_slide(slide, placeholders, slide_data)
        return

    # Section divider slides (same handler as closing)
    if content_type == 'section_divider':
        populate_closing_slide(slide, placeholders, slide_data)
        return

    # Hero-statement slides (single punchy statement, centered)
    if content_type == 'hero-statement':
        populate_closing_slide(slide, placeholders, slide_data)
        return

    # Timeline slides
    if content_type == 'timeline':
        populate_timeline_slide(slide, placeholders, slide_data)
        return
    
    # Deliverables slides - only route if actual deliverables/cards exist
    if content_type == 'deliverables':
        # Only use deliverables handler if there's actual deliverable/card data
        if extract_deliverables(slide_data):
            populate_deliverables_slide(slide, placeholders, slide_data)
            return
        # Otherwise fall through to content handler

    # Table slides (full-width tables)
    if content_type == 'table':
        populate_content_slide(slide, placeholders, slide_data)
        return

    # Comparison tables (side-by-side tables)
    if content_type == 'comparison_tables':
        populate_comparison_tables_slide(slide, placeholders, slide_data)
        return

    # Story card (full-bleed image with text overlay)
    if content_type == 'story_card':
        populate_story_card_slide(slide, placeholders, slide_data)
        return

    # Grid 2x2 content type - routes to appropriate 2x2 handler based on layout
    if content_type == 'grid_2x2':
        if 'image-top-4-body' in layout_name:
            populate_grid_2x2_image_top_4_body(slide, placeholders, slide_data)
        elif 'image-top-2-body' in layout_name:
            populate_grid_2x2_image_top_2_body_a(slide, placeholders, slide_data)
        else:
            # Default to 4-body variant
            populate_grid_2x2_image_top_4_body(slide, placeholders, slide_data)
        return

    # Grid 2x4 content type - 4 horizontal columns
    if content_type == 'grid_2x4':
        populate_grid_2x4(slide, placeholders, slide_data)
        return

    # Grid 3x2 content types - route to appropriate handler
    if content_type in ('grid_3x2', 'grid_3x2_6body'):
        if 'text-top-4-body' in layout_name:
            populate_grid_3x2_text_top_4_body(slide, placeholders, slide_data)
        elif 'image-top-3-body' in layout_name:
            populate_grid_3x2_image_top_3_body(slide, placeholders, slide_data)
        else:
            # Default to 6-body variant
            populate_grid_3x2_image_top_6_body_a(slide, placeholders, slide_data)
        return

    # Grid 3x2 with 3 body areas
    if content_type == 'grid_3x2_3body':
        populate_grid_3x2_image_top_3_body(slide, placeholders, slide_data)
        return

    # Grid 3x2 with text header row
    if content_type == 'grid_3x2_text_top':
        populate_grid_3x2_text_top_4_body(slide, placeholders, slide_data)
        return

    # Grid 2x2 with 4 body areas
    if content_type == 'grid_2x2_4body':
        populate_grid_2x2_image_top_4_body(slide, placeholders, slide_data)
        return

    # Grid 2x2 with 2 body areas
    if content_type in ('grid_2x2_2body', 'grid_2x2_2body_b'):
        populate_grid_2x2_image_top_2_body_a(slide, placeholders, slide_data)
        return

    # Table with image placeholder
    if content_type == 'table_with_image':
        populate_table_with_image_slide(slide, placeholders, slide_data)
        return

    # Grid 3x2 with images top, 3 body areas (layout 36)
    if layout_name == 'grid-3x2-image-top-3-body' or visual_type == 'grid-3x2-image-top-3-body':
        populate_grid_3x2_image_top_3_body(slide, placeholders, slide_data)
        return

    # Grid 3x2 with images top, 6 body areas (layout 37)
    if layout_name == 'grid-3x2-image-top-6-body-a' or visual_type == 'grid-3x2-image-top-6-body-a':
        populate_grid_3x2_image_top_6_body_a(slide, placeholders, slide_data)
        return

    # Grid 2x2 with images top, 2 body areas (layout 42 or 44)
    if layout_name in ('grid-2x2-image-top-2-body-a', 'grid-2x2-image-top-2-body-b') or \
       visual_type in ('grid-2x2-image-top-2-body-a', 'grid-2x2-image-top-2-body-b'):
        populate_grid_2x2_image_top_2_body_a(slide, placeholders, slide_data)
        return

    # Grid 2x2 with images top, 4 body areas (layout 43)
    if layout_name == 'grid-2x2-image-top-4-body' or visual_type == 'grid-2x2-image-top-4-body':
        populate_grid_2x2_image_top_4_body(slide, placeholders, slide_data)
        return

    # Grid 2x4 - 4 horizontal columns (layout 13)
    if layout_name == 'grid-2x4' or visual_type == 'grid-2x4':
        populate_grid_2x4(slide, placeholders, slide_data)
        return

    # Grid 3x2 with text top, 4 body areas (layout 16)
    if layout_name == 'grid-3x2-text-top-4-body' or visual_type == 'grid-3x2-text-top-4-body':
        populate_grid_3x2_text_top_4_body(slide, placeholders, slide_data)
        return

    # Content with images top, 4 body areas (layout 32)
    if layout_name == 'content-image-top-4-body' or visual_type == 'content-image-top-4-body':
        populate_content_image_top_4_body(slide, placeholders, slide_data)
        return

    # 1_content-image-right-a variant (layout 58) - same as content-image-right-a but no picture placeholder
    if layout_name == '1_content-image-right-a':
        populate_content_slide(slide, placeholders, slide_data)
        return

    # Auto-detect deliverables: body contains objects with title/description
    body = content.get('body', [])
    if body and isinstance(body, list) and len(body) > 0 and is_deliverable_object(body[0]):
        populate_deliverables_slide(slide, placeholders, slide_data)
        return

    # Data contrast visual type (side-by-side metrics with interpretation)
    if visual_type == 'data-contrast' and slide_data.get('data_contrast'):
        populate_data_contrast_slide(slide, placeholders, slide_data)
        return

    # Cards visual type (cards-N) - use column handler (cards converted to columns)
    if visual_type and visual_type.startswith('cards-') and slide_data.get('cards'):
        populate_column_slide(slide, placeholders, slide_data)
        return

    # Framework/comparison/process slides (column layouts)
    # Handles both generic types and specific column counts
    framework_types = (
        'framework', 'framework_2col', 'framework_3col', 'framework_4col', 'framework_5col',
        'comparison', 'process'
    )
    if content_type in framework_types:
        populate_column_slide(slide, placeholders, slide_data)
        return
    
    # =========================================================================
    # SECONDARY ROUTING: Based on visual_type and data presence
    # =========================================================================
    
    # Timeline by visual_type
    if visual_type == 'timeline-horizontal' or extract_timeline(slide_data):
        populate_timeline_slide(slide, placeholders, slide_data)
        return
    
    # Deliverables by presence of deliverables array
    if extract_deliverables(slide_data):
        populate_deliverables_slide(slide, placeholders, slide_data)
        return
    
    # Column layouts by visual_type or layout name
    column_visual_types = (
        'process-2-phase', 'process-3-phase', 'process-4-phase', 'process-5-phase',
        'comparison-2', 'comparison-3', 'comparison-4', 'comparison-5',
        'cards-2', 'cards-3', 'cards-4', 'cards-5'
    )
    if visual_type in column_visual_types or extract_columns(slide_data):
        populate_column_slide(slide, placeholders, slide_data)
        return
    
    # Column layouts by layout name
    if 'column' in layout_name.lower():
        populate_column_slide(slide, placeholders, slide_data)
        return
    
    # Title-cover by layout name
    if layout_name == 'title-cover':
        populate_title_cover_slide(slide, placeholders, slide_data)
        return
    
    # Contact by layout name
    if layout_name == 'contact-black' or layout_name == 'contact-white':
        populate_contact_slide(slide, placeholders, slide_data)
        return
    
    # =========================================================================
    # DEFAULT: Standard content slide
    # =========================================================================
    populate_content_slide(slide, placeholders, slide_data)


# =============================================================================
# ERROR ISOLATION HELPERS
# =============================================================================

def _clear_slide_content(slide):
    """Clear user-added content from a failed slide.

    Preserves layout-inherited placeholder shapes (cleared of text) while
    removing any non-placeholder shapes that were added during population.
    """
    # Clear placeholder text frames
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()

    # Collect and remove non-placeholder shapes
    shapes_to_remove = []
    for shape in slide.shapes:
        if not shape.is_placeholder:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        remove_shape(shape)


def _add_error_annotation(slide, slide_num: int, error: Exception, slide_data: dict):
    """Add a visible error text box to a failed slide.

    The annotation shows the slide number, error type, message, and an optional
    fix suggestion. The entire function is wrapped in try/except to prevent
    cascading failures if annotation itself fails.
    """
    try:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        content_type = slide_data.get('content_type', 'unknown')
        suggestion = _suggest_fix(error, slide_data)

        # Add error text box
        txBox = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.0), Inches(10.0), Inches(3.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Title paragraph
        p = tf.paragraphs[0]
        p.text = f"Slide {slide_num} Failed ({content_type})"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
        p.alignment = PP_ALIGN.LEFT

        # Error detail paragraph
        error_msg = str(error)[:200]
        p2 = tf.add_paragraph()
        p2.text = f"{type(error).__name__}: {error_msg}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.alignment = PP_ALIGN.LEFT

        # Suggestion paragraph (only if non-empty)
        if suggestion:
            p3 = tf.add_paragraph()
            p3.text = f"Fix: {suggestion}"
            p3.font.size = Pt(11)
            p3.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
            p3.alignment = PP_ALIGN.LEFT

        # Set text box border
        txBox.line.color.rgb = RGBColor(0xCC, 0x00, 0x00)
        txBox.line.width = Pt(2)

    except Exception:
        pass  # Don't cascade annotation failures


def _extract_error_location(error: Exception, slide_data: dict) -> str:
    """Extract location context from an error and slide data.

    Returns a string describing where in the slide content the error occurred,
    based on keywords in the error message.
    """
    content_type = slide_data.get('content_type', 'unknown')
    layout_name = slide_data.get('layout', {}).get('name', 'unknown')
    visual_type = slide_data.get('visual_type', '')
    msg = str(error).lower()

    # Try to identify specific location from error message
    if 'column' in msg:
        columns = slide_data.get('columns', [])
        return f"columns section ({len(columns)} columns defined)"
    if 'table' in msg:
        tables = slide_data.get('tables', []) or slide_data.get('table_blocks', [])
        return f"table data ({len(tables)} tables defined)"
    if 'timeline' in msg:
        timeline = slide_data.get('timeline', {})
        entries = timeline.get('entries', []) if timeline else []
        return f"timeline ({len(entries)} entries)"
    if 'deliverable' in msg:
        deliverables = slide_data.get('deliverables', [])
        return f"deliverables ({len(deliverables)} items)"
    if 'placeholder' in msg or 'idx' in msg:
        return f"layout '{layout_name}' placeholder access"
    if 'image' in msg or 'picture' in msg:
        return f"image insertion in layout '{layout_name}'"
    if 'body' in msg:
        body = slide_data.get('content', {}).get('body', [])
        return f"body content ({len(body) if isinstance(body, list) else 1} items)"

    return f"{content_type} slide using layout '{layout_name}'"


def _suggest_fix(error: Exception, slide_data: dict) -> str:
    """Suggest fixes based on error type and message.

    Provides actionable guidance for common error patterns encountered
    during slide population.
    """
    msg = str(error)
    msg_lower = msg.lower()
    content_type = slide_data.get('content_type', '')
    visual_type = slide_data.get('visual_type', '')
    columns = slide_data.get('columns', [])

    # KeyError - missing field
    if isinstance(error, KeyError):
        field = msg.strip("'\"")
        if field == 'body':
            return "Add 'body' array to content, even if empty: \"body\": []"
        if field == 'header':
            return "Each column needs a 'header' field. Add \"header\": \"Title\" to column."
        if field == 'title':
            return "Add 'title' to content: \"title\": \"Slide Title\""
        return f"Add missing field '{field}' to the slide data in layout plan."

    # IndexError - array bounds
    if isinstance(error, IndexError):
        if 'column' in msg_lower or columns:
            expected = 4 if '4' in visual_type else 3 if '3' in visual_type else 2
            return f"Visual type '{visual_type}' expects {expected} columns, but only {len(columns)} provided."
        if 'body' in msg_lower:
            return "Body array has fewer items than expected. Add more bullet points."
        return "Array index out of bounds. Check that data arrays have enough items for the layout."

    # TypeError - wrong type
    if isinstance(error, TypeError):
        if 'NoneType' in msg:
            if 'subscriptable' in msg_lower:
                return "Trying to access property of null. Check that parent object exists in layout plan."
            return "Required field is null. Ensure all required fields have values."
        if 'list' in msg_lower and 'str' in msg_lower:
            return "Expected a list but got a string. Wrap single items in array: [\"item\"]"
        if 'str' in msg_lower and 'dict' in msg_lower:
            return "Expected object but got string. Use {\"key\": \"value\"} format."
        return "Data type mismatch. Check that arrays are arrays and objects are objects."

    # AttributeError - wrong structure
    if isinstance(error, AttributeError):
        if 'get' in msg_lower:
            return "Expected object/dict but got different type. Check data structure in layout plan."
        return "Unexpected data structure. Run layout plan through validator first."

    # Placeholder errors
    if 'placeholder' in msg_lower:
        layout_name = slide_data.get('layout', {}).get('name', 'unknown')
        layout_index = slide_data.get('layout', {}).get('index', -1)
        return f"Layout '{layout_name}' (index {layout_index}) missing expected placeholder. Check template has this layout."

    # Image errors
    if 'image' in msg_lower or 'picture' in msg_lower:
        return "Image insertion failed. Check file path exists and is a valid image format (PNG, JPG, etc)."

    # Font/formatting errors
    if 'font' in msg_lower or 'format' in msg_lower:
        return "Text formatting error. Check that styled_runs have valid 'text' and 'style' fields."

    # Generic fallback with content type hint
    if content_type:
        return f"Error in {content_type} slide. Verify layout plan entry matches expected schema for this type."

    return "Review the layout plan entry for this slide. Run with --dry-run to validate before generating."


# =============================================================================
# FALLBACK CONFIGURATION
# =============================================================================

DEFAULT_FALLBACK_CONFIG = {
    "fallback_layout_index": 45,
    "fallback_layout_name": "content-centered-a",
    "content_type_fallbacks": {
        "_default": "content",
    },
    "warn_on_fallback": True,
}


def _load_fallback_config(template_path: str) -> dict:
    """Load fallback config from alongside template, or use defaults.

    Looks for fallback_config.json in the same directory as the template.
    If not found or malformed, returns hardcoded defaults with a warning.
    """
    config_path = Path(template_path).parent / "fallback_config.json"
    if config_path.exists():
        try:
            loaded = json.loads(config_path.read_text(encoding="utf-8"))
            # Merge with defaults (loaded values override)
            config = {**DEFAULT_FALLBACK_CONFIG, **loaded}
            # Merge nested content_type_fallbacks
            if "content_type_fallbacks" in loaded:
                config["content_type_fallbacks"] = {
                    **DEFAULT_FALLBACK_CONFIG.get("content_type_fallbacks", {}),
                    **loaded["content_type_fallbacks"],
                }
            return config
        except (json.JSONDecodeError, OSError) as e:
            _collect_warning(f"Failed to load fallback config: {e}, using defaults")
    return dict(DEFAULT_FALLBACK_CONFIG)


def _resolve_fallbacks(slide_data: dict, config: dict, max_layout_index: int) -> tuple[dict, list[str]]:
    """Resolve fallback values for unknown content types and invalid layout indices.

    Called once per slide BEFORE layout access and populate_slide. Resolves
    exactly once (no re-resolution to avoid loops).

    Args:
        slide_data: The slide specification dict from layout plan
        config: Loaded fallback config (from _load_fallback_config)
        max_layout_index: Actual template max layout index (len(slide_layouts) - 1)

    Returns:
        Tuple of (possibly-modified slide_data copy, list of fallback warning strings)
    """
    warnings = []
    content_type = slide_data.get('content_type', '')
    layout_info = slide_data.get('layout', {})
    layout_index = layout_info.get('index', 0)

    # --- Content type fallback ---
    if content_type and content_type not in VALID_CONTENT_TYPES:
        ct_fallbacks = config.get('content_type_fallbacks', {})
        # Look up specific mapping first, then _default, then hardcoded "content"
        fallback_type = ct_fallbacks.get(content_type, ct_fallbacks.get('_default', 'content'))
        # Safety: if resolved type is also unknown, use "content" (prevent loops)
        if fallback_type not in VALID_CONTENT_TYPES:
            fallback_type = 'content'

        # Build helpful message
        valid_types = sorted(VALID_CONTENT_TYPES)
        similar = [t for t in valid_types if content_type.split('_')[0] in t or t.split('_')[0] in content_type]
        hint = f" (did you mean: {', '.join(similar[:3])}?)" if similar else ""
        warnings.append(
            f"Unknown content_type '{content_type}'{hint} -> using generic '{fallback_type}' handler"
        )
        slide_data = {**slide_data, 'content_type': fallback_type, '_original_content_type': content_type}

    # --- Layout index fallback ---
    if layout_index > max_layout_index:
        fallback_index = config.get('fallback_layout_index', 45)
        fallback_name = config.get('fallback_layout_name', 'content-centered-a')
        # Validate fallback index itself
        if fallback_index > max_layout_index:
            fallback_index = 0
            fallback_name = 'first-layout-fallback'

        original_name = layout_info.get('name', 'unknown')
        warnings.append(
            f"Layout '{original_name}' (index {layout_index}) not in template (max: {max_layout_index}) "
            f"-> using '{fallback_name}' instead. Re-profile template if layouts were added."
        )
        slide_data = {**slide_data, 'layout': {
            **layout_info,
            'index': fallback_index,
            'name': fallback_name,
            '_original_index': layout_index,
            '_original_name': original_name,
        }}

    return slide_data, warnings


def _get_placeholder_with_fallback(placeholders: list, target_type: str, slide_num: int) -> tuple:
    """Get placeholder by type with fallback chain.

    Fallback order: target_type -> BODY -> largest available text placeholder.
    Collects a warning via _collect_warning when fallback is used.

    Args:
        placeholders: List of placeholder dicts from discover_placeholders()
        target_type: The desired placeholder type (e.g., 'TITLE', 'BODY', 'SUBTITLE')
        slide_num: Slide number for warning messages

    Returns:
        Tuple of (placeholder_dict or None, bool indicating if fallback was used)
    """
    # Try exact match
    matches = [p for p in placeholders if p.get('type') == target_type]
    if matches:
        return matches[0], False

    # Fallback to BODY (most versatile)
    body_phs = [p for p in placeholders if p.get('type') == 'BODY']
    if body_phs:
        # Prefer largest body placeholder
        largest = max(body_phs, key=lambda p: p.get('width', 0) * p.get('height', 0))
        _collect_warning(
            f"Slide {slide_num}: Layout missing {target_type} placeholder - content placed in BODY area instead"
        )
        return largest, True

    # Fallback to any text placeholder available
    text_phs = [p for p in placeholders if p.get('has_text_frame', False)]
    if text_phs:
        largest = max(text_phs, key=lambda p: p.get('width', 0) * p.get('height', 0))
        available_types = [p.get('type', 'unknown') for p in text_phs]
        _collect_warning(
            f"Slide {slide_num}: No {target_type}/BODY placeholder - using {largest.get('type', 'text')} "
            f"(available: {', '.join(set(available_types))})"
        )
        return largest, True

    # No suitable placeholder
    available = [p.get('type', 'unknown') for p in placeholders]
    _collect_warning(
        f"Slide {slide_num}: Cannot place {target_type} content - no text placeholders in layout "
        f"(found: {', '.join(set(available)) if available else 'none'})"
    )
    return None, True


# =============================================================================
# PRESENTATION GENERATION
# =============================================================================

def generate_presentation(
    layout_plan: dict,
    template_path: str,
    output_path: str,
    max_errors: int | None = None,
    perf_ctx: "PerfContext | None" = None,
    split_overflow: bool = False,
) -> GenerationResult:
    """Generate presentation from layout plan.

    Returns a validated GenerationResult with per-slide tracking of status,
    warnings, and errors.

    Args:
        layout_plan: Validated layout plan dictionary
        template_path: Path to the template .pptx file
        output_path: Path for the output .pptx file
        max_errors: Optional maximum number of ERROR-level issues before
                    treating further errors as FATAL and aborting
        perf_ctx: Optional performance context for instrumentation
        split_overflow: If True, automatically split content that overflows
                       placeholder bounds onto continuation slides
    """
    global _degradation_ctx

    # Phase 6: Initialize degradation context for tracking issues
    if _HAS_GRACEFUL_DEGRADATION:
        _degradation_ctx = DegradationContext(
            component="generate_pptx",
            abort_on_fatal=True,
            max_errors=max_errors,
        )
    else:
        _degradation_ctx = None

    # Load template
    if perf_ctx and _HAS_PERFORMANCE:
        with perf_ctx.phase("load_template"):
            prs = Presentation(template_path)
    else:
        prs = Presentation(template_path)

    # Remove all existing slides
    while len(prs.slides) > 0:
        compat_delete_slide(prs, 0)

    slides_data = layout_plan.get('slides', [])
    slide_results = []
    gen_warnings = []

    # Phase 5: Clear shrink warnings for this generation run
    _shrink_warnings.clear()

    # Phase 7: Expand slides_data with content overflow splitting
    if split_overflow and _HAS_CONTENT_SPLITTER:
        expanded_slides = []
        continuation_count = 0
        for sd in slides_data:
            # Skip types that don't benefit from splitting
            content_type = sd.get('content_type', '')
            if content_type in ('branding', 'title_slide', 'contact', 'closing', 'quote', 'section_divider'):
                expanded_slides.append(sd)
                continue

            # Check for overflow and create continuation slides if needed
            split_slides = create_continuation_slides(sd)

            if len(split_slides) > 1:
                # Renumber continuation slides
                base_num = sd.get('slide_number', 0)
                for i, split_sd in enumerate(split_slides):
                    if i == 0:
                        expanded_slides.append(split_sd)
                    else:
                        # Continuation slide gets fractional slide number
                        split_sd['slide_number'] = base_num + (i * 0.1)
                        split_sd['_original_slide_number'] = base_num
                        expanded_slides.append(split_sd)
                        continuation_count += 1
            else:
                expanded_slides.append(split_slides[0])

        if continuation_count > 0:
            gen_warnings.append(
                f"Content overflow: Created {continuation_count} continuation slide(s) "
                f"to fit overflowing content"
            )
        slides_data = expanded_slides

    # Load fallback config once
    fallback_config = _load_fallback_config(template_path)
    actual_max_index = len(prs.slide_layouts) - 1

    # Validate configured fallback layout exists in this template
    cfg_fallback_idx = fallback_config.get('fallback_layout_index', 45)
    if cfg_fallback_idx > actual_max_index:
        fallback_config['fallback_layout_index'] = 0
        msg = (
            f"Configured fallback layout index {cfg_fallback_idx} exceeds template "
            f"max ({actual_max_index}), using index 0"
        )
        gen_warnings.append(msg)
        if _degradation_ctx:
            _degradation_ctx.add_warning(
                category="config",
                message=msg,
                fallback_action="Using layout index 0 as fallback",
            )

    for slide_data in slides_data:
        slide_num = slide_data.get('slide_number', 0)

        # Phase 6: Check if we should abort (too many errors)
        if _degradation_ctx and not _degradation_ctx.can_continue():
            # Mark remaining slides as skipped due to abort
            remaining_count = len(slides_data) - len(slide_results)
            gen_warnings.append(
                f"Generation aborted after {_degradation_ctx.error_count} errors. "
                f"{remaining_count} slides skipped."
            )
            break

        # Phase 3: Resolve fallbacks BEFORE layout access
        slide_data, fallback_warnings = _resolve_fallbacks(
            slide_data, fallback_config, actual_max_index
        )

        # Record fallback warnings in degradation context
        if _degradation_ctx and fallback_warnings:
            for fw in fallback_warnings:
                _degradation_ctx.add_warning(
                    category="fallback",
                    message=fw,
                    slide_number=slide_num,
                )

        # Extract values (after fallback resolution, values may have changed)
        layout_info = slide_data.get('layout', {})
        layout_index = layout_info.get('index', 0)
        layout_name = layout_info.get('name', 'unknown')
        content_type = slide_data.get('content_type', 'unknown')

        # Skip branding/master-base slides entirely (don't create empty slides)
        if content_type == 'branding' or layout_name == 'master-base':
            slide_results.append(SlideResult(
                slide_number=slide_num,
                status=SlideStatus.SKIPPED,
                layout_name=layout_name,
                layout_index=layout_index,
                content_type=content_type,
                warnings=["Branding slide skipped (no content to populate)"],
            ))
            continue

        # Get layout from template
        try:
            layout = prs.slide_layouts[layout_index]
        except IndexError:
            msg = f"Slide {slide_num}: Layout index {layout_index} ({layout_name}) not found, skipping"
            gen_warnings.append(msg)
            if _degradation_ctx:
                _degradation_ctx.add_error(
                    category="layout",
                    message=msg,
                    slide_number=slide_num,
                    location="layout.index",
                    fallback_action="Slide skipped",
                )
            slide_results.append(SlideResult(
                slide_number=slide_num,
                status=SlideStatus.SKIPPED,
                layout_name=layout_name,
                layout_index=layout_index,
                content_type=content_type,
                errors=[SlideError(
                    slide_number=slide_num,
                    error_type="LayoutNotFound",
                    message=f"Layout index {layout_index} ({layout_name}) not found in template",
                    location="layout.index",
                )],
            ))
            continue

        # Add slide with layout
        slide = prs.slides.add_slide(layout)

        # Populate content
        slide_errors = []
        status = SlideStatus.SUCCESS
        _current_slide_warnings.clear()
        _shrink_warnings.clear()  # Phase 5: Clear per-slide shrink warnings

        # Phase 4: Pre-render text overflow check
        if _HAS_VISUAL_VALIDATOR:
            content = slide_data.get('content', {})
            title_text = content.get('title', '')
            body_items = content.get('body', [])
            # Check title overflow if title exists
            if title_text:
                overflow = check_text_overflow(
                    text=str(title_text),
                    font_name=get_brand_font('header'),
                    font_size_pt=28,
                    width_inches=8.0,
                    height_inches=1.5,
                )
                if overflow.get('overflows'):
                    _collect_warning(
                        f"Title may overflow: {overflow['pct_used']:.0f}% of space used "
                        f"({overflow['lines_needed']} lines needed, {overflow['lines_available']} available)"
                    )
            # Check body overflow if body exists
            if body_items:
                body_text = '\n'.join(str(b) for b in body_items if isinstance(b, str))
                if body_text:
                    overflow = check_text_overflow(
                        text=body_text,
                        font_name=get_brand_font('body'),
                        font_size_pt=16,
                        width_inches=8.0,
                        height_inches=4.0,
                    )
                    if overflow.get('overflows'):
                        _collect_warning(
                            f"Body may overflow: {overflow['pct_used']:.0f}% of space used "
                            f"({overflow['lines_needed']} lines needed, {overflow['lines_available']} available)"
                        )

        # Time per-slide population
        slide_timer = None
        if perf_ctx and _HAS_PERFORMANCE:
            slide_timer = PerfTimer(f"slide_{slide_num}").start()

        try:
            populate_slide(slide, slide_data)
        except Exception as e:
            status = SlideStatus.FAILED
            error_location = _extract_error_location(e, slide_data)
            error_suggestion = _suggest_fix(e, slide_data)
            slide_errors.append(SlideError(
                slide_number=slide_num,
                error_type=type(e).__name__,
                message=str(e),
                location=error_location,
                suggestion=error_suggestion,
            ))
            # Record in degradation context
            if _degradation_ctx:
                _degradation_ctx.add_error(
                    category="populate",
                    message=str(e),
                    slide_number=slide_num,
                    location=error_location,
                    suggestion=error_suggestion,
                    fallback_action="Created error annotation slide",
                    exception=e,
                )
            # Clear corrupted content and add error annotation
            _clear_slide_content(slide)
            _add_error_annotation(slide, slide_num, e, slide_data)

        # Record slide timing
        if slide_timer:
            slide_timer.stop()
            perf_ctx.record_phase(
                "slide",
                slide_timer.duration_ms,
                slide_number=slide_num,
                content_type=content_type,
            )

        # Merge fallback warnings, populate warnings, and shrink warnings
        slide_warnings = list(_current_slide_warnings) + list(_shrink_warnings) + fallback_warnings

        # Phase 3: Fallback-triggered slides get PARTIAL status
        if fallback_warnings and status == SlideStatus.SUCCESS:
            status = SlideStatus.PARTIAL

        # Set slide number (protected independently)
        # Use integer slide number for display (continuation slides use fractional internally)
        try:
            display_num = int(slide_num) if isinstance(slide_num, float) else slide_num
            set_slide_number(slide, display_num)
        except Exception:
            pass  # Non-critical, don't affect slide status

        # Phase 7: Handle fractional slide numbers from continuation slides
        # SlideResult expects integer, so we use the display number (rounded)
        display_slide_num = int(slide_num) if isinstance(slide_num, float) else slide_num

        slide_results.append(SlideResult(
            slide_number=display_slide_num,
            status=status,
            layout_name=layout_name,
            layout_index=layout_index,
            content_type=content_type,
            warnings=slide_warnings,
            errors=slide_errors,
        ))

    # Save presentation
    if perf_ctx and _HAS_PERFORMANCE:
        with perf_ctx.phase("save_pptx"):
            prs.save(output_path)
    else:
        prs.save(output_path)

    # Compute counts from results
    slides_succeeded = sum(1 for r in slide_results if r.status in (SlideStatus.SUCCESS, SlideStatus.PARTIAL))
    slides_failed = sum(1 for r in slide_results if r.status == SlideStatus.FAILED)
    slides_skipped = sum(1 for r in slide_results if r.status == SlideStatus.SKIPPED)

    # Collect errors from failed/skipped slides
    all_errors = []
    for sr in slide_results:
        if sr.status in (SlideStatus.FAILED, SlideStatus.SKIPPED):
            all_errors.extend(sr.errors)

    # Phase 6: Get degradation result and reset context
    degradation_result = None
    if _degradation_ctx:
        degradation_result = _degradation_ctx.get_result()
        # Mark as partial output if we had errors but continued
        if degradation_result.error_count > 0:
            degradation_result.partial_output = True

    # Reset global context
    _degradation_ctx = None

    result = GenerationResult(
        success=(slides_failed == 0 and slides_skipped == 0),
        output_path=output_path,
        slides_total=len(slides_data),
        slides_succeeded=slides_succeeded,
        slides_failed=slides_failed,
        slides_skipped=slides_skipped,
        results=slide_results,
        warnings=gen_warnings,
        errors=all_errors,
    )

    # Attach degradation result to GenerationResult for detailed issue tracking
    # This is available via result.degradation_result if the caller wants it
    if degradation_result:
        result._degradation_result = degradation_result

    return result


def get_degradation_result(gen_result: GenerationResult) -> "DegradationResult | None":
    """Get the degradation result from a GenerationResult if available.

    Returns the DegradationResult attached during generation, or None if
    graceful degradation was not enabled.
    """
    return getattr(gen_result, '_degradation_result', None)


def validate_layout_plan_against_template(layout_plan: dict, template_path: str) -> list:
    """Validate that layout plan can be generated with the given template.

    Checks:
    - All layout indices exist in the template
    - Content types have appropriate content
    - Column counts match visual types

    Returns list of error strings. Empty list means validation passed.
    """
    errors = []

    # Load template to check layout indices
    prs = Presentation(template_path)
    max_layout_index = len(prs.slide_layouts) - 1

    for slide_data in layout_plan.get('slides', []):
        slide_num = slide_data.get('slide_number', '?')
        layout_info = slide_data.get('layout', {})
        layout_index = layout_info.get('index')
        layout_name = layout_info.get('name', 'unknown')
        content_type = slide_data.get('content_type', 'unknown')
        visual_type = slide_data.get('visual_type', '')

        # Check layout index exists
        if layout_index is not None:
            if layout_index < 0 or layout_index > max_layout_index:
                errors.append(
                    f"Slide {slide_num}: layout index {layout_index} ({layout_name}) "
                    f"out of range (template has 0-{max_layout_index})"
                )

        # Check content completeness for specific types
        content = slide_data.get('content', {})
        title = content.get('title', '') or slide_data.get('content', {}).get('headline', '')

        if content_type in ('framework', 'comparison', 'framework_2col', 'framework_3col', 'framework_4col', 'framework_5col'):
            columns = slide_data.get('columns', [])
            if not columns:
                errors.append(
                    f"Slide {slide_num} '{title}': {content_type} requires 'columns' array. "
                    f"Add columns with header/body to layout plan."
                )

            # Validate column count matches visual type
            if visual_type and columns:
                expected = None
                for n in [2, 3, 4, 5]:
                    if f'-{n}' in visual_type or f'_{n}' in content_type:
                        expected = n
                        break

                if expected and len(columns) != expected:
                    errors.append(
                        f"Slide {slide_num} '{title}': {visual_type} expects {expected} columns, "
                        f"got {len(columns)}. Add {expected - len(columns)} more column(s) or change visual_type."
                    )

        elif content_type == 'timeline':
            timeline = slide_data.get('timeline', {})
            entries = timeline.get('entries', []) if isinstance(timeline, dict) else timeline
            if not entries:
                errors.append(
                    f"Slide {slide_num} '{title}': timeline requires 'timeline.entries' array. "
                    f"Add entries with date/title/description."
                )

        elif content_type == 'deliverables':
            deliverables = slide_data.get('deliverables', []) or content.get('deliverables', [])
            body = content.get('body', [])
            if not deliverables and not (isinstance(body, list) and body and isinstance(body[0], dict)):
                errors.append(
                    f"Slide {slide_num} '{title}': deliverables requires items with title/description. "
                    f"Add 'deliverables' array or structured 'body' items."
                )

        elif content_type in ('table', 'comparison_tables'):
            tables = slide_data.get('tables', []) or slide_data.get('table_blocks', [])
            table = content.get('table') or content.get('tables')
            if not tables and not table:
                errors.append(
                    f"Slide {slide_num} '{title}': {content_type} requires table data. "
                    f"Add 'tables' array with headers/rows or 'table' object."
                )

        elif content_type in ('grid_2x2', 'grid_2x4', 'grid_3x2', 'grid_3x2_6body', 'grid_2x2_2body', 'grid_2x2_2body_b'):
            columns = slide_data.get('columns', [])
            body = content.get('body', [])
            if not columns and not body:
                errors.append(
                    f"Slide {slide_num} '{title}': {content_type} requires content. "
                    f"Add 'columns' array or 'body' items for grid cells."
                )

    return errors


def validate_image_paths(layout_plan: dict) -> list[dict]:
    """Validate that all image paths in the layout plan exist.

    Checks file_path references in:
    - columns[].file_path
    - cards[].file_path
    - deliverables[].file_path
    - extras.image_file / content.image_file
    - extras.background / content.background (when it's a file path)

    Returns list of dicts with keys: slide_number, source, file_path.
    Empty list means all images exist.
    """
    import os
    from pathlib import Path as PathLib

    missing_images = []
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.emf', '.wmf'}

    def image_exists(file_path: str) -> bool:
        """Check if an image file exists (absolute or relative to CWD)."""
        path = PathLib(file_path)
        if path.is_absolute():
            return path.exists()
        return PathLib(os.getcwd()).joinpath(file_path).exists()

    for slide in layout_plan.get('slides', []):
        slide_num = slide.get('slide_number', 0)

        # Check columns
        for col in slide.get('columns', []) or []:
            file_path = col.get('file_path')
            if file_path and not image_exists(file_path):
                missing_images.append({
                    'slide_number': slide_num,
                    'source': 'column',
                    'file_path': file_path
                })

        # Check cards
        for card in slide.get('cards', []) or []:
            file_path = card.get('file_path')
            if file_path and not image_exists(file_path):
                missing_images.append({
                    'slide_number': slide_num,
                    'source': 'card',
                    'file_path': file_path
                })

        # Check deliverables
        for deliv in slide.get('deliverables', []) or []:
            file_path = deliv.get('file_path')
            if file_path and not image_exists(file_path):
                missing_images.append({
                    'slide_number': slide_num,
                    'source': 'deliverable',
                    'file_path': file_path
                })

        # Check extras.image_file and content.image_file
        extras = slide.get('extras', {}) or {}
        content = slide.get('content', {}) or {}

        image_file = extras.get('image_file') or content.get('image_file')
        if image_file and not image_exists(image_file):
            missing_images.append({
                'slide_number': slide_num,
                'source': 'image_file',
                'file_path': image_file
            })

        # Check background (only if it looks like a file path)
        background = extras.get('background') or content.get('background')
        if background and isinstance(background, str):
            ext = PathLib(background).suffix.lower()
            if ext in image_extensions and not image_exists(background):
                missing_images.append({
                    'slide_number': slide_num,
                    'source': 'background',
                    'file_path': background
                })

    return missing_images


def _check_fonts_before_generation(layout_plan: dict) -> None:
    """Check font availability before generation and log fallback usage.

    Performs a preflight check of fonts that will be used during generation:
    - Brand fonts from template config (header/body)
    - Fonts specified in typography markers in the layout plan

    Logs warnings to stderr when fallback fonts will be used.
    """
    try:
        from font_fallback import check_font_availability
    except ImportError:
        # Font fallback module not available, skip check
        return

    # Collect fonts to check
    fonts_to_check = set()

    # Add brand fonts from template config
    config = get_template_config()
    if config:
        fonts_to_check.add(config.brand.header_font)
        fonts_to_check.add(config.brand.body_font)
    else:
        fonts_to_check.add('Aptos')  # IC default

    # Extract fonts from typography markers in layout plan
    for slide in layout_plan.get('slides', []):
        _extract_fonts_from_slide(slide, fonts_to_check)

    # Check availability
    if fonts_to_check:
        result = check_font_availability(list(fonts_to_check), log_fallbacks=False)

        # Print font status
        print("\nFont check:", file=sys.stderr)
        if result['available']:
            print(f"  Available: {', '.join(sorted(result['available']))}", file=sys.stderr)

        if result['fallbacks']:
            for font_name, info in result['fallbacks'].items():
                print(
                    f"  Fallback: '{font_name}' → '{info['resolved_to']}'",
                    file=sys.stderr
                )

        if result['unavailable']:
            print(
                f"  Warning: No fallback for: {', '.join(result['unavailable'])}",
                file=sys.stderr
            )


def _extract_fonts_from_slide(slide: dict, fonts: set) -> None:
    """Extract font names from typography markers in a slide.

    Searches for {font:FontName} markers in:
    - content.title
    - content.body items
    - columns[].header / columns[].body
    - cards[].header / cards[].body
    - tables data
    """
    import re

    font_pattern = re.compile(r'\{font:([^}]+)\}')

    def extract_from_text(text: str | None):
        if text:
            for match in font_pattern.finditer(str(text)):
                fonts.add(match.group(1))

    def extract_from_list(items: list | None):
        if items:
            for item in items:
                if isinstance(item, str):
                    extract_from_text(item)
                elif isinstance(item, dict):
                    extract_from_text(item.get('text'))

    # Content section
    content = slide.get('content', {}) or {}
    extract_from_text(content.get('title'))
    extract_from_list(content.get('body'))

    # Columns
    for col in slide.get('columns', []) or []:
        extract_from_text(col.get('header'))
        if isinstance(col.get('body'), str):
            extract_from_text(col.get('body'))
        else:
            extract_from_list(col.get('body'))

    # Cards
    for card in slide.get('cards', []) or []:
        extract_from_text(card.get('header'))
        if isinstance(card.get('body'), str):
            extract_from_text(card.get('body'))
        else:
            extract_from_list(card.get('body'))

    # Tables
    for table in slide.get('tables', []) or []:
        for row in table.get('data', []) or []:
            for cell in row:
                extract_from_text(cell)


def _format_error_report(result: GenerationResult) -> str:
    """Format a human-readable error report from GenerationResult.

    Shows generation outcome, slide counts, failed/skipped slide details
    with error types and fix suggestions, and per-slide warnings.
    """
    lines = []

    # Header
    if result.success:
        lines.append("Generation SUCCEEDED")
    else:
        lines.append("Generation COMPLETED WITH ERRORS")

    # Summary counts
    lines.append(
        f"  Total: {result.slides_total} | OK: {result.slides_succeeded} "
        f"| Failed: {result.slides_failed} | Skipped: {result.slides_skipped}"
    )

    # Failed/Skipped slide details
    problem_slides = [sr for sr in result.results if sr.status in (SlideStatus.FAILED, SlideStatus.SKIPPED)]
    if problem_slides:
        lines.append("")
        lines.append("Failed/Skipped slides:")
        for sr in problem_slides:
            lines.append(f"  Slide {sr.slide_number} [{sr.status.value}]: {sr.content_type} ({sr.layout_name})")
            for err in sr.errors:
                lines.append(f"    {err.error_type}: {err.message}")
                if err.suggestion:
                    lines.append(f"    Fix: {err.suggestion}")

    # Partial slides (fallback used)
    partial_slides = [sr for sr in result.results if sr.status == SlideStatus.PARTIAL]
    if partial_slides:
        lines.append("")
        lines.append("Partial slides (fallback used):")
        for sr in partial_slides:
            lines.append(f"  Slide {sr.slide_number} [{sr.status.value}]: {sr.content_type} ({sr.layout_name})")
            for w in sr.warnings:
                if 'fallback' in w.lower() or '->' in w:
                    lines.append(f"    Fallback: {w}")

    # Warnings section
    slides_with_warnings = [(sr.slide_number, sr.warnings) for sr in result.results if sr.warnings]
    if slides_with_warnings or result.warnings:
        lines.append("")
        lines.append("Warnings:")
        # Generation-level warnings
        for w in result.warnings:
            lines.append(f"  {w}")
        # Per-slide warnings
        for slide_num, warnings in slides_with_warnings:
            for w in warnings:
                lines.append(f"  Slide {slide_num}: {w}")

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description='Generate PowerPoint from layout plan JSON (template-agnostic with --config)'
    )
    parser.add_argument('layout_plan', help='Path to layout plan JSON file')
    parser.add_argument('--template', '-t',
        default='template/inner-chapter.pptx',
        help='Path to template .pptx file (default: template/inner-chapter.pptx)')
    parser.add_argument('--config', '-c',
        help='Template config JSON file (for template-agnostic generation)')
    parser.add_argument('--output', '-o', help='Output .pptx file path (required unless --dry-run)')
    parser.add_argument('--dry-run', action='store_true',
                        help='Validate layout plan against template without generating output')
    parser.add_argument('--json', action='store_true',
                        help='Output GenerationResult as JSON to stdout')
    parser.add_argument('--reference', '-r',
                        help='Reference PPTX for content validation (optional)')
    parser.add_argument('--validate-threshold', type=float, default=90.0,
                        help='Validation threshold for pass/fail (default: 90)')
    parser.add_argument('--validate-images', action='store_true',
                        help='Preflight check: validate all image paths exist before generation')
    parser.add_argument('--perf', action='store_true',
                        help='Enable performance instrumentation and report timing')
    parser.add_argument('--perf-json', type=str, metavar='FILE',
                        help='Write performance data to JSON file')
    parser.add_argument('--split-overflow', action='store_true',
                        help='Automatically split overflowing content onto continuation slides')

    args = parser.parse_args()

    # Initialize performance context if requested
    perf_ctx = None
    if (args.perf or args.perf_json) and _HAS_PERFORMANCE:
        perf_ctx = PerfContext("generate_pptx")
        perf_ctx.start()

    # Validate --output is provided unless dry-run
    if not args.dry_run and not args.output:
        parser.error("--output is required unless --dry-run is specified")

    # Load template config if provided
    if args.config:
        if perf_ctx:
            config_timer = PerfTimer("load_config").start()
        config_path = Path(args.config)
        if not config_path.exists():
            print(f"Error: Config file not found: {args.config}", file=sys.stderr)
            sys.exit(1)
        try:
            config_data = json.loads(config_path.read_text(encoding='utf-8'))
            if _HAS_TEMPLATE_CONFIG:
                config = TemplateConfig.model_validate(config_data)
                set_template_config(config)
                print(f"Using template config: {config.template_name}", file=sys.stderr)
            else:
                print("Warning: TemplateConfig not available, using IC defaults", file=sys.stderr)
        except Exception as e:
            print(f"Error loading config: {e}", file=sys.stderr)
            sys.exit(1)
        if perf_ctx:
            config_timer.stop()
            perf_ctx.record_phase("load_config", config_timer.duration_ms)

    # Validate inputs
    if not Path(args.layout_plan).exists():
        print(f"Error: Layout plan not found: {args.layout_plan}", file=sys.stderr)
        sys.exit(1)

    if not Path(args.template).exists():
        print(f"Error: Template not found: {args.template}", file=sys.stderr)
        sys.exit(1)

    # Load layout plan
    print(f"Generating presentation...", file=sys.stderr)
    print(f"  Layout plan: {args.layout_plan}", file=sys.stderr)
    print(f"  Template: {args.template}", file=sys.stderr)

    try:
        if perf_ctx:
            with perf_ctx.phase("load_layout_plan"):
                layout_plan = load_layout_plan(args.layout_plan)
        else:
            layout_plan = load_layout_plan(args.layout_plan)
    except LayoutPlanError as e:
        print(f"\nError: {e}", file=sys.stderr)
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"\nError: Invalid JSON in layout plan: {e}", file=sys.stderr)
        sys.exit(1)

    # Print slide summary
    print(f"\nSlide routing:", file=sys.stderr)
    for s in layout_plan.get('slides', []):
        ct = s.get('content_type', 'unknown') or 'unknown'
        vt = s.get('visual_type') or '-'
        ln = s.get('layout', {}).get('name', 'unknown')
        print(f"  {s.get('slide_number', 0):>2}: {ct:<15} | {vt:<20} | {ln}", file=sys.stderr)

    # Preflight font availability check
    if perf_ctx:
        with perf_ctx.phase("check_fonts"):
            _check_fonts_before_generation(layout_plan)
    else:
        _check_fonts_before_generation(layout_plan)

    # Preflight image validation
    if args.validate_images:
        print(f"\n{'='*60}", file=sys.stderr)
        print("IMAGE VALIDATION (preflight)", file=sys.stderr)
        print(f"{'='*60}", file=sys.stderr)

        missing_images = validate_image_paths(layout_plan)

        if missing_images:
            print(f"\nPreflight FAILED: {len(missing_images)} missing image(s):", file=sys.stderr)
            for item in missing_images:
                print(f"  Slide {item['slide_number']}: [{item['source']}] {item['file_path']}", file=sys.stderr)
            print("\nFix: Ensure all image files exist before generation.", file=sys.stderr)
            sys.exit(1)
        else:
            print("\nAll image paths validated successfully.", file=sys.stderr)

    # Dry-run validation
    if args.dry_run:
        print(f"\n{'='*60}", file=sys.stderr)
        print(f"DRY RUN - Validating layout plan against template", file=sys.stderr)
        print(f"{'='*60}", file=sys.stderr)

        validation_errors = validate_layout_plan_against_template(layout_plan, args.template)

        if validation_errors:
            print(f"\nValidation FAILED with {len(validation_errors)} error(s):", file=sys.stderr)
            for err in validation_errors:
                print(f"  - {err}", file=sys.stderr)
            sys.exit(1)
        else:
            print(f"\nValidation PASSED", file=sys.stderr)
            print(f"  Slides to create: {len(layout_plan.get('slides', []))}", file=sys.stderr)
            print(f"  All layout indices found in template", file=sys.stderr)
            sys.exit(0)

    # Generate presentation
    if perf_ctx:
        with perf_ctx.phase("generate_presentation"):
            result = generate_presentation(
                layout_plan, args.template, args.output,
                perf_ctx=perf_ctx, split_overflow=args.split_overflow,
            )
    else:
        result = generate_presentation(
            layout_plan, args.template, args.output,
            split_overflow=args.split_overflow,
        )

    # Phase 4: Validation (post-generation)
    output_dir = str(Path(args.output).parent) if args.output else '.'

    # Use comprehensive validation if reference provided, otherwise just visual
    if args.reference and Path(args.reference).exists():
        try:
            from validation_orchestrator import run_comprehensive_validation, format_validation_report

            validation_summary = run_comprehensive_validation(
                pptx_path=args.output,
                template_path=args.template,
                layout_plan=layout_plan,
                reference_path=args.reference,
                threshold=args.validate_threshold,
                output_dir=output_dir
            )

            # Update result with validation info
            if validation_summary.details.get('visual', {}).get('composite_path'):
                result = result.model_copy(update={
                    'thumbnail_path': validation_summary.details['visual']['composite_path'],
                    'visual_warnings': validation_summary.warnings,
                })
                print(f"  Thumbnails: {validation_summary.details['visual']['composite_path']}", file=sys.stderr)

            # Print validation summary
            print(f"\n  Validation Score: {validation_summary.overall_score:.1f}% "
                  f"({'PASSED' if validation_summary.passed else 'FAILED'})", file=sys.stderr)

            if validation_summary.flagged_slides:
                print(f"  Flagged slides: {len(validation_summary.flagged_slides)}", file=sys.stderr)
                for fs in validation_summary.flagged_slides[:5]:  # Show first 5
                    print(f"    Slide {fs['slide']}: {fs['issue']} [{fs['severity']}]", file=sys.stderr)
                if len(validation_summary.flagged_slides) > 5:
                    print(f"    ... and {len(validation_summary.flagged_slides) - 5} more", file=sys.stderr)

            if validation_summary.fix_suggestions:
                print(f"\n  Fix suggestions:", file=sys.stderr)
                for sug in validation_summary.fix_suggestions[:3]:  # Show first 3
                    print(f"    [{sug['priority'].upper()}] {sug['suggestion']}", file=sys.stderr)
                if len(validation_summary.fix_suggestions) > 3:
                    print(f"    ... and {len(validation_summary.fix_suggestions) - 3} more", file=sys.stderr)

            if validation_summary.warnings:
                print(f"\n  Warnings:", file=sys.stderr)
                for w in validation_summary.warnings[:3]:
                    print(f"    {w}", file=sys.stderr)

        except Exception as e:
            # Validation failure never blocks output
            print(f"  Comprehensive validation skipped: {e}", file=sys.stderr)
    else:
        # Fallback to simple visual validation
        try:
            from visual_validator import validate_visual
            visual_result = validate_visual(
                pptx_path=args.output,
                template_path=args.template,
                layout_plan=layout_plan,
                output_dir=output_dir,
            )
            if visual_result.get('composite_path'):
                result = result.model_copy(update={
                    'thumbnail_path': visual_result['composite_path'],
                    'visual_warnings': visual_result.get('warnings', []),
                })
                print(f"  Thumbnails: {visual_result['composite_path']}", file=sys.stderr)
            if visual_result.get('warnings'):
                print(f"\n  Visual warnings:", file=sys.stderr)
                for w in visual_result['warnings']:
                    print(f"    {w}", file=sys.stderr)
        except Exception as e:
            # Visual validation failure never blocks output
            print(f"  Visual validation skipped: {e}", file=sys.stderr)

    # Human-readable output
    print(f"\n  Output: {args.output}", file=sys.stderr)
    print(_format_error_report(result), file=sys.stderr)

    # Machine-readable output
    if args.json:
        print(result.model_dump_json(indent=2))

    # Performance report
    if perf_ctx:
        perf_ctx.stop()
        report = perf_ctx.get_report()

        # Print performance summary
        if args.perf:
            print(f"\n{'='*60}", file=sys.stderr)
            print(report.format_summary(), file=sys.stderr)

            # Identify bottlenecks
            bottlenecks = identify_bottlenecks(report)
            if bottlenecks:
                print(f"\nIdentified Bottlenecks:", file=sys.stderr)
                for b in bottlenecks:
                    print(f"  - {b}", file=sys.stderr)
            print(f"{'='*60}\n", file=sys.stderr)

        # Write performance JSON if requested
        if args.perf_json:
            perf_data = report.to_dict()
            Path(args.perf_json).write_text(
                json.dumps(perf_data, indent=2, ensure_ascii=False),
                encoding='utf-8'
            )
            print(f"Performance data written to: {args.perf_json}", file=sys.stderr)

    sys.exit(0 if result.success else 1)


if __name__ == '__main__':
    main()