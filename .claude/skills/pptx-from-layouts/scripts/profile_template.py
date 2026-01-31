#!/usr/bin/env python3
"""
profile_template.py: All-in-one template profiler with caching support.

Extracts layout information from a PowerPoint template and generates:
1. Full profile JSON (with placeholder positions)
2. Condensed digest JSON (for layout matching)

Usage:
    python profile_template.py template.pptx --name my-template --output-dir ./template/
    python profile_template.py template.pptx --name inner-chapter  # outputs to current dir
    python profile_template.py template.pptx --name my-template --use-cache  # use cached if valid
    python profile_template.py --clear-cache  # clear all cached profiles
"""

import argparse
import hashlib
import json
import os
import re
import sys
from pathlib import Path
from collections import defaultdict

_SCRIPTS_DIR = str(Path(__file__).resolve().parents[3] / "scripts")
if _SCRIPTS_DIR not in sys.path:
    sys.path.insert(0, _SCRIPTS_DIR)

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx_compat import get_theme_part, get_part_blob, etree


# =============================================================================
# UTILITIES
# =============================================================================

def to_kebab_case(name: str) -> str:
    """Convert name to kebab-case."""
    # Replace spaces and underscores with hyphens
    name = re.sub(r'[\s_]+', '-', name)
    # Insert hyphen before uppercase letters and lowercase them
    name = re.sub(r'([a-z])([A-Z])', r'\1-\2', name)
    # Remove any non-alphanumeric characters except hyphens
    name = re.sub(r'[^a-zA-Z0-9-]', '', name)
    # Collapse multiple hyphens
    name = re.sub(r'-+', '-', name)
    # Strip leading/trailing hyphens and lowercase
    return name.strip('-').lower()


def emu_to_inches(emu) -> float | None:
    """Convert EMUs to inches."""
    if emu is None:
        return None
    return round(emu / 914400, 3)


def get_placeholder_type_name(ph_type) -> str:
    """Convert placeholder type enum to readable string."""
    type_map = {
        PP_PLACEHOLDER.TITLE: "TITLE",
        PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
        PP_PLACEHOLDER.BODY: "BODY",
        PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
        PP_PLACEHOLDER.PICTURE: "PICTURE",
        PP_PLACEHOLDER.CHART: "CHART",
        PP_PLACEHOLDER.TABLE: "TABLE",
        PP_PLACEHOLDER.OBJECT: "OBJECT",
        PP_PLACEHOLDER.FOOTER: "FOOTER",
        PP_PLACEHOLDER.DATE: "DATE",
        PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
        PP_PLACEHOLDER.HEADER: "HEADER",
        PP_PLACEHOLDER.MEDIA_CLIP: "MEDIA_CLIP",
        PP_PLACEHOLDER.ORG_CHART: "ORG_CHART",
        PP_PLACEHOLDER.BITMAP: "BITMAP",
        PP_PLACEHOLDER.VERTICAL_BODY: "VERTICAL_BODY",
        PP_PLACEHOLDER.VERTICAL_TITLE: "VERTICAL_TITLE",
        PP_PLACEHOLDER.VERTICAL_OBJECT: "VERTICAL_OBJECT",
    }
    return type_map.get(ph_type, f"UNKNOWN({ph_type})")


def normalize_signature_key(sig_dict: dict) -> str:
    """Convert signature dict to normalized string format."""
    if not sig_dict:
        return "empty"
    sorted_items = sorted(sig_dict.items())
    return "_".join(f"{k}:{v}" for k, v in sorted_items)


def get_category(layout_name: str) -> str:
    """Determine category based on layout name."""
    name = layout_name.lower()
    
    if name == "master-base":
        return "master"
    if name.startswith("title"):
        return "title"
    if name.startswith("team-"):
        return "team"
    if name.startswith("contact"):
        return "contact"
    if name.startswith("legal"):
        return "legal"
    if name.startswith("grid-"):
        return "grid"
    if name.startswith("column-"):
        return "column"
    if name.startswith("image-"):
        return "image"
    return "content"


# =============================================================================
# EXTRACTION
# =============================================================================

def extract_color_scheme(prs) -> dict:
    """Extract theme color scheme from presentation."""
    colors = {}
    try:
        theme_part = get_theme_part(prs)

        if theme_part:
            theme_xml = etree.fromstring(get_part_blob(theme_part))
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            
            clrScheme = theme_xml.find('.//a:clrScheme', ns)
            if clrScheme is not None:
                colors["scheme_name"] = (clrScheme.get('name') or '').strip()
                for child in clrScheme:
                    tag_name = child.tag.split('}')[-1]
                    for color_elem in child:
                        color_tag = color_elem.tag.split('}')[-1]
                        if color_tag == 'srgbClr':
                            colors[tag_name] = color_elem.get('val')
                        elif color_tag == 'sysClr':
                            colors[tag_name] = color_elem.get('lastClr', color_elem.get('val'))
    except Exception as e:
        colors["_error"] = str(e)
    return colors


def extract_font_scheme(prs) -> dict:
    """Extract font scheme from presentation."""
    fonts = {}
    try:
        theme_part = get_theme_part(prs)

        if theme_part:
            theme_xml = etree.fromstring(get_part_blob(theme_part))
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            
            fontScheme = theme_xml.find('.//a:fontScheme', ns)
            if fontScheme is not None:
                fonts["scheme_name"] = fontScheme.get('name')
                
                majorFont = fontScheme.find('a:majorFont', ns)
                if majorFont is not None:
                    latin = majorFont.find('a:latin', ns)
                    if latin is not None:
                        fonts["major_font"] = latin.get('typeface')
                    ea = majorFont.find('a:ea', ns)
                    if ea is not None:
                        fonts["major_font_ea"] = ea.get('typeface')
                
                minorFont = fontScheme.find('a:minorFont', ns)
                if minorFont is not None:
                    latin = minorFont.find('a:latin', ns)
                    if latin is not None:
                        fonts["minor_font"] = latin.get('typeface')
                    ea = minorFont.find('a:ea', ns)
                    if ea is not None:
                        fonts["minor_font_ea"] = ea.get('typeface')
    except Exception as e:
        fonts["_error"] = str(e)
    return fonts


def extract_paragraph_info(paragraph) -> dict | None:
    """Extract paragraph formatting details."""
    info = {}
    try:
        if paragraph.level is not None:
            info["level"] = paragraph.level
    except Exception:
        pass
    return info if info else None


def profile_placeholder(shape) -> dict:
    """Extract detailed info about a placeholder shape."""
    ph = shape.placeholder_format
    info = {
        "idx": ph.idx,
        "type": get_placeholder_type_name(ph.type),
        "position": {
            "left_inches": emu_to_inches(shape.left),
            "top_inches": emu_to_inches(shape.top),
            "width_inches": emu_to_inches(shape.width),
            "height_inches": emu_to_inches(shape.height),
        }
    }
    
    if hasattr(shape, 'text_frame'):
        try:
            tf = shape.text_frame
            if tf.paragraphs:
                para_info = extract_paragraph_info(tf.paragraphs[0])
                if para_info:
                    info["default_paragraph"] = para_info
        except Exception:
            pass
    
    return info


def profile_layout(layout, index: int) -> dict:
    """Profile a single slide layout."""
    placeholders = []
    other_shapes = []
    
    for shape in layout.shapes:
        if shape.is_placeholder:
            placeholders.append(profile_placeholder(shape))
        else:
            shape_info = {
                "name": shape.name,
                "shape_type": f"{str(shape.shape_type).split('.')[-1]} ({shape.shape_type})" if shape.shape_type else "UNKNOWN",
                "position": {
                    "left_inches": emu_to_inches(shape.left),
                    "top_inches": emu_to_inches(shape.top),
                    "width_inches": emu_to_inches(shape.width),
                    "height_inches": emu_to_inches(shape.height),
                }
            }
            other_shapes.append(shape_info)
    
    # Build signature
    signature = defaultdict(int)
    for ph in placeholders:
        ph_type = ph["type"].lower()
        signature[ph_type] += 1
    
    return {
        "index": index,
        "name": layout.name,
        "category": get_category(layout.name),
        "placeholders": placeholders,
        "other_shapes": other_shapes,
        "signature": dict(signature)
    }


def profile_slide_master(slide_master) -> list:
    """Profile the slide master for shapes."""
    master_shapes = []
    
    for shape in slide_master.shapes:
        if not shape.is_placeholder:
            shape_info = {
                "name": shape.name,
                "shape_type": f"{str(shape.shape_type).split('.')[-1]} ({shape.shape_type})" if shape.shape_type else "UNKNOWN",
                "position": {
                    "left_inches": emu_to_inches(shape.left),
                    "top_inches": emu_to_inches(shape.top),
                    "width_inches": emu_to_inches(shape.width),
                    "height_inches": emu_to_inches(shape.height),
                }
            }
            master_shapes.append(shape_info)
    
    return master_shapes


def extract_profile(pptx_path: str) -> dict:
    """Extract full profile from PowerPoint template."""
    prs = Presentation(pptx_path)
    
    profile = {
        "template_file": str(Path(pptx_path).resolve()),
        "slide_dimensions": {
            "width_inches": emu_to_inches(prs.slide_width),
            "height_inches": emu_to_inches(prs.slide_height),
        },
        "master_styles": {
            "color_scheme": extract_color_scheme(prs),
            "font_scheme": extract_font_scheme(prs),
        },
        "master_shapes": profile_slide_master(prs.slide_master),
        "layouts": [],
        "signature_lookup": {}
    }
    
    for idx, layout in enumerate(prs.slide_layouts):
        layout_profile = profile_layout(layout, idx)
        profile["layouts"].append(layout_profile)
    
    # Build signature lookup with normalized keys
    sig_lookup = defaultdict(list)
    for layout in profile["layouts"]:
        sig_key = normalize_signature_key(layout["signature"])
        sig_lookup[sig_key].append({
            "index": layout["index"],
            "name": layout["name"]
        })
    
    profile["signature_lookup"] = dict(sig_lookup)
    
    return profile


# =============================================================================
# DIGEST CREATION
# =============================================================================

def create_digest(profile: dict) -> dict:
    """Create condensed digest from full profile."""
    digest = {
        "_meta": {
            "source": profile.get('template_file', 'unknown'),
            "description": "Condensed template digest for layout matching.",
            "slide_dimensions": profile.get('slide_dimensions')
        },
        "layouts": [],
        "signature_lookup": profile.get('signature_lookup', {}),
        "capabilities": {
            "max_body": 0,
            "max_picture": 0,
            "categories": set(),
            "body_counts_available": set()
        }
    }
    
    for layout in profile.get('layouts', []):
        condensed = {
            "index": layout.get('index'),
            "name": layout.get('name'),
            "category": layout.get('category', 'unknown'),
            "signature": layout.get('signature', {})
        }
        
        if layout.get('use_case'):
            condensed['use_case'] = layout['use_case']
        
        digest['layouts'].append(condensed)
        
        sig = layout.get('signature', {})
        body = sig.get('body', 0)
        picture = sig.get('picture', 0)
        
        digest['capabilities']['max_body'] = max(digest['capabilities']['max_body'], body)
        digest['capabilities']['max_picture'] = max(digest['capabilities']['max_picture'], picture)
        digest['capabilities']['categories'].add(condensed['category'])
        if body > 0:
            digest['capabilities']['body_counts_available'].add(body)
    
    digest['capabilities']['categories'] = sorted(digest['capabilities']['categories'])
    digest['capabilities']['body_counts_available'] = sorted(digest['capabilities']['body_counts_available'])
    
    return digest


# =============================================================================
# CACHING INFRASTRUCTURE
# =============================================================================

def get_cache_dir() -> Path:
    """Get the cache directory for template profiles."""
    # Use XDG cache dir on Linux/Mac, fallback to .cache in home
    xdg_cache = os.environ.get('XDG_CACHE_HOME')
    if xdg_cache:
        cache_base = Path(xdg_cache)
    else:
        cache_base = Path.home() / '.cache'

    cache_dir = cache_base / 'pptx-template-profiles'
    cache_dir.mkdir(parents=True, exist_ok=True)
    return cache_dir


def get_file_hash(file_path: Path) -> str:
    """Calculate MD5 hash of a file for cache validation."""
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        # Read in chunks for large files
        for chunk in iter(lambda: f.read(65536), b''):
            hasher.update(chunk)
    return hasher.hexdigest()


def get_cache_metadata_path(template_path: Path) -> Path:
    """Get the path to cache metadata file for a template."""
    cache_dir = get_cache_dir()
    # Use filename + hash of absolute path as cache key
    path_hash = hashlib.md5(str(template_path.resolve()).encode()).hexdigest()[:8]
    return cache_dir / f"{template_path.stem}_{path_hash}_cache.json"


def is_cache_valid(template_path: Path, cache_meta_path: Path) -> bool:
    """Check if cached profile is still valid for the template."""
    if not cache_meta_path.exists():
        return False

    try:
        with open(cache_meta_path, 'r', encoding='utf-8') as f:
            meta = json.load(f)

        # Check modification time
        current_mtime = template_path.stat().st_mtime
        if meta.get('mtime') != current_mtime:
            return False

        # Check file hash for extra safety
        current_hash = get_file_hash(template_path)
        if meta.get('hash') != current_hash:
            return False

        # Check that cached files exist
        if not Path(meta.get('profile_path', '')).exists():
            return False
        if meta.get('digest_path') and not Path(meta['digest_path']).exists():
            return False

        return True
    except (json.JSONDecodeError, KeyError, OSError):
        return False


def save_cache_metadata(template_path: Path, profile_path: Path,
                        digest_path: Path | None = None):
    """Save cache metadata for a template profile."""
    cache_meta_path = get_cache_metadata_path(template_path)

    meta = {
        'template_path': str(template_path.resolve()),
        'mtime': template_path.stat().st_mtime,
        'hash': get_file_hash(template_path),
        'profile_path': str(profile_path.resolve()),
        'digest_path': str(digest_path.resolve()) if digest_path else None,
        'cached_at': str(Path(cache_meta_path).stat().st_mtime if cache_meta_path.exists() else 0)
    }

    with open(cache_meta_path, 'w', encoding='utf-8') as f:
        json.dump(meta, f, indent=2)


def get_cached_paths(template_path: Path) -> tuple[Path | None, Path | None]:
    """Get cached profile and digest paths if cache is valid.

    Returns (profile_path, digest_path) or (None, None) if cache invalid.
    """
    cache_meta_path = get_cache_metadata_path(template_path)

    if not is_cache_valid(template_path, cache_meta_path):
        return None, None

    with open(cache_meta_path, 'r', encoding='utf-8') as f:
        meta = json.load(f)

    profile_path = Path(meta['profile_path']) if meta.get('profile_path') else None
    digest_path = Path(meta['digest_path']) if meta.get('digest_path') else None

    return profile_path, digest_path


def clear_all_cache():
    """Clear all cached template profiles."""
    cache_dir = get_cache_dir()
    count = 0
    for cache_file in cache_dir.glob('*_cache.json'):
        try:
            # Also try to read and delete the referenced files
            with open(cache_file, 'r', encoding='utf-8') as f:
                meta = json.load(f)
            for key in ['profile_path', 'digest_path']:
                if meta.get(key) and Path(meta[key]).exists():
                    # Only delete if in cache directory
                    if str(cache_dir) in meta[key]:
                        Path(meta[key]).unlink()
            cache_file.unlink()
            count += 1
        except (json.JSONDecodeError, OSError):
            cache_file.unlink()
            count += 1
    return count


# =============================================================================
# MAIN
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description='Profile a PowerPoint template and generate profile + digest JSON files'
    )
    parser.add_argument('template', nargs='?', help='Path to .pptx template file')
    parser.add_argument('--name', '-n',
                        help='Template name (kebab-case, e.g., "inner-chapter")')
    parser.add_argument('--output-dir', '-o', default='.',
                        help='Output directory (default: current directory)')
    parser.add_argument('--profile-only', action='store_true',
                        help='Only generate full profile, skip digest')
    parser.add_argument('--digest-only', action='store_true',
                        help='Only generate digest (requires existing profile)')
    parser.add_argument('--use-cache', action='store_true',
                        help='Use cached profile if valid (based on file modification time)')
    parser.add_argument('--clear-cache', action='store_true',
                        help='Clear all cached profiles and exit')
    parser.add_argument('--generate-config', action='store_true',
                        help='Also generate template config JSON for template-agnostic generation')

    args = parser.parse_args()

    # Handle clear-cache mode first (no other args needed)
    if args.clear_cache:
        count = clear_all_cache()
        print(f"Cleared {count} cached profile(s)", file=sys.stderr)
        return

    # Validate required arguments
    if not args.template:
        parser.error("template is required unless using --clear-cache")
    if not args.name:
        parser.error("--name is required unless using --clear-cache")

    # Validate and normalize name
    name = to_kebab_case(args.name)
    if name != args.name:
        print(f"Note: Normalized name to '{name}'", file=sys.stderr)

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    profile_path = output_dir / f"{name}-profile.json"
    digest_path = output_dir / f"{name}-digest.json"
    template_path = Path(args.template)

    # Handle cache mode
    if args.use_cache and template_path.exists():
        cached_profile, cached_digest = get_cached_paths(template_path)
        if cached_profile and cached_profile.exists():
            print(f"Using cached profile: {cached_profile}", file=sys.stderr)
            # Copy cached files to output location if different
            if cached_profile.resolve() != profile_path.resolve():
                import shutil
                shutil.copy2(cached_profile, profile_path)
                if cached_digest and cached_digest.exists():
                    shutil.copy2(cached_digest, digest_path)
            print(f"Cache hit - skipped expensive extraction", file=sys.stderr)
            return
        else:
            print(f"Cache miss - regenerating profile", file=sys.stderr)

    # Handle digest-only mode
    if args.digest_only:
        if not profile_path.exists():
            print(f"Error: Profile not found: {profile_path}", file=sys.stderr)
            print("Run without --digest-only to generate profile first.", file=sys.stderr)
            sys.exit(1)
        
        profile = json.loads(profile_path.read_text(encoding='utf-8'))
        digest = create_digest(profile)
        
        digest_json = json.dumps(digest, indent=2, ensure_ascii=False)
        digest_path.write_text(digest_json, encoding='utf-8')
        
        print(f"Created: {digest_path}", file=sys.stderr)
        print(f"  Size: {len(digest_json):,} bytes", file=sys.stderr)
        return
    
    # Check template exists
    if not template_path.exists():
        print(f"Error: Template not found: {args.template}", file=sys.stderr)
        sys.exit(1)

    # Extract profile
    print(f"Extracting profile from: {template_path}", file=sys.stderr)
    profile = extract_profile(str(template_path))
    
    # Write profile
    profile_json = json.dumps(profile, indent=2, ensure_ascii=False)
    profile_path.write_text(profile_json, encoding='utf-8')
    
    print(f"\nCreated: {profile_path}", file=sys.stderr)
    print(f"  Layouts: {len(profile['layouts'])}", file=sys.stderr)
    print(f"  Size: {len(profile_json):,} bytes", file=sys.stderr)
    
    # Check for signature collisions
    collisions = {k: v for k, v in profile['signature_lookup'].items() if len(v) > 1}
    if collisions:
        print(f"\n  Signature collisions detected: {len(collisions)}", file=sys.stderr)
        for sig, layouts in collisions.items():
            names = [l['name'] for l in layouts]
            print(f"    {sig}: {', '.join(names)}", file=sys.stderr)
        print("  Consider adding use_case descriptions to differentiate these.", file=sys.stderr)
    
    # Create digest unless profile-only
    final_digest_path = None
    if not args.profile_only:
        digest = create_digest(profile)
        digest_json = json.dumps(digest, indent=2, ensure_ascii=False)
        digest_path.write_text(digest_json, encoding='utf-8')
        final_digest_path = digest_path

        reduction = (1 - len(digest_json) / len(profile_json)) * 100

        print(f"\nCreated: {digest_path}", file=sys.stderr)
        print(f"  Size: {len(digest_json):,} bytes ({reduction:.0f}% smaller)", file=sys.stderr)
        print(f"  Max body: {digest['capabilities']['max_body']}", file=sys.stderr)
        print(f"  Categories: {', '.join(digest['capabilities']['categories'])}", file=sys.stderr)

    # Save cache metadata for future runs
    save_cache_metadata(template_path, profile_path, final_digest_path)
    print(f"\nCache metadata saved for future --use-cache runs", file=sys.stderr)

    # Generate template config if requested
    if args.generate_config:
        from generate_config import generate_config
        config = generate_config(profile_data=profile, template_name=name)
        config_path = output_dir / f"{name}-config.json"
        config_json = json.dumps(config, indent=2, ensure_ascii=False)
        config_path.write_text(config_json, encoding='utf-8')
        print(f"\nCreated: {config_path}", file=sys.stderr)
        print(f"  Layout mappings: {len(config['layout_mappings'])}", file=sys.stderr)
        print(f"  Content type routes: {len(config['content_type_routing'])}", file=sys.stderr)

    # Print summary
    print(f"\n{'='*60}", file=sys.stderr)
    print("TEMPLATE PROFILE SUMMARY", file=sys.stderr)
    print(f"{'='*60}", file=sys.stderr)
    print(f"Template: {name}", file=sys.stderr)
    print(f"Dimensions: {profile['slide_dimensions']['width_inches']}\" x {profile['slide_dimensions']['height_inches']}\"", file=sys.stderr)
    print(f"Layouts: {len(profile['layouts'])}", file=sys.stderr)
    
    # Category breakdown
    categories = defaultdict(int)
    for layout in profile['layouts']:
        categories[layout['category']] += 1
    
    print(f"\nLayouts by category:", file=sys.stderr)
    for cat, count in sorted(categories.items()):
        print(f"  {cat}: {count}", file=sys.stderr)


if __name__ == '__main__':
    main()
