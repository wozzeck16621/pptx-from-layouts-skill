#!/usr/bin/env python3
"""
PPTX Quality Checker - Standalone validation without reference file.

Checks for common issues that make presentations look unprofessional:
- Empty or near-empty slides
- Missing titles
- Placeholder text left in ("[Image: ...]", "Click to add...")
- Text overflow (content exceeding shape bounds)
- Unused layouts (template coverage)
- Consistency issues (font sizes, bullet styles)
- Accessibility issues (low contrast, small text)

Tier 2 validators (optional, based on inputs):
- Typography hierarchy (PPTX): Font size progression validation
- Visual type appropriateness (layout plan): Content-visual type matching
- Column balance (layout plan): Multi-column content distribution

Usage:
    python quality_check.py presentation.pptx
    python quality_check.py presentation.pptx --template template.pptx
    python quality_check.py presentation.pptx --layout-plan layout.json
    python quality_check.py presentation.pptx --json --output report.json
"""

import argparse
import json
import re
import sys
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple
from enum import Enum
import multiprocessing

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# PIL for accurate text measurement
try:
    from PIL import ImageFont
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# Import Tier 2 validators
from typography_hierarchy import TypographyHierarchyValidator
from visual_type_checker import VisualTypeChecker
from column_balance import ColumnBalanceAnalyzer
from whitespace_analyzer import WhitespaceAnalyzer
from layout_coverage_analyzer import LayoutCoverageAnalyzer


# ============================================================================
# Parallel Validator Execution Functions (module-level for multiprocessing)
# ============================================================================

def _run_typography_validator(pptx_path: str) -> Dict[str, Any]:
    """Run typography hierarchy validator in a subprocess.

    Args:
        pptx_path: Path to the PPTX file to validate.

    Returns:
        Dictionary with validator results including score, passed, issues, etc.
    """
    try:
        validator = TypographyHierarchyValidator(pptx_path)
        report = validator.check()
        return {
            'validator': 'typography_hierarchy',
            'score': report.score,
            'passed': report.passed,
            'summary': report.summary,
            'global_stats': report.global_stats,
            'issues': [
                {
                    'slide_number': i.slide_number,
                    'category': i.category,
                    'severity': i.severity.value,
                    'message': i.message,
                    'suggestion': i.suggestion,
                    'details': i.details,
                }
                for i in report.issues
            ],
            'error': None,
        }
    except Exception as e:
        return {
            'validator': 'typography_hierarchy',
            'error': str(e),
            'score': None,
            'passed': None,
            'issues': [],
        }


def _run_whitespace_validator(pptx_path: str) -> Dict[str, Any]:
    """Run whitespace analyzer in a subprocess.

    Args:
        pptx_path: Path to the PPTX file to validate.

    Returns:
        Dictionary with validator results including score, passed, issues, etc.
    """
    try:
        analyzer = WhitespaceAnalyzer(pptx_path)
        report = analyzer.check()
        return {
            'validator': 'whitespace',
            'score': report.score,
            'passed': report.passed,
            'summary': report.summary,
            'global_stats': report.global_stats,
            'issues': [
                {
                    'slide_number': i.slide_number,
                    'category': i.category,
                    'severity': i.severity.value,
                    'message': i.message,
                    'suggestion': i.suggestion,
                    'details': i.details,
                }
                for i in report.issues
            ],
            'error': None,
        }
    except Exception as e:
        return {
            'validator': 'whitespace',
            'error': str(e),
            'score': None,
            'passed': None,
            'issues': [],
        }


def _run_visual_type_validator(layout_plan_path: str) -> Dict[str, Any]:
    """Run visual type checker in a subprocess.

    Args:
        layout_plan_path: Path to the layout plan JSON file.

    Returns:
        Dictionary with validator results including score, passed, issues, etc.
    """
    try:
        checker = VisualTypeChecker(layout_plan_path)
        report = checker.check()
        return {
            'validator': 'visual_type',
            'score': report.score,
            'passed': report.passed,
            'summary': report.summary,
            'visual_type_distribution': report.visual_type_distribution,
            'issues': [
                {
                    'slide_number': i.slide_number,
                    'visual_type': i.visual_type,
                    'content_type': i.content_type,
                    'severity': i.severity.value,
                    'message': i.message,
                    'suggestion': i.suggestion,
                    'details': i.details,
                }
                for i in report.issues
            ],
            'error': None,
        }
    except Exception as e:
        return {
            'validator': 'visual_type',
            'error': str(e),
            'score': None,
            'passed': None,
            'issues': [],
        }


def _run_column_balance_validator(layout_plan_path: str) -> Dict[str, Any]:
    """Run column balance analyzer in a subprocess.

    Args:
        layout_plan_path: Path to the layout plan JSON file.

    Returns:
        Dictionary with validator results including score, passed, issues, etc.
    """
    try:
        analyzer = ColumnBalanceAnalyzer(layout_plan_path)
        report = analyzer.check()
        return {
            'validator': 'column_balance',
            'score': report.score,
            'passed': report.passed,
            'summary': report.summary,
            'slides_with_columns': report.slides_with_columns,
            'issues': [
                {
                    'slide_number': i.slide_number,
                    'column_numbers': i.column_numbers,
                    'visual_type': i.visual_type,
                    'severity': i.severity.value,
                    'message': i.message,
                    'suggestion': i.suggestion,
                    'details': i.details,
                }
                for i in report.issues
            ],
            'error': None,
        }
    except Exception as e:
        return {
            'validator': 'column_balance',
            'error': str(e),
            'score': None,
            'passed': None,
            'issues': [],
        }


class Severity(Enum):
    """Issue severity levels."""
    ERROR = "error"      # Must fix before delivery
    WARNING = "warning"  # Should fix, may be intentional
    INFO = "info"        # Nice to know, minor improvement


class ExitCode(Enum):
    """Exit codes based on worst issue severity encountered.

    Exit codes are ordered by severity, allowing consumers to check:
    - exit code == 0: No issues (clean pass)
    - exit code == 1: INFO-level issues only (minor)
    - exit code >= 2: WARNING or ERROR level issues (should review)
    - exit code == 3: ERROR-level issues (must fix)
    """
    CLEAN = 0       # No issues found
    INFO = 1        # Only INFO-level issues
    WARNING = 2     # WARNING-level issues (worst)
    ERROR = 3       # ERROR-level issues (worst)


@dataclass
class Issue:
    """A quality issue found in the presentation."""
    slide_number: int
    category: str
    severity: Severity
    message: str
    suggestion: str
    details: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Tier2Results:
    """Results from Tier 2 validators."""
    typography_hierarchy: Optional[Dict[str, Any]] = None
    visual_type: Optional[Dict[str, Any]] = None
    column_balance: Optional[Dict[str, Any]] = None
    whitespace: Optional[Dict[str, Any]] = None


@dataclass
class HeuristicScoreBreakdown:
    """Breakdown of heuristic scores by component.

    Each component contributes to the composite score based on its weight.
    Weights sum to 100% when all components are available.
    """
    # Core Tier 1 checks (content quality)
    core_quality: Optional[float] = None       # Weight: 40%
    # Tier 2 components
    typography_hierarchy: Optional[float] = None   # Weight: 15%
    whitespace_analysis: Optional[float] = None    # Weight: 15%
    visual_type: Optional[float] = None           # Weight: 15%
    column_balance: Optional[float] = None        # Weight: 15%

    # Weights for each component (must sum to 100)
    WEIGHTS = {
        'core_quality': 40,
        'typography_hierarchy': 15,
        'whitespace_analysis': 15,
        'visual_type': 15,
        'column_balance': 15,
    }

    def compute_composite_score(self) -> float:
        """Compute weighted composite score from available components.

        Returns a score from 0-100. Components that are not available
        (None) are excluded and their weights are redistributed to
        available components.
        """
        available_scores = []
        available_weights = []

        # Collect available components
        components = [
            ('core_quality', self.core_quality),
            ('typography_hierarchy', self.typography_hierarchy),
            ('whitespace_analysis', self.whitespace_analysis),
            ('visual_type', self.visual_type),
            ('column_balance', self.column_balance),
        ]

        for name, score in components:
            if score is not None:
                available_scores.append(score)
                available_weights.append(self.WEIGHTS[name])

        if not available_scores:
            return 0.0

        # Normalize weights to sum to 100
        total_weight = sum(available_weights)
        if total_weight == 0:
            return 0.0

        normalized_weights = [w / total_weight * 100 for w in available_weights]

        # Compute weighted average
        composite = sum(s * w / 100 for s, w in zip(available_scores, normalized_weights))
        return max(0.0, min(100.0, composite))

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        result = {
            'composite_score': round(self.compute_composite_score(), 1),
            'components': {},
            'weights': self.WEIGHTS.copy(),
        }

        if self.core_quality is not None:
            result['components']['core_quality'] = round(self.core_quality, 1)
        if self.typography_hierarchy is not None:
            result['components']['typography_hierarchy'] = round(self.typography_hierarchy, 1)
        if self.whitespace_analysis is not None:
            result['components']['whitespace_analysis'] = round(self.whitespace_analysis, 1)
        if self.visual_type is not None:
            result['components']['visual_type'] = round(self.visual_type, 1)
        if self.column_balance is not None:
            result['components']['column_balance'] = round(self.column_balance, 1)

        return result


@dataclass
class QualityReport:
    """Complete quality report for a presentation."""
    file_path: str
    slide_count: int
    issues: List[Issue]
    layout_coverage: Dict[str, Any]
    summary: Dict[str, int]
    score: float
    passed: bool
    recommendations: List[str]
    tier2_results: Optional[Tier2Results] = None
    heuristic_breakdown: Optional[HeuristicScoreBreakdown] = None

    def to_dict(self) -> Dict[str, Any]:
        result = {
            'file_path': self.file_path,
            'slide_count': self.slide_count,
            'score': self.score,
            'passed': self.passed,
            'summary': self.summary,
            'layout_coverage': self.layout_coverage,
            'recommendations': self.recommendations,
            'issues': [
                {
                    'slide_number': i.slide_number,
                    'category': i.category,
                    'severity': i.severity.value,
                    'message': i.message,
                    'suggestion': i.suggestion,
                    'details': i.details
                }
                for i in self.issues
            ]
        }
        if self.tier2_results:
            result['tier2'] = {}
            if self.tier2_results.typography_hierarchy:
                result['tier2']['typography_hierarchy'] = self.tier2_results.typography_hierarchy
            if self.tier2_results.visual_type:
                result['tier2']['visual_type'] = self.tier2_results.visual_type
            if self.tier2_results.column_balance:
                result['tier2']['column_balance'] = self.tier2_results.column_balance
            if self.tier2_results.whitespace:
                result['tier2']['whitespace'] = self.tier2_results.whitespace
        if self.heuristic_breakdown:
            result['heuristic_scoring'] = self.heuristic_breakdown.to_dict()
        return result


class TextMeasurer:
    """Measures text dimensions using PIL for accurate overflow detection.

    Uses PIL's ImageFont to measure actual text width based on font metrics.
    Falls back to character-count heuristic when PIL is unavailable.
    """

    # Font search paths by platform
    FONT_SEARCH_PATHS = {
        'darwin': [
            '/System/Library/Fonts',
            '/Library/Fonts',
            '~/Library/Fonts',
        ],
        'linux': [
            '/usr/share/fonts/truetype',
            '/usr/local/share/fonts',
            '~/.local/share/fonts',
            '~/.fonts',
        ],
        'win32': [
            'C:/Windows/Fonts',
        ],
    }

    # Font file mappings for common fonts
    FONT_FILES = {
        # Aptos family (Microsoft/IC brand)
        'aptos': 'Aptos.ttf',
        'aptos display': 'Aptos-Display.ttf',
        'aptos narrow': 'Aptos-Narrow.ttf',
        # Calibri family
        'calibri': 'Calibri.ttf',
        'calibri light': 'calibril.ttf',
        'calibri bold': 'calibrib.ttf',
        # Arial family
        'arial': 'Arial.ttf',
        'arial bold': 'Arial Bold.ttf',
        'arial narrow': 'Arial Narrow.ttf',
        # Fallbacks
        'helvetica': 'Helvetica.ttc',
        'verdana': 'Verdana.ttf',
    }

    # Fallback chain for when requested font is not found
    FALLBACK_CHAIN = ['arial', 'helvetica', 'verdana', 'calibri']

    def __init__(self):
        """Initialize the text measurer."""
        self._font_cache: Dict[tuple, Any] = {}
        self._platform = sys.platform
        self._search_paths = self._get_search_paths()

    def _get_search_paths(self) -> List[Path]:
        """Get font search paths for current platform."""
        paths = []
        platform_paths = self.FONT_SEARCH_PATHS.get(self._platform, [])

        for p in platform_paths:
            expanded = Path(p).expanduser()
            if expanded.exists():
                paths.append(expanded)

        return paths

    def _find_font_file(self, font_name: str) -> Optional[Path]:
        """Find the font file for a given font name."""
        if not font_name:
            font_name = 'arial'

        # Normalize font name
        normalized = font_name.lower().strip()

        # Try exact match first
        if normalized in self.FONT_FILES:
            filename = self.FONT_FILES[normalized]
            for search_path in self._search_paths:
                # Try direct path
                font_path = search_path / filename
                if font_path.exists():
                    return font_path
                # Try case-insensitive search in directory
                try:
                    for f in search_path.iterdir():
                        if f.name.lower() == filename.lower():
                            return f
                except (PermissionError, OSError):
                    continue

        # Try partial match
        for search_path in self._search_paths:
            try:
                for f in search_path.rglob('*.tt[fc]'):
                    if normalized in f.stem.lower():
                        return f
            except (PermissionError, OSError):
                continue

        # Try fallback chain
        for fallback in self.FALLBACK_CHAIN:
            if fallback != normalized and fallback in self.FONT_FILES:
                filename = self.FONT_FILES[fallback]
                for search_path in self._search_paths:
                    font_path = search_path / filename
                    if font_path.exists():
                        return font_path

        return None

    def _get_font(self, font_name: str, size_pt: float) -> Optional[Any]:
        """Get a PIL ImageFont for the specified font and size.

        Args:
            font_name: Name of the font (e.g., 'Arial', 'Aptos')
            size_pt: Font size in points

        Returns:
            PIL ImageFont object or None if PIL unavailable
        """
        if not HAS_PIL:
            return None

        cache_key = (font_name, size_pt)
        if cache_key in self._font_cache:
            return self._font_cache[cache_key]

        font_path = self._find_font_file(font_name)

        try:
            if font_path:
                font = ImageFont.truetype(str(font_path), int(size_pt))
            else:
                # Use PIL's default font scaled (less accurate but functional)
                font = ImageFont.load_default()
            self._font_cache[cache_key] = font
            return font
        except (OSError, IOError):
            # Font loading failed, cache None to avoid repeated attempts
            self._font_cache[cache_key] = None
            return None

    def measure_text_width(self, text: str, font_name: str, size_pt: float) -> float:
        """Measure the width of text in points, accounting for kerning.

        Uses PIL's getlength() method which includes kerning pairs for accurate
        text width measurement. Falls back to character-count heuristic when
        PIL is unavailable.

        Args:
            text: The text to measure
            font_name: Name of the font
            size_pt: Font size in points

        Returns:
            Width of the text in points (kerning-aware when PIL available)
        """
        if not text:
            return 0.0

        if not HAS_PIL:
            # Fallback: estimate based on character count and font size
            # Average character width is roughly 0.5-0.6 of font size
            return len(text) * size_pt * 0.55

        font = self._get_font(font_name, size_pt)
        if not font:
            return len(text) * size_pt * 0.55

        try:
            # Use getlength() for kerning-aware measurement (Pillow 8.0+)
            # This accounts for kerning pairs like "AV", "To", "We" etc.
            if hasattr(font, 'getlength'):
                return font.getlength(text)
            # Fallback to getbbox for older Pillow versions
            bbox = font.getbbox(text)
            if bbox:
                return bbox[2] - bbox[0]
            return len(text) * size_pt * 0.55
        except Exception:
            return len(text) * size_pt * 0.55

    def measure_text_height(self, text: str, font_name: str, size_pt: float,
                           line_spacing: float = 1.2) -> float:
        """Measure the height of text in points, accounting for line wrapping.

        Args:
            text: The text to measure
            font_name: Name of the font
            size_pt: Font size in points
            line_spacing: Line height multiplier (default 1.2)

        Returns:
            Height of the text in points
        """
        if not text:
            return 0

        # Count actual newlines
        lines = text.split('\n')
        num_lines = len(lines)

        # Line height = font size * line spacing
        line_height = size_pt * line_spacing

        return num_lines * line_height

    def _wrap_paragraph(self, paragraph: str, font_name: str, size_pt: float,
                        available_width_pt: float) -> List[str]:
        """Word-wrap a paragraph using font metrics.

        Performs proper word-by-word wrapping using measured text widths,
        accounting for kerning. Handles long words that exceed line width
        by breaking them with hyphens.

        Args:
            paragraph: Single paragraph text (no newlines)
            font_name: Name of the font
            size_pt: Font size in points
            available_width_pt: Available width in points

        Returns:
            List of wrapped lines
        """
        if not paragraph.strip():
            return ['']

        words = paragraph.split()
        if not words:
            return ['']

        lines: List[str] = []
        current_line = ''

        # Measure hyphen width for word breaking
        hyphen_width = self.measure_text_width('-', font_name, size_pt)

        for word in words:
            if not current_line:
                # Starting a new line
                word_width = self.measure_text_width(word, font_name, size_pt)

                if word_width <= available_width_pt:
                    # Word fits on its own line
                    current_line = word
                else:
                    # Word is too long - break it with hyphen
                    broken_lines = self._break_long_word(
                        word, font_name, size_pt, available_width_pt, hyphen_width
                    )
                    if len(broken_lines) > 1:
                        # Add all complete lines except the last
                        lines.extend(broken_lines[:-1])
                        # Last fragment becomes current line
                        current_line = broken_lines[-1]
                    else:
                        current_line = broken_lines[0]
            else:
                # Try to add word to current line
                test_line = current_line + ' ' + word
                test_width = self.measure_text_width(test_line, font_name, size_pt)

                if test_width <= available_width_pt:
                    # Word fits on current line
                    current_line = test_line
                else:
                    # Word doesn't fit - start new line
                    lines.append(current_line)

                    word_width = self.measure_text_width(word, font_name, size_pt)
                    if word_width <= available_width_pt:
                        current_line = word
                    else:
                        # Word is too long - break it
                        broken_lines = self._break_long_word(
                            word, font_name, size_pt, available_width_pt, hyphen_width
                        )
                        if len(broken_lines) > 1:
                            lines.extend(broken_lines[:-1])
                            current_line = broken_lines[-1]
                        else:
                            current_line = broken_lines[0]

        # Don't forget the last line
        if current_line:
            lines.append(current_line)

        return lines if lines else ['']

    def _break_long_word(self, word: str, font_name: str, size_pt: float,
                         available_width_pt: float, hyphen_width: float) -> List[str]:
        """Break a long word into multiple lines with hyphens.

        Args:
            word: The word to break
            font_name: Name of the font
            size_pt: Font size in points
            available_width_pt: Available width in points
            hyphen_width: Pre-measured width of hyphen character

        Returns:
            List of word fragments (with hyphens except for last fragment)
        """
        if not word:
            return ['']

        fragments: List[str] = []
        remaining = word
        max_width_with_hyphen = available_width_pt - hyphen_width

        while remaining:
            if self.measure_text_width(remaining, font_name, size_pt) <= available_width_pt:
                # Remaining text fits
                fragments.append(remaining)
                break

            # Find break point - binary search for longest prefix that fits
            # with a hyphen
            low, high = 1, len(remaining)
            best_break = 1  # At minimum, take one character

            while low <= high:
                mid = (low + high) // 2
                prefix = remaining[:mid]
                prefix_width = self.measure_text_width(prefix, font_name, size_pt)

                if prefix_width <= max_width_with_hyphen:
                    best_break = mid
                    low = mid + 1
                else:
                    high = mid - 1

            # Ensure we make progress (at least 1 char per line)
            if best_break == 0:
                best_break = 1

            fragment = remaining[:best_break] + '-'
            fragments.append(fragment)
            remaining = remaining[best_break:]

        return fragments if fragments else [word]

    def estimate_lines_needed(self, text: str, font_name: str, size_pt: float,
                             available_width_pt: float) -> int:
        """Estimate how many lines are needed to display text.

        Uses proper word-by-word wrapping with font metrics (kerning-aware).
        Handles explicit newlines as paragraph breaks and breaks long words
        with hyphens when they exceed available width.

        Args:
            text: The text to fit
            font_name: Name of the font
            size_pt: Font size in points
            available_width_pt: Available width in points

        Returns:
            Estimated number of lines needed
        """
        if not text or available_width_pt <= 0:
            return 0

        total_lines = 0

        # Process each paragraph (split by newlines)
        for paragraph in text.split('\n'):
            # Wrap paragraph using font metrics
            wrapped_lines = self._wrap_paragraph(
                paragraph, font_name, size_pt, available_width_pt
            )
            total_lines += len(wrapped_lines)

        return max(1, total_lines)

    def text_fits_in_shape(self, text: str, font_name: str, size_pt: float,
                          shape_width_pt: float, shape_height_pt: float,
                          margin_pt: float = 10, line_spacing: float = 1.2) -> tuple:
        """Check if text fits within shape bounds.

        Args:
            text: The text to fit
            font_name: Name of the font
            size_pt: Font size in points
            shape_width_pt: Shape width in points
            shape_height_pt: Shape height in points
            margin_pt: Internal margin in points
            line_spacing: Line height multiplier

        Returns:
            Tuple of (fits: bool, overflow_ratio: float, details: dict)
            overflow_ratio > 1.0 means text overflows
            overflow_ratio of -1.0 indicates unmeasurable shape (skip check)
        """
        if not text:
            return (True, 0.0, {'lines_needed': 0, 'lines_available': 0})

        available_width = shape_width_pt - (2 * margin_pt)
        available_height = shape_height_pt - (2 * margin_pt)

        # Handle zero or negative dimensions - return special marker to skip
        # These occur with table headers, dividers, or malformed shapes
        if available_width <= 0 or available_height <= 0:
            return (True, -1.0, {
                'error': 'shape too small to measure',
                'skip_reason': 'zero_or_negative_dimensions',
                'shape_width_pt': shape_width_pt,
                'shape_height_pt': shape_height_pt,
            })

        # Estimate lines needed
        lines_needed = self.estimate_lines_needed(
            text, font_name, size_pt, available_width
        )

        # Calculate lines available based on line height
        line_height = size_pt * line_spacing

        # If shape can't fit even one line at this font size, it's unmeasurable
        # This avoids false positives from shapes sized for different fonts
        min_height_for_one_line = line_height * 0.8  # Allow 20% tolerance
        if available_height < min_height_for_one_line:
            return (True, -1.0, {
                'error': 'shape height insufficient for font size',
                'skip_reason': 'shape_too_short_for_font',
                'available_height_pt': available_height,
                'min_line_height_pt': line_height,
            })

        lines_available = max(1, int(available_height / line_height))

        # Calculate overflow ratio
        overflow_ratio = lines_needed / lines_available

        fits = lines_needed <= lines_available

        return (fits, overflow_ratio, {
            'lines_needed': lines_needed,
            'lines_available': lines_available,
            'text_length': len(text),
            'available_width_pt': available_width,
            'available_height_pt': available_height,
        })


class QualityChecker:
    """Checks presentation quality without needing a reference file."""

    # Placeholder patterns that shouldn't appear in final presentations
    PLACEHOLDER_PATTERNS = [
        r'\[Image:\s*[^\]]*\]',           # [Image: description]
        r'\[Insert\s+[^\]]*\]',           # [Insert image here]
        r'Click to add\s+\w+',            # Click to add title/text
        r'Add\s+\w+\s+here',              # Add text here
        r'Lorem ipsum',                    # Lorem ipsum placeholder
        r'TODO:?\s*',                      # TODO markers
        r'PLACEHOLDER',                    # Explicit placeholder
        r'XXX',                            # XXX markers
        r'\{\{[^}]+\}\}',                  # {{template}} markers
    ]

    # Minimum readable font size in points
    MIN_FONT_SIZE_PT = 10

    # Maximum text that should fit in a shape (rough estimate)
    MAX_CHARS_PER_SHAPE = 2000

    def __init__(self, pptx_path: str, template_path: str = None, layout_plan_path: str = None):
        self.pptx_path = Path(pptx_path)
        self.template_path = Path(template_path) if template_path else None
        self.layout_plan_path = Path(layout_plan_path) if layout_plan_path else None

        if not self.pptx_path.exists():
            raise FileNotFoundError(f"Presentation not found: {pptx_path}")

        self.prs = Presentation(str(self.pptx_path))
        self.template_prs = None
        if self.template_path and self.template_path.exists():
            self.template_prs = Presentation(str(self.template_path))

        self.issues: List[Issue] = []
        self.tier2_results: Optional[Tier2Results] = None

    def check(self, parallel: bool = True) -> QualityReport:
        """Run all quality checks and return report.

        Args:
            parallel: If True, run Tier 2 validators in parallel using
                     ProcessPoolExecutor. If False, run sequentially
                     (useful for debugging).

        Returns:
            QualityReport with all issues, scores, and recommendations.
        """
        self.issues = []
        self.tier2_results = Tier2Results()

        # Run Tier 1 checks (core quality)
        self._check_empty_slides()
        self._check_missing_titles()
        self._check_placeholder_text()
        self._check_text_overflow()
        self._check_table_overflow()
        self._check_font_sizes()
        self._check_consistency()
        self._check_slide_numbers()
        self._check_image_paths()
        self._check_slide_density()
        self._check_layout_repetitiveness()

        # Capture Tier 1 issue counts before Tier 2 validators add more
        tier1_errors = sum(1 for i in self.issues if i.severity == Severity.ERROR)
        tier1_warnings = sum(1 for i in self.issues if i.severity == Severity.WARNING)
        tier1_info = sum(1 for i in self.issues if i.severity == Severity.INFO)

        # Calculate core quality score (Tier 1 only)
        # Errors: -10 each, Warnings: -3 each, Info: -1 each (capped at 0)
        core_score = 100 - (tier1_errors * 10) - (tier1_warnings * 3) - (tier1_info * 1)
        core_score = max(0.0, min(100.0, core_score))

        # Run Tier 2 checks (parallel by default)
        self._run_tier2_validators(parallel=parallel)

        # Get layout coverage
        layout_coverage = self._analyze_layout_coverage()

        # Calculate summary (includes Tier 2 issues)
        summary = {
            'errors': sum(1 for i in self.issues if i.severity == Severity.ERROR),
            'warnings': sum(1 for i in self.issues if i.severity == Severity.WARNING),
            'info': sum(1 for i in self.issues if i.severity == Severity.INFO),
        }

        # Build heuristic score breakdown
        heuristic_breakdown = self._build_heuristic_breakdown(core_score)

        # Use composite score as the main score
        score = heuristic_breakdown.compute_composite_score()

        # Generate recommendations
        recommendations = self._generate_recommendations(summary, layout_coverage)

        return QualityReport(
            file_path=str(self.pptx_path),
            slide_count=len(self.prs.slides),
            issues=self.issues,
            layout_coverage=layout_coverage,
            summary=summary,
            score=score,
            passed=summary['errors'] == 0,
            recommendations=recommendations,
            tier2_results=self.tier2_results if self._has_tier2_results() else None,
            heuristic_breakdown=heuristic_breakdown
        )

    def _build_heuristic_breakdown(self, core_score: float) -> HeuristicScoreBreakdown:
        """Build heuristic score breakdown from Tier 1 and Tier 2 results."""
        breakdown = HeuristicScoreBreakdown(core_quality=core_score)

        # Extract Tier 2 scores if available
        if self.tier2_results:
            if self.tier2_results.typography_hierarchy:
                typo = self.tier2_results.typography_hierarchy
                if typo.get('score') is not None and not typo.get('error'):
                    breakdown.typography_hierarchy = float(typo['score'])

            if self.tier2_results.whitespace:
                ws = self.tier2_results.whitespace
                if ws.get('score') is not None and not ws.get('error'):
                    breakdown.whitespace_analysis = float(ws['score'])

            if self.tier2_results.visual_type:
                vt = self.tier2_results.visual_type
                if vt.get('score') is not None and not vt.get('error'):
                    breakdown.visual_type = float(vt['score'])

            if self.tier2_results.column_balance:
                cb = self.tier2_results.column_balance
                if cb.get('score') is not None and not cb.get('error'):
                    breakdown.column_balance = float(cb['score'])

        return breakdown

    def _has_tier2_results(self) -> bool:
        """Check if any Tier 2 results are available."""
        return (
            self.tier2_results is not None and
            (self.tier2_results.typography_hierarchy is not None or
             self.tier2_results.visual_type is not None or
             self.tier2_results.column_balance is not None or
             self.tier2_results.whitespace is not None)
        )

    def _run_tier2_validators(self, parallel: bool = True) -> None:
        """Run Tier 2 validators and integrate their results.

        Args:
            parallel: If True, run validators in parallel using ProcessPoolExecutor.
                     If False, run sequentially (useful for debugging).
        """
        if parallel:
            self._run_tier2_validators_parallel()
        else:
            self._run_tier2_validators_sequential()

    def _run_tier2_validators_parallel(self) -> None:
        """Run Tier 2 validators in parallel using multiprocessing.

        Uses ProcessPoolExecutor to run validators concurrently:
        - Typography Hierarchy (PPTX-based)
        - Whitespace Analyzer (PPTX-based)
        - Visual Type Checker (layout plan-based, if provided)
        - Column Balance (layout plan-based, if provided)

        Results are collected and integrated back into the main checker.
        """
        # Determine which validators to run
        validators_to_run: List[Tuple[callable, str]] = [
            (_run_typography_validator, str(self.pptx_path)),
            (_run_whitespace_validator, str(self.pptx_path)),
        ]

        # Add layout plan validators if layout plan is provided
        if self.layout_plan_path and self.layout_plan_path.exists():
            validators_to_run.append(
                (_run_visual_type_validator, str(self.layout_plan_path))
            )
            validators_to_run.append(
                (_run_column_balance_validator, str(self.layout_plan_path))
            )

        # Run validators in parallel
        # Use spawn context for clean subprocess state (avoids issues on macOS)
        max_workers = min(len(validators_to_run), multiprocessing.cpu_count())

        results: List[Dict[str, Any]] = []
        try:
            with ProcessPoolExecutor(max_workers=max_workers) as executor:
                # Submit all tasks
                futures = {
                    executor.submit(func, arg): func.__name__
                    for func, arg in validators_to_run
                }

                # Collect results as they complete
                for future in as_completed(futures):
                    validator_name = futures[future]
                    try:
                        result = future.result()
                        results.append(result)
                    except Exception as e:
                        # Handle subprocess failures
                        results.append({
                            'validator': validator_name.replace('_run_', '').replace('_validator', ''),
                            'error': f"Subprocess failed: {str(e)}",
                            'score': None,
                            'passed': None,
                            'issues': [],
                        })
        except Exception:
            # Fallback to sequential execution if parallel fails
            # (e.g., on systems where multiprocessing is restricted)
            self._run_tier2_validators_sequential()
            return

        # Integrate results
        self._integrate_tier2_results(results)

    def _run_tier2_validators_sequential(self) -> None:
        """Run Tier 2 validators sequentially (fallback mode).

        Used when parallel execution is disabled or fails.
        """
        results: List[Dict[str, Any]] = []

        # Typography Hierarchy
        results.append(_run_typography_validator(str(self.pptx_path)))

        # Whitespace Analyzer
        results.append(_run_whitespace_validator(str(self.pptx_path)))

        # Layout plan validators (only if layout plan provided)
        if self.layout_plan_path and self.layout_plan_path.exists():
            results.append(_run_visual_type_validator(str(self.layout_plan_path)))
            results.append(_run_column_balance_validator(str(self.layout_plan_path)))

        # Integrate results
        self._integrate_tier2_results(results)

    def _integrate_tier2_results(self, results: List[Dict[str, Any]]) -> None:
        """Integrate results from parallel/sequential validator execution.

        Args:
            results: List of validator result dictionaries.
        """
        for result in results:
            validator = result.get('validator', '')

            if validator == 'typography_hierarchy':
                self._integrate_typography_result(result)
            elif validator == 'whitespace':
                self._integrate_whitespace_result(result)
            elif validator == 'visual_type':
                self._integrate_visual_type_result(result)
            elif validator == 'column_balance':
                self._integrate_column_balance_result(result)

    def _integrate_typography_result(self, result: Dict[str, Any]) -> None:
        """Integrate typography hierarchy validator result."""
        if result.get('error'):
            self.tier2_results.typography_hierarchy = {
                'error': result['error'],
                'score': None,
                'passed': None,
            }
        else:
            self.tier2_results.typography_hierarchy = {
                'score': result['score'],
                'passed': result['passed'],
                'summary': result.get('summary', {}),
                'global_stats': result.get('global_stats', {}),
            }

            # Convert issues to main Issue format
            for issue_data in result.get('issues', []):
                self.issues.append(Issue(
                    slide_number=issue_data['slide_number'],
                    category=f"typography_{issue_data['category']}",
                    severity=Severity(issue_data['severity']),
                    message=issue_data['message'],
                    suggestion=issue_data['suggestion'],
                    details=issue_data.get('details', {})
                ))

    def _integrate_whitespace_result(self, result: Dict[str, Any]) -> None:
        """Integrate whitespace analyzer result."""
        if result.get('error'):
            self.tier2_results.whitespace = {
                'error': result['error'],
                'score': None,
                'passed': None,
            }
        else:
            self.tier2_results.whitespace = {
                'score': result['score'],
                'passed': result['passed'],
                'summary': result.get('summary', {}),
                'global_stats': result.get('global_stats', {}),
            }

            # Convert issues to main Issue format
            for issue_data in result.get('issues', []):
                self.issues.append(Issue(
                    slide_number=issue_data['slide_number'],
                    category=f"whitespace_{issue_data['category']}",
                    severity=Severity(issue_data['severity']),
                    message=issue_data['message'],
                    suggestion=issue_data['suggestion'],
                    details=issue_data.get('details', {})
                ))

    def _integrate_visual_type_result(self, result: Dict[str, Any]) -> None:
        """Integrate visual type checker result."""
        if result.get('error'):
            self.tier2_results.visual_type = {
                'error': result['error'],
                'score': None,
                'passed': None,
            }
        else:
            self.tier2_results.visual_type = {
                'score': result['score'],
                'passed': result['passed'],
                'summary': result.get('summary', {}),
                'visual_type_distribution': result.get('visual_type_distribution', {}),
            }

            # Convert issues to main Issue format
            for issue_data in result.get('issues', []):
                self.issues.append(Issue(
                    slide_number=issue_data['slide_number'],
                    category=f"visual_type_{issue_data['visual_type']}",
                    severity=Severity(issue_data['severity']),
                    message=issue_data['message'],
                    suggestion=issue_data['suggestion'],
                    details={**issue_data.get('details', {}), 'content_type': issue_data.get('content_type', '')}
                ))

    def _integrate_column_balance_result(self, result: Dict[str, Any]) -> None:
        """Integrate column balance analyzer result."""
        if result.get('error'):
            self.tier2_results.column_balance = {
                'error': result['error'],
                'score': None,
                'passed': None,
            }
        else:
            self.tier2_results.column_balance = {
                'score': result['score'],
                'passed': result['passed'],
                'summary': result.get('summary', {}),
                'slides_with_columns': result.get('slides_with_columns', 0),
            }

            # Convert issues to main Issue format
            for issue_data in result.get('issues', []):
                self.issues.append(Issue(
                    slide_number=issue_data['slide_number'],
                    category="column_balance",
                    severity=Severity(issue_data['severity']),
                    message=issue_data['message'],
                    suggestion=issue_data['suggestion'],
                    details={
                        **issue_data.get('details', {}),
                        'column_numbers': issue_data.get('column_numbers', []),
                        'visual_type': issue_data.get('visual_type', '')
                    }
                ))

    def _check_empty_slides(self):
        """Check for slides with no content."""
        for i, slide in enumerate(self.prs.slides, 1):
            text_content = []
            has_table = False
            has_image = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    # Ignore slide numbers and very short text
                    if text and len(text) > 3:
                        text_content.append(text)
                if shape.has_table:
                    has_table = True
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True

            total_text = ' '.join(text_content)

            if not total_text and not has_table and not has_image:
                self.issues.append(Issue(
                    slide_number=i,
                    category='empty_slide',
                    severity=Severity.ERROR,
                    message="Slide appears to be empty",
                    suggestion="Add content or remove this slide",
                    details={'has_shapes': len(list(slide.shapes))}
                ))
            elif len(total_text) < 20 and not has_table and not has_image:
                self.issues.append(Issue(
                    slide_number=i,
                    category='minimal_content',
                    severity=Severity.WARNING,
                    message=f"Slide has very little content ({len(total_text)} chars)",
                    suggestion="Consider adding more detail or combining with another slide",
                    details={'text_length': len(total_text), 'text_preview': total_text[:50]}
                ))

    # Layouts that use BODY placeholders for title text instead of TITLE placeholders
    # These are valid layouts and should not trigger missing_title warnings
    TITLE_EXEMPT_LAYOUTS = {'title-cover', 'master-base', 'content-image-right-text-left'}

    def _check_missing_titles(self):
        """Check for slides without titles.

        Handles layouts that use BODY (type 2) placeholders for title text
        instead of TITLE (type 1) placeholders by exempting known layouts
        and checking for title-like content in BODY placeholders.
        """
        for i, slide in enumerate(self.prs.slides, 1):
            # Skip first slide (often title slide with different structure)
            if i == 1:
                continue

            # Skip last slide (often contact/thank you with no title)
            if i == len(self.prs.slides):
                continue

            # Skip layouts that intentionally use BODY placeholders for titles
            layout_name = slide.slide_layout.name if slide.slide_layout else ''
            if layout_name in self.TITLE_EXEMPT_LAYOUTS:
                continue

            has_title = False
            for shape in slide.shapes:
                if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                    if shape.placeholder_format.type == 1:  # TITLE
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            has_title = True
                            break

            if not has_title:
                self.issues.append(Issue(
                    slide_number=i,
                    category='missing_title',
                    severity=Severity.WARNING,
                    message="Slide has no title",
                    suggestion="Add a descriptive title to improve navigation",
                    details={'layout': layout_name}
                ))

    def _check_placeholder_text(self):
        """Check for placeholder text that shouldn't be in final presentation."""
        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text

                for pattern in self.PLACEHOLDER_PATTERNS:
                    matches = re.findall(pattern, text, re.IGNORECASE)
                    if matches:
                        self.issues.append(Issue(
                            slide_number=i,
                            category='placeholder_text',
                            severity=Severity.ERROR,
                            message=f"Placeholder text found: '{matches[0]}'",
                            suggestion="Replace with actual content or remove",
                            details={'pattern': pattern, 'match': matches[0]}
                        ))
                        break  # One issue per shape

    def _check_text_overflow(self):
        """Check for text that likely overflows shape bounds.

        Uses PIL-based text measurement when available for accurate overflow detection.
        Falls back to character-count heuristic when PIL is unavailable.
        """
        # Create text measurer (uses PIL when available)
        measurer = TextMeasurer()

        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Skip known placeholder patterns
                if text.startswith('[Image:') or text.startswith('[Insert'):
                    continue

                # Skip slide number placeholders (type 13 = SLIDE_NUMBER)
                if (shape.is_placeholder and
                    hasattr(shape, 'placeholder_format') and
                    shape.placeholder_format.type == 13):
                    continue

                # Get shape dimensions
                if not shape.width or not shape.height:
                    continue

                # Convert EMU to points (914400 EMU = 1 inch, 72 points = 1 inch)
                shape_width_pt = shape.width / 914400 * 72
                shape_height_pt = shape.height / 914400 * 72

                # Get font info from shape (use first run's font as representative)
                font_name = 'Arial'  # Default
                font_size_pt = 12.0  # Default

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_name = run.font.name
                        if run.font.size:
                            font_size_pt = run.font.size.pt
                        break  # Use first run's font
                    if font_name != 'Arial' or font_size_pt != 12.0:
                        break

                # Use PIL-based measurement
                fits, overflow_ratio, details = measurer.text_fits_in_shape(
                    text=text,
                    font_name=font_name,
                    size_pt=font_size_pt,
                    shape_width_pt=shape_width_pt,
                    shape_height_pt=shape_height_pt,
                    margin_pt=10,
                    line_spacing=1.2
                )

                # Skip unmeasurable shapes (zero-height, too small for font, etc.)
                # These produce -1.0 overflow_ratio as a skip marker
                if overflow_ratio < 0:
                    continue

                # Report overflow with 10% threshold (overflow_ratio > 1.10)
                # 10% threshold reduces false positives from:
                # - Line count rounding (int truncation of lines_available)
                # - Font metric variations between measured and actual fonts
                # - PowerPoint's auto-shrink and auto-fit behaviors
                if overflow_ratio > 1.10:
                    lines_n = details.get('lines_needed', 0)
                    lines_a = details.get('lines_available', 0)
                    overflow_pct = round((overflow_ratio - 1.0) * 100, 1)
                    msg = (f"Text exceeds shape bounds by {overflow_pct}% "
                           f"({lines_n} lines, {lines_a} fit)")
                    self.issues.append(Issue(
                        slide_number=i,
                        category='text_overflow',
                        severity=Severity.WARNING,
                        message=msg,
                        suggestion="Reduce text, shrink font, or enlarge shape",
                        details={
                            'text_length': len(text),
                            'lines_needed': details.get('lines_needed', 0),
                            'lines_available': details.get('lines_available', 0),
                            'overflow_ratio': round(overflow_ratio, 2),
                            'font_name': font_name,
                            'font_size_pt': font_size_pt,
                            'text_preview': text[:80] + '...' if len(text) > 80 else text,
                            'using_pil': HAS_PIL,
                        }
                    ))

    def _check_table_overflow(self):
        """Check for table cells with text that overflows cell bounds.

        Uses PIL-based text measurement when available for accurate overflow detection.
        Falls back to character-count heuristic when PIL is unavailable.
        """
        measurer = TextMeasurer()

        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_table:
                    continue

                table = shape.table

                for row_idx, row in enumerate(table.rows):
                    # Get row height in points
                    row_height = row.height
                    if not row_height:
                        continue
                    row_height_pt = row_height / 914400 * 72

                    for col_idx in range(len(table.columns)):
                        # Get column width in points
                        col_width = table.columns[col_idx].width
                        if not col_width:
                            continue
                        col_width_pt = col_width / 914400 * 72

                        cell = table.cell(row_idx, col_idx)
                        text = cell.text_frame.text.strip() if cell.text_frame else ""
                        if not text:
                            continue

                        # Skip placeholder patterns
                        if text.startswith('[Image:') or text.startswith('[Insert'):
                            continue

                        # Get font info from cell (use first run's font as representative)
                        font_name = 'Arial'
                        font_size_pt = 10.0  # Tables often use smaller fonts

                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    font_name = run.font.name
                                if run.font.size:
                                    font_size_pt = run.font.size.pt
                                break
                            if font_name != 'Arial' or font_size_pt != 10.0:
                                break

                        # Use PIL-based measurement with reduced margin for table cells
                        # Table cells typically have smaller internal padding (~5pt)
                        fits, overflow_ratio, details = measurer.text_fits_in_shape(
                            text=text,
                            font_name=font_name,
                            size_pt=font_size_pt,
                            shape_width_pt=col_width_pt,
                            shape_height_pt=row_height_pt,
                            margin_pt=5,  # Smaller margin for table cells
                            line_spacing=1.1  # Tighter line spacing in tables
                        )

                        # Skip unmeasurable cells (zero-height, too small, etc.)
                        if overflow_ratio < 0:
                            continue

                        # Report overflow with 10% threshold (overflow_ratio > 1.10)
                        # 10% threshold reduces false positives from line count rounding
                        if overflow_ratio > 1.10:
                            lines_n = details.get('lines_needed', 0)
                            lines_a = details.get('lines_available', 0)
                            overflow_pct = round((overflow_ratio - 1.0) * 100, 1)
                            msg = (f"Table cell ({row_idx + 1},{col_idx + 1}) text exceeds bounds "
                                   f"by {overflow_pct}% ({lines_n} lines, {lines_a} fit)")
                            self.issues.append(Issue(
                                slide_number=i,
                                category='table_cell_overflow',
                                severity=Severity.WARNING,
                                message=msg,
                                suggestion="Reduce text, shrink font, widen column, or increase row height",
                                details={
                                    'row': row_idx + 1,
                                    'col': col_idx + 1,
                                    'text_length': len(text),
                                    'lines_needed': details.get('lines_needed', 0),
                                    'lines_available': details.get('lines_available', 0),
                                    'overflow_ratio': round(overflow_ratio, 2),
                                    'font_name': font_name,
                                    'font_size_pt': font_size_pt,
                                    'cell_width_pt': round(col_width_pt, 1),
                                    'cell_height_pt': round(row_height_pt, 1),
                                    'text_preview': text[:50] + '...' if len(text) > 50 else text,
                                    'using_pil': HAS_PIL,
                                }
                            ))

    def _check_font_sizes(self):
        """Check for fonts that are too small to read."""
        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            size_pt = run.font.size.pt
                            if size_pt < self.MIN_FONT_SIZE_PT:
                                self.issues.append(Issue(
                                    slide_number=i,
                                    category='small_font',
                                    severity=Severity.INFO,
                                    message=f"Font size {size_pt}pt may be hard to read",
                                    suggestion=f"Consider increasing to at least {self.MIN_FONT_SIZE_PT}pt",
                                    details={
                                        'font_size': size_pt,
                                        'text_preview': run.text[:30] if run.text else ''
                                    }
                                ))
                                break  # One issue per shape

    def _check_consistency(self):
        """Check for consistency issues across slides."""
        # Track font sizes used for titles
        title_sizes = {}
        body_sizes = {}

        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                is_title = (shape.is_placeholder and
                           hasattr(shape, 'placeholder_format') and
                           shape.placeholder_format.type == 1)

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            size = run.font.size.pt
                            if is_title:
                                title_sizes[size] = title_sizes.get(size, 0) + 1
                            else:
                                body_sizes[size] = body_sizes.get(size, 0) + 1

        # Flag if many different title sizes used
        if len(title_sizes) > 3:
            self.issues.append(Issue(
                slide_number=0,  # Global issue
                category='inconsistent_titles',
                severity=Severity.INFO,
                message=f"Title font sizes vary ({len(title_sizes)} different sizes used)",
                suggestion="Consider standardizing title sizes for consistency",
                details={'sizes': dict(title_sizes)}
            ))

    def _check_slide_numbers(self):
        """Check that slide numbers are present and sequential."""
        # This is informational - some designs intentionally omit numbers
        slides_with_numbers = 0

        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                    if shape.placeholder_format.type == 13:  # SLIDE_NUMBER
                        slides_with_numbers += 1
                        break

        total = len(self.prs.slides)
        if slides_with_numbers > 0 and slides_with_numbers < total - 2:
            # Some slides have numbers but not all (excluding first/last)
            self.issues.append(Issue(
                slide_number=0,
                category='inconsistent_numbering',
                severity=Severity.INFO,
                message=f"Slide numbers on {slides_with_numbers}/{total} slides",
                suggestion="Add or remove slide numbers consistently",
                details={'slides_with_numbers': slides_with_numbers, 'total': total}
            ))

    def _check_image_paths(self):
        """Check that image paths referenced in layout plan exist.

        Validates file_path references in:
        - columns[].file_path
        - cards[].file_path
        - deliverables[].file_path
        - extras.image_file / content.image_file
        - extras.background / content.background (when it's a file path)

        Only runs when a layout plan is provided.
        """
        if not self.layout_plan_path or not self.layout_plan_path.exists():
            return

        try:
            layout_plan = json.loads(self.layout_plan_path.read_text())
        except (json.JSONDecodeError, IOError) as e:
            self.issues.append(Issue(
                slide_number=0,
                category='invalid_layout_plan',
                severity=Severity.WARNING,
                message=f"Could not parse layout plan: {e}",
                suggestion="Ensure the layout plan is valid JSON",
                details={'error': str(e)}
            ))
            return

        slides = layout_plan.get('slides', [])
        # Supported image extensions for background path detection
        image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.emf', '.wmf'}

        for slide in slides:
            slide_num = slide.get('slide_number', 0)
            missing_images = []

            # Check columns
            for col in slide.get('columns', []) or []:
                file_path = col.get('file_path')
                if file_path and not self._image_exists(file_path):
                    missing_images.append(('column', file_path))

            # Check cards
            for card in slide.get('cards', []) or []:
                file_path = card.get('file_path')
                if file_path and not self._image_exists(file_path):
                    missing_images.append(('card', file_path))

            # Check deliverables
            for deliv in slide.get('deliverables', []) or []:
                file_path = deliv.get('file_path')
                if file_path and not self._image_exists(file_path):
                    missing_images.append(('deliverable', file_path))

            # Check extras.image_file
            extras = slide.get('extras', {}) or {}
            content = slide.get('content', {}) or {}

            image_file = extras.get('image_file') or content.get('image_file')
            if image_file and not self._image_exists(image_file):
                missing_images.append(('image_file', image_file))

            # Check background (only if it looks like a file path)
            background = extras.get('background') or content.get('background')
            if background and isinstance(background, str):
                # Check if it has an image extension (not just a color or description)
                ext = Path(background).suffix.lower()
                if ext in image_extensions and not self._image_exists(background):
                    missing_images.append(('background', background))

            # Report issues for this slide
            for source, path in missing_images:
                self.issues.append(Issue(
                    slide_number=slide_num,
                    category='missing_image',
                    severity=Severity.ERROR,
                    message=f"Image not found: {path}",
                    suggestion=f"Ensure the image file exists or remove the {source} file_path reference",
                    details={'source': source, 'file_path': path}
                ))

    def _image_exists(self, file_path: str) -> bool:
        """Check if an image file exists.

        Handles both absolute and relative paths (relative to CWD).
        """
        path = Path(file_path)
        if path.is_absolute():
            return path.exists()
        # Try relative to current working directory
        return Path.cwd().joinpath(file_path).exists()

    # Density thresholds (percentage of slide area)
    MAX_CONTENT_DENSITY = 85  # Slides above this are flagged as crowded

    # Layouts exempt from density checks (branding slides with decorative elements)
    DENSITY_EXEMPT_LAYOUTS = {'master-base'}

    def _check_slide_density(self):
        """Check for slides with excessive content density.

        Calculates the percentage of slide area covered by content shapes
        and flags slides that exceed the density threshold.

        Excludes:
        - master-base layouts (branding slides with decorative elements)
        - Slide number, footer, and date placeholders (metadata, not content)

        This is an informational check (INFO severity) since high density
        can be intentional for certain visual types (tables, processes, etc.)
        but may indicate content that should be split across multiple slides.
        """
        # Get slide dimensions in points
        slide_width_pt = self.prs.slide_width / 914400 * 72 if self.prs.slide_width else 720
        slide_height_pt = self.prs.slide_height / 914400 * 72 if self.prs.slide_height else 540
        slide_area = slide_width_pt * slide_height_pt

        for i, slide in enumerate(self.prs.slides, 1):
            # Skip exempt layouts (branding slides)
            layout_name = slide.slide_layout.name if slide.slide_layout else ''
            if layout_name in self.DENSITY_EXEMPT_LAYOUTS:
                continue

            # Calculate content area
            total_shape_area = 0.0

            for shape in slide.shapes:
                # Skip metadata placeholders (slide number, footer, date)
                if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    if ph_type in (13, 14, 15):  # SLIDE_NUMBER, FOOTER, DATE
                        continue

                # Calculate shape area
                if shape.left is not None and shape.top is not None:
                    width_pt = shape.width / 914400 * 72 if shape.width else 0
                    height_pt = shape.height / 914400 * 72 if shape.height else 0
                    total_shape_area += width_pt * height_pt

            # Calculate density percentage
            density_pct = (total_shape_area / slide_area) * 100 if slide_area > 0 else 0

            if density_pct > self.MAX_CONTENT_DENSITY:
                self.issues.append(Issue(
                    slide_number=i,
                    category='high_density',
                    severity=Severity.INFO,
                    message=f"Slide has high content density ({density_pct:.0f}%)",
                    suggestion="Consider splitting content or using a simpler layout",
                    details={
                        'density_pct': round(density_pct, 1),
                        'threshold': self.MAX_CONTENT_DENSITY,
                        'layout': layout_name,
                    }
                ))

    # Threshold for consecutive same-layout warning
    CONSECUTIVE_SAME_LAYOUT_THRESHOLD = 3
    # Layouts exempt from repetitiveness checks (expected to repeat)
    REPETITIVENESS_EXEMPT_LAYOUTS = {'title-cover', 'master-base', 'content-image-right-a'}

    def _check_layout_repetitiveness(self):
        """Check for consecutive slides using the same layout.

        Flags when the same layout is used for 3+ consecutive slides,
        which can make presentations feel monotonous. This is especially
        important for data-heavy presentations with multiple tables.

        Exempts common content layouts (title-cover, master-base,
        content-image-right-a) which naturally repeat without issue.
        """
        if len(self.prs.slides) < self.CONSECUTIVE_SAME_LAYOUT_THRESHOLD:
            return

        # Build list of (slide_number, layout_name, has_table)
        slide_layouts = []
        for i, slide in enumerate(self.prs.slides, 1):
            layout_name = slide.slide_layout.name if slide.slide_layout else 'unknown'
            has_table = any(shape.has_table for shape in slide.shapes)
            slide_layouts.append((i, layout_name, has_table))

        # Find consecutive runs of same layout
        i = 0
        while i < len(slide_layouts):
            current_layout = slide_layouts[i][1]

            # Skip exempt layouts
            if current_layout in self.REPETITIVENESS_EXEMPT_LAYOUTS:
                i += 1
                continue

            # Count consecutive slides with same layout
            run_start = i
            run_length = 1
            tables_in_run = 1 if slide_layouts[i][2] else 0

            while (i + run_length < len(slide_layouts) and
                   slide_layouts[i + run_length][1] == current_layout):
                if slide_layouts[i + run_length][2]:
                    tables_in_run += 1
                run_length += 1

            # Flag if run exceeds threshold
            if run_length >= self.CONSECUTIVE_SAME_LAYOUT_THRESHOLD:
                start_slide = slide_layouts[run_start][0]
                end_slide = slide_layouts[run_start + run_length - 1][0]

                # Higher severity if multiple tables in run
                if tables_in_run >= 2:
                    message = (f"{run_length} consecutive slides use '{current_layout}' "
                              f"(slides {start_slide}-{end_slide}, {tables_in_run} with tables)")
                    suggestion = ("Consider varying layouts for table slides to reduce monotony. "
                                "Try 'content-centered-a' or 'comparison-tables' for variety.")
                else:
                    message = (f"{run_length} consecutive slides use '{current_layout}' "
                              f"(slides {start_slide}-{end_slide})")
                    suggestion = "Consider using different layouts to add visual variety"

                self.issues.append(Issue(
                    slide_number=0,  # Global issue
                    category='layout_repetitiveness',
                    severity=Severity.INFO,
                    message=message,
                    suggestion=suggestion,
                    details={
                        'layout': current_layout,
                        'run_length': run_length,
                        'start_slide': start_slide,
                        'end_slide': end_slide,
                        'tables_in_run': tables_in_run,
                    }
                ))

            i += run_length

    def _analyze_layout_coverage(self) -> Dict[str, Any]:
        """Analyze which layouts are used and suggest better-fitting alternatives."""
        layouts_used: Dict[str, int] = {}
        layouts_available: List[str] = []

        # Count layouts used
        for slide in self.prs.slides:
            if slide.slide_layout:
                name = slide.slide_layout.name
                layouts_used[name] = layouts_used.get(name, 0) + 1

        # Get available layouts from template or presentation
        source = self.template_prs if self.template_prs else self.prs
        for layout in source.slide_layouts:
            layouts_available.append(layout.name)

        # Identify unused layouts
        unused = [layout for layout in layouts_available if layout not in layouts_used]

        # Run layout coverage analysis for suggestions
        layout_suggestions = []
        content_patterns = {}
        try:
            analyzer = LayoutCoverageAnalyzer(
                str(self.pptx_path),
                str(self.template_path) if self.template_path else None
            )
            suggestions = analyzer.analyze()

            # Convert suggestions to serializable format
            for s in suggestions:
                layout_suggestions.append({
                    'slide_number': s.slide_number,
                    'current_layout': s.current_layout,
                    'suggested_layout': s.suggested_layout,
                    'reason': s.reason,
                    'confidence': round(s.confidence, 2),
                    'content_pattern': s.content_pattern.value,
                })

                # Add INFO-level issue for each suggestion
                self.issues.append(Issue(
                    slide_number=s.slide_number,
                    category='layout_suggestion',
                    severity=Severity.INFO,
                    message=f"Consider using '{s.suggested_layout}' instead of '{s.current_layout}'",
                    suggestion=s.reason,
                    details={
                        'current_layout': s.current_layout,
                        'suggested_layout': s.suggested_layout,
                        'content_pattern': s.content_pattern.value,
                        'confidence': round(s.confidence, 2),
                    }
                ))

            # Get content pattern distribution
            summary = analyzer.get_summary()
            content_patterns = summary.get('content_patterns_detected', {})

        except Exception as e:
            # Log but don't fail if layout analysis fails
            layout_suggestions = []
            content_patterns = {'error': str(e)}

        return {
            'used': layouts_used,
            'available_count': len(layouts_available),
            'used_count': len(layouts_used),
            'unused': unused[:10],  # Limit output
            'coverage_pct': round(len(layouts_used) / len(layouts_available) * 100, 1) if layouts_available else 0,
            'suggestions': layout_suggestions,
            'content_patterns': content_patterns,
        }

    def _generate_recommendations(self, summary: Dict[str, int], coverage: Dict[str, Any]) -> List[str]:
        """Generate actionable recommendations based on issues found."""
        recommendations = []

        if summary['errors'] > 0:
            recommendations.append(
                f"Fix {summary['errors']} error(s) before delivery - these are critical issues."
            )

        if summary['warnings'] > 3:
            recommendations.append(
                f"Review {summary['warnings']} warning(s) - most should be addressed."
            )

        # Category-specific recommendations
        categories = {}
        for issue in self.issues:
            categories[issue.category] = categories.get(issue.category, 0) + 1

        if categories.get('placeholder_text', 0) > 0:
            recommendations.append(
                "Replace all placeholder text with actual content before delivery."
            )

        if categories.get('empty_slide', 0) > 0:
            recommendations.append(
                "Remove empty slides or add content to them."
            )

        if categories.get('text_overflow', 0) > 2:
            recommendations.append(
                "Several slides have text overflow - consider condensing content or using multiple slides."
            )

        if categories.get('missing_image', 0) > 0:
            recommendations.append(
                "One or more images referenced in the layout plan are missing. "
                "Ensure all image files exist before generation."
            )

        if coverage['coverage_pct'] < 20:
            recommendations.append(
                f"Only using {coverage['used_count']} of {coverage['available_count']} layouts - "
                "explore more layout options for visual variety."
            )

        if not recommendations:
            recommendations.append("Presentation looks good! No major issues found.")

        return recommendations


def get_exit_code(summary: Dict[str, int]) -> ExitCode:
    """Determine exit code based on worst issue severity.

    Args:
        summary: Dict with 'errors', 'warnings', 'info' counts

    Returns:
        ExitCode based on the most severe issue type found
    """
    if summary.get('errors', 0) > 0:
        return ExitCode.ERROR
    if summary.get('warnings', 0) > 0:
        return ExitCode.WARNING
    if summary.get('info', 0) > 0:
        return ExitCode.INFO
    return ExitCode.CLEAN


def format_terminal_report(report: QualityReport) -> str:
    """Format report for terminal output with visual hierarchy."""
    lines = []

    # Header
    lines.append("")
    lines.append("=" * 60)
    lines.append("  PRESENTATION QUALITY REPORT")
    lines.append("=" * 60)
    lines.append(f"  File: {Path(report.file_path).name}")
    lines.append(f"  Slides: {report.slide_count}")
    lines.append("")

    # Score with visual indicator
    score_bar = "" * int(report.score / 5) + "" * (20 - int(report.score / 5))
    status = " PASS" if report.passed else " NEEDS WORK"
    lines.append(f"  Score: {report.score:.0f}/100  [{score_bar}]  {status}")
    lines.append("")

    # Summary
    lines.append("-" * 60)
    lines.append("  SUMMARY")
    lines.append("-" * 60)
    lines.append(f"    Errors:   {report.summary['errors']:3d}  (must fix)")
    lines.append(f"    Warnings: {report.summary['warnings']:3d}  (should fix)")
    lines.append(f"    Info:     {report.summary['info']:3d}  (nice to fix)")
    lines.append("")

    # Heuristic Score Breakdown (if available)
    if report.heuristic_breakdown:
        lines.append("-" * 60)
        lines.append("  HEURISTIC SCORE BREAKDOWN")
        lines.append("-" * 60)
        hb = report.heuristic_breakdown

        # Component scores with weight indicators
        def format_component(name: str, score: Optional[float], weight: int) -> str:
            if score is None:
                return f"    {name:22s}  --/100  (weight: {weight}%)"
            bar = "" * int(score / 10) + "" * (10 - int(score / 10))
            return f"    {name:22s} {score:3.0f}/100 [{bar}] (weight: {weight}%)"

        lines.append(format_component("Core Quality", hb.core_quality, hb.WEIGHTS['core_quality']))
        lines.append(format_component("Typography", hb.typography_hierarchy, hb.WEIGHTS['typography_hierarchy']))
        lines.append(format_component("Whitespace", hb.whitespace_analysis, hb.WEIGHTS['whitespace_analysis']))
        lines.append(format_component("Visual Type", hb.visual_type, hb.WEIGHTS['visual_type']))
        lines.append(format_component("Column Balance", hb.column_balance, hb.WEIGHTS['column_balance']))
        lines.append("")

        # Composite calculation explanation
        components = []
        if hb.core_quality is not None:
            components.append('core_quality')
        if hb.typography_hierarchy is not None:
            components.append('typography_hierarchy')
        if hb.whitespace_analysis is not None:
            components.append('whitespace_analysis')
        if hb.visual_type is not None:
            components.append('visual_type')
        if hb.column_balance is not None:
            components.append('column_balance')

        active_weight = sum(hb.WEIGHTS[c] for c in components)
        lines.append(f"    Active components: {len(components)}/5 ({active_weight}%  normalized to 100%)")
        lines.append("")

    # Issues by slide
    if report.issues:
        lines.append("-" * 60)
        lines.append("  ISSUES BY SLIDE")
        lines.append("-" * 60)

        # Group by slide
        by_slide: Dict[int, List[Issue]] = {}
        for issue in report.issues:
            slide = issue.slide_number
            if slide not in by_slide:
                by_slide[slide] = []
            by_slide[slide].append(issue)

        # Global issues first
        if 0 in by_slide:
            lines.append("  [Global]")
            for issue in by_slide[0]:
                icon = "" if issue.severity == Severity.ERROR else "!" if issue.severity == Severity.WARNING else ""
                lines.append(f"    {icon} {issue.message}")
            lines.append("")
            del by_slide[0]

        # Per-slide issues
        for slide_num in sorted(by_slide.keys()):
            issues = by_slide[slide_num]
            has_error = any(i.severity == Severity.ERROR for i in issues)
            icon = "" if has_error else "!"
            lines.append(f"  Slide {slide_num} {icon}")
            for issue in issues:
                sev_icon = "" if issue.severity == Severity.ERROR else "!" if issue.severity == Severity.WARNING else ""
                lines.append(f"    {sev_icon} {issue.message}")
                if issue.severity == Severity.ERROR:
                    lines.append(f"       {issue.suggestion}")
            lines.append("")

    # Layout coverage
    lines.append("-" * 60)
    lines.append("  LAYOUT COVERAGE")
    lines.append("-" * 60)
    cov = report.layout_coverage
    lines.append(f"    Using {cov['used_count']} of {cov['available_count']} layouts ({cov['coverage_pct']}%)")
    if cov['used']:
        lines.append("    Most used:")
        for name, count in sorted(cov['used'].items(), key=lambda x: -x[1])[:5]:
            lines.append(f"      {count}x {name}")
    lines.append("")

    # Layout Suggestions (if any)
    suggestions = cov.get('suggestions', [])
    if suggestions:
        lines.append("-" * 60)
        lines.append("  LAYOUT SUGGESTIONS")
        lines.append("-" * 60)
        lines.append(f"    Found {len(suggestions)} slide(s) that might benefit from different layouts:")
        lines.append("")
        for s in suggestions[:10]:  # Limit to top 10
            conf_pct = int(s['confidence'] * 100)
            conf_bar = "" * (conf_pct // 10) + "" * (10 - conf_pct // 10)
            lines.append(f"    Slide {s['slide_number']}: {s['current_layout']}")
            lines.append(f"       Suggest: {s['suggested_layout']}")
            lines.append(f"      Pattern: {s['content_pattern']} | Confidence: [{conf_bar}] {conf_pct}%")
            lines.append(f"      {s['reason']}")
            lines.append("")

    # Content patterns (if available)
    content_patterns = cov.get('content_patterns', {})
    if content_patterns and not content_patterns.get('error'):
        lines.append("    Content patterns detected:")
        for pattern, count in sorted(content_patterns.items(), key=lambda x: -x[1])[:6]:
            lines.append(f"      {count}x {pattern}")
        lines.append("")

    # Tier 2 Results (if available)
    if report.tier2_results:
        lines.append("-" * 60)
        lines.append("  TIER 2 VALIDATION")
        lines.append("-" * 60)

        # Typography Hierarchy
        if report.tier2_results.typography_hierarchy:
            typo = report.tier2_results.typography_hierarchy
            if typo.get('error'):
                lines.append(f"    Typography Hierarchy: ERROR - {typo['error']}")
            else:
                status = " PASS" if typo.get('passed') else " ISSUES"
                lines.append(f"    Typography Hierarchy: {typo.get('score', 0):.0f}/100 {status}")
                if typo.get('global_stats'):
                    stats = typo['global_stats']
                    if stats.get('title'):
                        lines.append(f"      Title: {stats['title']['min']:.0f}-{stats['title']['max']:.0f}pt")
                    if stats.get('body'):
                        lines.append(f"      Body: {stats['body']['min']:.0f}-{stats['body']['max']:.0f}pt")

        # Visual Type Appropriateness
        if report.tier2_results.visual_type:
            vt = report.tier2_results.visual_type
            if vt.get('error'):
                lines.append(f"    Visual Type Check: ERROR - {vt['error']}")
            else:
                status = " PASS" if vt.get('passed') else " ISSUES"
                lines.append(f"    Visual Type Check: {vt.get('score', 0):.0f}/100 {status}")
                if vt.get('visual_type_distribution'):
                    top_types = list(vt['visual_type_distribution'].items())[:3]
                    if top_types:
                        types_str = ", ".join(f"{t}({c})" for t, c in top_types)
                        lines.append(f"      Top types: {types_str}")

        # Column Balance
        if report.tier2_results.column_balance:
            cb = report.tier2_results.column_balance
            if cb.get('error'):
                lines.append(f"    Column Balance: ERROR - {cb['error']}")
            else:
                status = " PASS" if cb.get('passed') else " ISSUES"
                lines.append(f"    Column Balance: {cb.get('score', 0):.0f}/100 {status}")
                if cb.get('slides_with_columns'):
                    lines.append(f"      Multi-column slides: {cb['slides_with_columns']}")

        # Whitespace Analysis
        if report.tier2_results.whitespace:
            ws = report.tier2_results.whitespace
            if ws.get('error'):
                lines.append(f"    Whitespace Analysis: ERROR - {ws['error']}")
            else:
                status = " PASS" if ws.get('passed') else " ISSUES"
                lines.append(f"    Whitespace Analysis: {ws.get('score', 0):.0f}/100 {status}")
                if ws.get('global_stats'):
                    stats = ws['global_stats']
                    if stats.get('density'):
                        lines.append(f"      Density: {stats['density']['min']:.0f}-{stats['density']['max']:.0f}%")
                    if stats.get('margins', {}).get('left'):
                        margins = stats['margins']
                        lines.append(f"      Margins: {margins['left']['avg']:.0f}pt avg")

        lines.append("")

    # Recommendations
    lines.append("-" * 60)
    lines.append("  RECOMMENDATIONS")
    lines.append("-" * 60)
    for i, rec in enumerate(report.recommendations, 1):
        lines.append(f"    {i}. {rec}")
    lines.append("")
    lines.append("=" * 60)

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description='Check presentation quality without reference file'
    )
    parser.add_argument('presentation', help='PPTX file to check')
    parser.add_argument('--template', '-t', help='Template PPTX for layout coverage analysis')
    parser.add_argument('--layout-plan', '-l', help='Layout plan JSON for Tier 2 validation')
    parser.add_argument('--json', action='store_true', help='Output as JSON')
    parser.add_argument('--output', '-o', help='Output file path')
    parser.add_argument('--strict', action='store_true', help='Fail on any warnings')
    parser.add_argument('--no-parallel', action='store_true',
                        help='Disable parallel Tier 2 validation (useful for debugging)')

    args = parser.parse_args()

    try:
        checker = QualityChecker(args.presentation, args.template, args.layout_plan)
        report = checker.check(parallel=not args.no_parallel)
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error checking presentation: {e}", file=sys.stderr)
        sys.exit(1)

    # Output
    if args.json:
        output = json.dumps(report.to_dict(), indent=2)
    else:
        output = format_terminal_report(report)

    if args.output:
        Path(args.output).write_text(output)
        print(f"Report saved to: {args.output}")
    else:
        print(output)

    # Exit code based on worst issue severity
    # --strict mode: treat warnings same as errors (exit 2+ means failure)
    exit_code = get_exit_code(report.summary)
    if args.strict and exit_code.value >= ExitCode.WARNING.value:
        # In strict mode, any warning or error is a failure
        sys.exit(exit_code.value)
    sys.exit(exit_code.value)


if __name__ == '__main__':
    main()
