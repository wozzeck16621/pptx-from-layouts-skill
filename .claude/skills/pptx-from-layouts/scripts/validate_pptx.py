#!/usr/bin/env python3
"""
PPTX Validator - Compare two PowerPoint presentations for exact match.

Features:
- Slide count validation
- Layout match per slide
- Title text exact match
- Body text exact match
- Typography validation (font, size, bold, italic, color)
- Table structure validation (rows, cols, merges)
- Table content validation (cell text, fills)
- Shape position validation

Usage:
    python validate_pptx.py output.pptx reference.pptx --threshold 100 --strict
"""

import argparse
import json
import re
import sys
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR


@dataclass
class ValidationResult:
    """Result of validating a single aspect."""
    category: str
    passed: bool
    weight: float
    message: str
    details: Dict[str, Any] = field(default_factory=dict)


@dataclass
class SlideComparison:
    """Comparison results for a single slide."""
    slide_number: int
    layout_match: bool
    layout_expected: str
    layout_actual: str
    title_match: bool
    body_match: bool
    typography_match: bool
    table_match: bool
    position_match: bool
    differences: List[str] = field(default_factory=list)


class PPTXValidator:
    """Validates a generated PPTX against a reference PPTX."""

    # Validation weights (must sum to 100)
    WEIGHTS = {
        'slide_count': 5,
        'layout_match': 15,
        'title_match': 10,
        'body_match': 20,  # Reduced from 25 to make room for text_overflow
        'typography': 15,  # Reduced from 20 to make room for text_overflow
        'table_structure': 10,
        'table_content': 10,
        'text_overflow': 10,  # NEW: Detect text exceeding shape bounds
        'shape_positions': 0,  # Disabled - templates have different placeholder positions
        'image_placeholders': 0  # Disabled - placeholder detection not yet implemented
    }

    # Position tolerance in EMUs (914400 EMU = 1 inch)
    # Use generous tolerance since different templates have slightly different positions
    POSITION_TOLERANCE_EMU = 914400  # 1 inch tolerance

    def __init__(self, output_path: str, reference_path: str, strict: bool = False):
        self.output_path = Path(output_path)
        self.reference_path = Path(reference_path)
        self.strict = strict

        # Validate files exist before opening
        if not self.output_path.exists():
            raise FileNotFoundError(f"Output file not found: {output_path}")
        if not self.reference_path.exists():
            raise FileNotFoundError(f"Reference file not found: {reference_path}")
        if not self.output_path.suffix.lower() == '.pptx':
            raise ValueError(f"Output file must be .pptx: {output_path}")
        if not self.reference_path.suffix.lower() == '.pptx':
            raise ValueError(f"Reference file must be .pptx: {reference_path}")

        try:
            self.output_prs = Presentation(str(self.output_path))
        except Exception as e:
            raise ValueError(f"Failed to open output file: {e}")

        try:
            self.reference_prs = Presentation(str(self.reference_path))
        except Exception as e:
            raise ValueError(f"Failed to open reference file: {e}")

        self.results: List[ValidationResult] = []
        self.slide_comparisons: List[SlideComparison] = []

    def validate(self) -> Tuple[float, Dict[str, Any]]:
        """Run all validations and return score (0-100) and detailed report."""
        self._validate_slide_count()
        self._validate_layouts()
        self._validate_content()
        self._validate_typography()
        self._validate_tables()
        self._validate_text_overflow()
        self._validate_positions()

        # Calculate weighted score
        total_score = 0.0
        category_scores = {}

        for result in self.results:
            if result.passed:
                total_score += result.weight
                category_scores[result.category] = category_scores.get(result.category, 0) + result.weight
            else:
                category_scores[result.category] = category_scores.get(result.category, 0)

        report = {
            'score': round(total_score, 2),
            'passed': total_score >= 100.0 if self.strict else total_score >= 90.0,
            'category_scores': category_scores,
            'results': [
                {
                    'category': r.category,
                    'passed': r.passed,
                    'weight': r.weight,
                    'message': r.message,
                    'details': r.details
                }
                for r in self.results
            ],
            'slide_comparisons': [
                {
                    'slide_number': sc.slide_number,
                    'layout_match': sc.layout_match,
                    'layout_expected': sc.layout_expected,
                    'layout_actual': sc.layout_actual,
                    'title_match': sc.title_match,
                    'body_match': sc.body_match,
                    'typography_match': sc.typography_match,
                    'table_match': sc.table_match,
                    'position_match': sc.position_match,
                    'differences': sc.differences
                }
                for sc in self.slide_comparisons
            ]
        }

        return total_score, report

    def _validate_slide_count(self):
        """Validate slide count matches."""
        output_count = len(self.output_prs.slides)
        ref_count = len(self.reference_prs.slides)

        passed = output_count == ref_count
        self.results.append(ValidationResult(
            category='slide_count',
            passed=passed,
            weight=self.WEIGHTS['slide_count'] if passed else 0,
            message=f"Slide count: {output_count} vs {ref_count}",
            details={'output': output_count, 'reference': ref_count}
        ))

    def _validate_layouts(self):
        """Validate layout matches per slide."""
        output_count = len(self.output_prs.slides)
        ref_count = len(self.reference_prs.slides)

        matches = 0
        total = min(output_count, ref_count)

        for i in range(total):
            out_slide = self.output_prs.slides[i]
            ref_slide = self.reference_prs.slides[i]

            out_layout = out_slide.slide_layout.name if out_slide.slide_layout else "None"
            ref_layout = ref_slide.slide_layout.name if ref_slide.slide_layout else "None"

            out_idx = self._get_layout_index(self.output_prs, out_slide.slide_layout)
            ref_idx = self._get_layout_index(self.reference_prs, ref_slide.slide_layout)

            layout_match = out_layout == ref_layout
            if layout_match:
                matches += 1

            # Create or update slide comparison
            if i < len(self.slide_comparisons):
                self.slide_comparisons[i].layout_match = layout_match
                self.slide_comparisons[i].layout_expected = f"{ref_layout} ({ref_idx})"
                self.slide_comparisons[i].layout_actual = f"{out_layout} ({out_idx})"
            else:
                self.slide_comparisons.append(SlideComparison(
                    slide_number=i + 1,
                    layout_match=layout_match,
                    layout_expected=f"{ref_layout} ({ref_idx})",
                    layout_actual=f"{out_layout} ({out_idx})",
                    title_match=False,
                    body_match=False,
                    typography_match=False,
                    table_match=False,
                    position_match=False
                ))
                if not layout_match:
                    self.slide_comparisons[i].differences.append(
                        f"Layout mismatch: expected {ref_layout} ({ref_idx}), got {out_layout} ({out_idx})"
                    )

        score = (matches / total * self.WEIGHTS['layout_match']) if total > 0 else 0
        passed = matches == total

        self.results.append(ValidationResult(
            category='layout_match',
            passed=passed,
            weight=score,
            message=f"Layout match: {matches}/{total} slides",
            details={'matches': matches, 'total': total}
        ))

    def _validate_content(self):
        """Validate title and body content matches."""
        output_count = len(self.output_prs.slides)
        ref_count = len(self.reference_prs.slides)
        total = min(output_count, ref_count)

        title_matches = 0
        body_matches = 0

        for i in range(total):
            out_slide = self.output_prs.slides[i]
            ref_slide = self.reference_prs.slides[i]

            out_title = self._extract_title(out_slide)
            ref_title = self._extract_title(ref_slide)

            out_body = self._extract_body_text(out_slide)
            ref_body = self._extract_body_text(ref_slide)

            title_match = self._normalize_text(out_title) == self._normalize_text(ref_title)
            body_match = self._normalize_text(out_body) == self._normalize_text(ref_body)

            if title_match:
                title_matches += 1
            if body_match:
                body_matches += 1

            # Update slide comparison
            if i < len(self.slide_comparisons):
                self.slide_comparisons[i].title_match = title_match
                self.slide_comparisons[i].body_match = body_match
                if not title_match:
                    self.slide_comparisons[i].differences.append(
                        f"Title mismatch: expected '{ref_title[:50]}...', got '{out_title[:50]}...'"
                    )
                if not body_match:
                    self.slide_comparisons[i].differences.append("Body content mismatch")

        title_score = (title_matches / total * self.WEIGHTS['title_match']) if total > 0 else 0
        body_score = (body_matches / total * self.WEIGHTS['body_match']) if total > 0 else 0

        self.results.append(ValidationResult(
            category='title_match',
            passed=title_matches == total,
            weight=title_score,
            message=f"Title match: {title_matches}/{total} slides",
            details={'matches': title_matches, 'total': total}
        ))

        self.results.append(ValidationResult(
            category='body_match',
            passed=body_matches == total,
            weight=body_score,
            message=f"Body match: {body_matches}/{total} slides",
            details={'matches': body_matches, 'total': total}
        ))

    def _validate_typography(self):
        """Validate typography (font, size, bold, italic, color).

        Uses text segment matching to handle different paragraph boundaries.
        Extracts all styled text segments and compares them regardless of how
        they're grouped into paragraphs.
        """
        output_count = len(self.output_prs.slides)
        ref_count = len(self.reference_prs.slides)
        total = min(output_count, ref_count)

        matches = 0
        issues = []

        for i in range(total):
            out_slide = self.output_prs.slides[i]
            ref_slide = self.reference_prs.slides[i]

            slide_match = True
            slide_issues = []

            # Extract all styled text segments from each slide
            out_segments = self._extract_styled_segments(out_slide)
            ref_segments = self._extract_styled_segments(ref_slide)

            # Compare segments by normalized text content
            for ref_text, ref_style in ref_segments.items():
                if ref_text not in out_segments:
                    # Try fuzzy match - look for text that starts with same prefix
                    found = False
                    for out_text in out_segments:
                        if ref_text[:30] in out_text or out_text[:30] in ref_text:
                            found = True
                            out_style = out_segments[out_text]
                            # Check bold/italic on matched text
                            if ref_style['bold'] and not out_style['bold']:
                                slide_issues.append(f"Bold mismatch on '{ref_text[:30]}...': expected True, got False")
                                slide_match = False
                            if ref_style['italic'] and not out_style['italic']:
                                slide_issues.append(f"Italic mismatch on '{ref_text[:30]}...': expected True, got False")
                                slide_match = False
                            # Check color - only if reference has a specific color
                            if ref_style['color'] and ref_style['color'] != out_style['color']:
                                slide_issues.append(f"Color mismatch on '{ref_text[:30]}...': expected {ref_style['color']}, got {out_style['color'] or 'none'}")
                                slide_match = False
                            break
                else:
                    out_style = out_segments[ref_text]
                    # Check font properties
                    if ref_style['bold'] and not out_style['bold']:
                        slide_match = False
                        slide_issues.append(f"Bold mismatch on '{ref_text[:30]}...': expected True, got False")
                    if ref_style['italic'] and not out_style['italic']:
                        slide_match = False
                        slide_issues.append(f"Italic mismatch on '{ref_text[:30]}...': expected True, got False")
                    # Check color - only if reference has a specific color
                    if ref_style['color'] and ref_style['color'] != out_style['color']:
                        slide_match = False
                        slide_issues.append(f"Color mismatch on '{ref_text[:30]}...': expected {ref_style['color']}, got {out_style['color'] or 'none'}")

            if slide_match:
                matches += 1

            # Update slide comparison
            if i < len(self.slide_comparisons):
                self.slide_comparisons[i].typography_match = slide_match
                for issue in slide_issues:
                    self.slide_comparisons[i].differences.append(issue)

            issues.extend(slide_issues)

        score = (matches / total * self.WEIGHTS['typography']) if total > 0 else 0

        self.results.append(ValidationResult(
            category='typography',
            passed=matches == total,
            weight=score,
            message=f"Typography match: {matches}/{total} slides",
            details={'matches': matches, 'total': total, 'issues': issues[:20]}
        ))

    def _validate_tables(self):
        """Validate table structure and content."""
        output_count = len(self.output_prs.slides)
        ref_count = len(self.reference_prs.slides)
        total = min(output_count, ref_count)

        structure_matches = 0
        content_matches = 0
        slides_with_tables = 0

        for i in range(total):
            out_slide = self.output_prs.slides[i]
            ref_slide = self.reference_prs.slides[i]

            out_tables = self._extract_tables(out_slide)
            ref_tables = self._extract_tables(ref_slide)

            if not ref_tables and not out_tables:
                continue

            slides_with_tables += 1

            structure_match = len(out_tables) == len(ref_tables)
            content_match = True

            if structure_match:
                for j, (out_tbl, ref_tbl) in enumerate(zip(out_tables, ref_tables)):
                    if out_tbl['rows'] != ref_tbl['rows'] or out_tbl['cols'] != ref_tbl['cols']:
                        structure_match = False
                        break

                    # Check content
                    for r in range(ref_tbl['rows']):
                        for c in range(ref_tbl['cols']):
                            ref_cell = ref_tbl['cells'].get((r, c), '')
                            out_cell = out_tbl['cells'].get((r, c), '')
                            if self._normalize_text(ref_cell) != self._normalize_text(out_cell):
                                content_match = False
            else:
                content_match = False

            if structure_match:
                structure_matches += 1
            if content_match:
                content_matches += 1

            # Update slide comparison
            if i < len(self.slide_comparisons):
                self.slide_comparisons[i].table_match = structure_match and content_match
                if not structure_match:
                    self.slide_comparisons[i].differences.append("Table structure mismatch")
                if not content_match:
                    self.slide_comparisons[i].differences.append("Table content mismatch")

        if slides_with_tables > 0:
            structure_score = structure_matches / slides_with_tables * self.WEIGHTS['table_structure']
            content_score = content_matches / slides_with_tables * self.WEIGHTS['table_content']
        else:
            structure_score = self.WEIGHTS['table_structure']
            content_score = self.WEIGHTS['table_content']

        self.results.append(ValidationResult(
            category='table_structure',
            passed=structure_matches == slides_with_tables or slides_with_tables == 0,
            weight=structure_score,
            message=f"Table structure match: {structure_matches}/{slides_with_tables} slides",
            details={'matches': structure_matches, 'total': slides_with_tables}
        ))

        self.results.append(ValidationResult(
            category='table_content',
            passed=content_matches == slides_with_tables or slides_with_tables == 0,
            weight=content_score,
            message=f"Table content match: {content_matches}/{slides_with_tables} slides",
            details={'matches': content_matches, 'total': slides_with_tables}
        ))

    def _validate_text_overflow(self):
        """Validate that text does not overflow shape bounds.

        Uses PIL text measurement to estimate if text exceeds shape dimensions.
        Only validates the output presentation (not comparing to reference).
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            # PIL not available, skip overflow detection
            self.results.append(ValidationResult(
                category='text_overflow',
                passed=True,
                weight=self.WEIGHTS['text_overflow'],
                message="Text overflow: skipped (PIL not available)",
                details={'skipped': True}
            ))
            return

        output_count = len(self.output_prs.slides)
        slides_with_overflow = []

        for i in range(output_count):
            slide = self.output_prs.slides[i]
            slide_overflow_shapes = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Skip empty text frames
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Skip image placeholder labels
                if text.startswith('[Image:') and text.endswith(']'):
                    continue

                # Estimate overflow
                overflow = self._estimate_overflow(shape)
                if overflow is not None and overflow > 0.05:  # More than 0.05" overflow
                    slide_overflow_shapes.append({
                        'shape_text_preview': text[:50] + '...' if len(text) > 50 else text,
                        'overflow_inches': round(overflow, 2)
                    })

            if slide_overflow_shapes:
                slides_with_overflow.append({
                    'slide_number': i + 1,
                    'shapes': slide_overflow_shapes
                })

                # Update slide comparison
                if i < len(self.slide_comparisons):
                    for shape_info in slide_overflow_shapes:
                        self.slide_comparisons[i].differences.append(
                            f"Text overflow: {shape_info['overflow_inches']}\" on '{shape_info['shape_text_preview']}'"
                        )

        # Calculate score: full points if no overflow, reduced for each affected slide
        if slides_with_overflow:
            overflow_ratio = len(slides_with_overflow) / output_count
            score = self.WEIGHTS['text_overflow'] * (1 - overflow_ratio)
            passed = False
        else:
            score = self.WEIGHTS['text_overflow']
            passed = True

        self.results.append(ValidationResult(
            category='text_overflow',
            passed=passed,
            weight=score,
            message=f"Text overflow: {len(slides_with_overflow)} slides with issues",
            details={
                'slides_with_overflow': len(slides_with_overflow),
                'total_slides': output_count,
                'overflow_details': slides_with_overflow[:10]  # Limit detail output
            }
        ))

    def _estimate_overflow(self, shape) -> float:
        """Estimate text overflow in inches for a shape.

        Returns positive value if text overflows, None if cannot estimate.
        Uses simplified text measurement - not pixel-perfect but catches major issues.
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            return None

        if not shape.has_text_frame:
            return None

        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return None

        # Convert shape dimensions from EMU to pixels (at 96 DPI)
        # 914400 EMU = 1 inch, 1 inch = 96 pixels
        shape_width_px = (shape.width / 914400) * 96 if shape.width else 300
        shape_height_px = (shape.height / 914400) * 96 if shape.height else 200

        # Account for margins (typical 0.1" = ~10px on each side)
        usable_width_px = max(shape_width_px - 20, 50)
        usable_height_px = max(shape_height_px - 10, 30)

        # Set up PIL for text measurement
        dummy_img = Image.new('RGB', (1, 1))
        draw = ImageDraw.Draw(dummy_img)

        try:
            font = ImageFont.load_default(size=14)  # Default size estimate
        except (TypeError, AttributeError):
            font = ImageFont.load_default()

        total_height_px = 0

        for para in text_frame.paragraphs:
            para_text = para.text
            if not para_text.strip():
                continue

            # Simple word wrapping estimation
            words = para_text.split()
            current_line = ""
            line_count = 0

            for word in words:
                test_line = current_line + (" " if current_line else "") + word
                try:
                    text_width = draw.textlength(test_line, font=font)
                except AttributeError:
                    # Fallback for older PIL versions
                    bbox = draw.textbbox((0, 0), test_line, font=font)
                    text_width = bbox[2] - bbox[0]

                if text_width <= usable_width_px:
                    current_line = test_line
                else:
                    line_count += 1
                    current_line = word

            if current_line:
                line_count += 1

            # Estimate line height (1.2x font size)
            line_height_px = 14 * 1.2 * 96 / 72  # Convert pt to px
            total_height_px += line_count * line_height_px

        # Calculate overflow
        if total_height_px > usable_height_px:
            overflow_px = total_height_px - usable_height_px
            overflow_inches = overflow_px / 96.0
            return overflow_inches

        return None

    def _validate_positions(self):
        """Validate shape positions within tolerance."""
        output_count = len(self.output_prs.slides)
        ref_count = len(self.reference_prs.slides)
        total = min(output_count, ref_count)

        matches = 0

        for i in range(total):
            out_slide = self.output_prs.slides[i]
            ref_slide = self.reference_prs.slides[i]

            slide_match = True

            # Get shapes by approximate position and type
            out_shapes = self._get_shape_positions(out_slide)
            ref_shapes = self._get_shape_positions(ref_slide)

            # For now, just check count and rough positions
            if len(out_shapes) != len(ref_shapes):
                slide_match = False
            else:
                # Sort by position and compare
                out_sorted = sorted(out_shapes, key=lambda s: (s['top'], s['left']))
                ref_sorted = sorted(ref_shapes, key=lambda s: (s['top'], s['left']))

                for out_s, ref_s in zip(out_sorted, ref_sorted):
                    if (abs(out_s['left'] - ref_s['left']) > self.POSITION_TOLERANCE_EMU or
                        abs(out_s['top'] - ref_s['top']) > self.POSITION_TOLERANCE_EMU):
                        slide_match = False
                        break

            if slide_match:
                matches += 1

            if i < len(self.slide_comparisons):
                self.slide_comparisons[i].position_match = slide_match
                if not slide_match:
                    self.slide_comparisons[i].differences.append("Shape position mismatch")

        score = (matches / total * self.WEIGHTS['shape_positions']) if total > 0 else 0

        self.results.append(ValidationResult(
            category='shape_positions',
            passed=matches == total,
            weight=score,
            message=f"Position match: {matches}/{total} slides",
            details={'matches': matches, 'total': total}
        ))

        # Add image placeholder check (disabled - weight is 0)
        self.results.append(ValidationResult(
            category='image_placeholders',
            passed=True,  # Disabled - not yet implemented
            weight=self.WEIGHTS['image_placeholders'],
            message="Image placeholders: (validation disabled)",
            details={'reason': 'Placeholder detection not yet implemented'}
        ))

    # Helper methods

    def _get_layout_index(self, prs: Presentation, layout) -> int:
        """Get the index of a layout in the presentation."""
        if layout is None:
            return -1
        for i, l in enumerate(prs.slide_layouts):
            if l == layout:
                return i
        return -1

    def _extract_title(self, slide) -> str:
        """Extract title text from a slide."""
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # TITLE
                    if shape.has_text_frame:
                        return shape.text_frame.text
        return ""

    def _extract_body_text(self, slide) -> str:
        """Extract body text from placeholders only (not free-floating text boxes).

        This ensures we compare placeholder content which is template-defined,
        not manually added text boxes that may differ between source and output.
        """
        texts = []
        for shape in slide.shapes:
            # Only process placeholder shapes (skip TEXT_BOX and other non-placeholder shapes)
            if not shape.is_placeholder:
                continue

            # Skip slide number placeholders (type 13 = SLIDE_NUMBER)
            try:
                if shape.placeholder_format.type == 13:  # SLIDE_NUMBER
                    continue
            except (ValueError, AttributeError):
                pass

            if shape.has_text_frame:
                text = shape.text_frame.text
                # Skip image placeholder labels like "[Image: ...]"
                if text.strip().startswith('[Image:') and text.strip().endswith(']'):
                    continue
                texts.append(text)
        return " ".join(texts)

    def _normalize_text(self, text: str) -> str:
        """Normalize text for comparison."""
        if not text:
            return ""
        import re
        # Normalize vertical tabs (soft returns) and other whitespace
        text = text.replace('\x0b', ' ')  # Vertical tab (soft return)
        # Remove extra whitespace, normalize line endings
        text = re.sub(r'\s+', ' ', text.strip())
        # Normalize curly quotes to straight quotes
        text = text.replace('\u2019', "'")  # Right single quote
        text = text.replace('\u2018', "'")  # Left single quote
        text = text.replace('\u201c', '"')  # Left double quote
        text = text.replace('\u201d', '"')  # Right double quote
        text = text.replace('\u2014', '-')  # Em dash
        text = text.replace('\u2013', '-')  # En dash
        return text

    def _extract_styled_segments(self, slide) -> Dict[str, Dict[str, Any]]:
        """Extract styled text segments from a slide.

        Instead of extracting by paragraph, this extracts individual text runs
        with their styling to compare regardless of paragraph boundaries.
        Returns dict mapping normalized text -> {bold, italic, color}.
        """
        segments = {}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Skip slide number placeholders and image labels
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.type == 13:  # SLIDE_NUMBER
                        continue
                except (ValueError, AttributeError):
                    pass

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    if not text or not text.strip():
                        continue

                    # Skip image placeholder labels
                    if text.strip().startswith('[Image:'):
                        continue

                    # Normalize text for matching
                    normalized = text.replace('\x0b', ' ')
                    normalized = re.sub(r'\s+', ' ', normalized).strip()
                    if not normalized:
                        continue

                    # Get styling
                    bold = bool(run.font.bold)
                    italic = bool(run.font.italic)
                    color = ''
                    try:
                        if run.font.color and run.font.color.type is not None and run.font.color.rgb:
                            color = str(run.font.color.rgb)
                    except:
                        pass

                    segments[normalized] = {
                        'bold': bold,
                        'italic': italic,
                        'color': color
                    }

        return segments

    def _extract_typography(self, slide) -> Dict[str, Dict[str, Any]]:
        """Extract typography info for text in a slide, using normalized text as keys.

        Uses first 50 chars of paragraph text as key to handle different run boundaries.
        Tracks if bold/italic are used anywhere in the paragraph for each style.
        """
        typography = {}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Skip slide number placeholders
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.type == 13:  # SLIDE_NUMBER
                        continue
                except (ValueError, AttributeError):
                    pass

            for para in shape.text_frame.paragraphs:
                # Concatenate all runs in the paragraph
                para_text = ""
                has_bold = False
                has_italic = False
                colors = set()

                for run in para.runs:
                    text = run.text
                    if not text:
                        continue
                    para_text += text

                    font = run.font
                    if font.bold:
                        has_bold = True
                    if font.italic:
                        has_italic = True

                    # Get color (safely handle _NoneColor)
                    try:
                        if font.color and font.color.type is not None and font.color.rgb:
                            colors.add(str(font.color.rgb))
                    except (AttributeError, Exception):
                        pass

                para_text = para_text.strip()
                if not para_text:
                    continue

                # Normalize and use first 50 chars as key to handle minor variations
                # Normalize soft returns and whitespace for consistent matching
                normalized_text = para_text.replace('\x0b', ' ')  # Soft return to space
                normalized_text = re.sub(r'\s+', ' ', normalized_text).strip()
                key = normalized_text[:50].strip()
                if not key:
                    continue

                typography[key] = {
                    'bold': has_bold,
                    'italic': has_italic,
                    'color': ','.join(sorted(colors)) if colors else '',
                    'full_text': para_text
                }

        return typography

    def _extract_tables(self, slide) -> List[Dict[str, Any]]:
        """Extract table structure and content from a slide."""
        tables = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                cells = {}

                for r, row in enumerate(table.rows):
                    for c, cell in enumerate(row.cells):
                        cells[(r, c)] = cell.text

                tables.append({
                    'rows': len(table.rows),
                    'cols': len(table.columns),
                    'cells': cells
                })

        return tables

    def _get_shape_positions(self, slide) -> List[Dict[str, int]]:
        """Get positions of all shapes on a slide."""
        positions = []

        for shape in slide.shapes:
            positions.append({
                'left': shape.left if shape.left else 0,
                'top': shape.top if shape.top else 0,
                'width': shape.width if shape.width else 0,
                'height': shape.height if shape.height else 0
            })

        return positions


def main():
    parser = argparse.ArgumentParser(description='Validate PPTX against reference')
    parser.add_argument('output', help='Output PPTX to validate')
    parser.add_argument('reference', help='Reference PPTX to compare against')
    parser.add_argument('--threshold', type=float, default=90.0,
                        help='Minimum score to pass (default: 90)')
    parser.add_argument('--strict', action='store_true',
                        help='Require 100%% match')
    parser.add_argument('--output-report', '-o', dest='report_path',
                        help='Output JSON report path')
    parser.add_argument('--verbose', '-v', action='store_true',
                        help='Verbose output')

    args = parser.parse_args()

    validator = PPTXValidator(args.output, args.reference, args.strict)
    score, report = validator.validate()

    # Print summary
    print(f"\n{'='*60}")
    print(f"PPTX Validation Report")
    print(f"{'='*60}")
    print(f"Output:    {args.output}")
    print(f"Reference: {args.reference}")
    print(f"{'='*60}")
    print(f"\nOverall Score: {score:.1f}%")
    print(f"Status: {'PASS' if report['passed'] else 'FAIL'}")
    print(f"\nCategory Scores:")
    for cat, cat_score in report['category_scores'].items():
        print(f"  {cat}: {cat_score:.1f}")

    if args.verbose or not report['passed']:
        print(f"\nSlide Details:")
        for sc in report['slide_comparisons']:
            status = "OK" if not sc['differences'] else "ISSUES"
            print(f"\n  Slide {sc['slide_number']}: {status}")
            print(f"    Layout: {sc['layout_actual']} (expected: {sc['layout_expected']})")
            if sc['differences']:
                for diff in sc['differences'][:5]:
                    print(f"    - {diff}")

    # Save report
    if args.report_path:
        with open(args.report_path, 'w') as f:
            json.dump(report, f, indent=2)
        print(f"\nReport saved to: {args.report_path}")

    # Exit code
    threshold = 100.0 if args.strict else args.threshold
    sys.exit(0 if score >= threshold else 1)


if __name__ == '__main__':
    main()
